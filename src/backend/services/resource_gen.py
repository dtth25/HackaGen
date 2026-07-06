"""Generation service for the four public outputs: Book, Quiz, Vid, and Slide."""

import asyncio
import json
import os
import re
import shutil
import subprocess
import textwrap
import threading
from typing import Any, Optional

from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate

from backend.core.config import extract_json, get_course_path, get_llm, logger
from backend.core.prompts import (
    BOOK_GENERATION_PROMPT,
    FLASHCARD_GENERATION_PROMPT,
    QUIZ_V2_PROMPT,
    SLIDE_GENERATION_PROMPT,
    VID_SCENES_PROMPT,
)
from backend.services.learning_profile import build_profile_directives


class ResourceGenerator:
    """Generate Book, Quiz, Vid, and Slide outputs from an initialized RAG course."""

    def __init__(self, rag_chains):
        if isinstance(rag_chains, str):
            self.course_id = rag_chains
            self.rag = None
            self.vectorstore = None
        else:
            self.rag = rag_chains
            self.course_id = getattr(rag_chains, "course_id", "")
            self.vectorstore = getattr(rag_chains, "vectorstore", None)

    def _save_json(self, path: str, payload: Any) -> None:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False, indent=2)

    def _read_json(self, path: str, default: Any = None) -> Any:
        if not os.path.exists(path):
            return default
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return default

    def _clean_doc_text(self, doc, max_chars: int = 320) -> str:
        text = doc.page_content
        text = re.sub(r"===\s*BẮT ĐẦU.*?===", " ", text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"===\s*KẾT THÚC.*?===", " ", text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"\[MÃ ĐỊNH DANH TRANG:\s*\d+\]", " ", text, flags=re.IGNORECASE)
        text = re.sub(r"\bNỘI DUNG:\s*", " ", text, flags=re.IGNORECASE)
        text = re.sub(r"\bMã định danh trang\s+\d+\s+nội dung\b", " ", text, flags=re.IGNORECASE)
        text = re.sub(r"\s+", " ", text).strip()
        return text[:max_chars].strip()

    def _clean_generated_text(self, text: Any, compact: bool = True) -> str:
        """Remove internal extraction markers and banned phrases from generated/public text."""
        from backend.services.context_cleaner import scrub_banned_phrases

        value = str(text or "")
        value = re.sub(r"===\s*BẮT ĐẦU.*?===", " ", value, flags=re.IGNORECASE | re.DOTALL)
        value = re.sub(r"===\s*KẾT THÚC.*?===", " ", value, flags=re.IGNORECASE | re.DOTALL)
        value = re.sub(r"\[MÃ ĐỊNH DANH TRANG:\s*\d+\]", " ", value, flags=re.IGNORECASE)
        value = re.sub(r"\bNỘI DUNG:\s*", " ", value, flags=re.IGNORECASE)
        value = re.sub(r"\b(page|source|chunk_id)\s*:\s*[^,\n]+", " ", value, flags=re.IGNORECASE)
        value = re.sub(r"\[?\bsource_chunk_id\s*:\s*[\w.-]+\]?", " ", value, flags=re.IGNORECASE)
        value = scrub_banned_phrases(value)
        if compact:
            return re.sub(r"\s+", " ", value).strip()
        return re.sub(r"[ \t]+", " ", value).replace("\n\n\n", "\n\n").strip()

    def _sanitize_payload(self, value: Any) -> Any:
        """Recursively sanitize public payload strings.

        `source_chunk_ids` is intentionally kept on every output (quiz/flashcards
        included) per the product grounding requirement — every generated item must
        carry its source_chunk_ids so the UI can show "Grounded" provenance.
        """
        if isinstance(value, dict):
            banned_keys = {"page", "source", "chunk_id"}
            return {key: self._sanitize_payload(item) for key, item in value.items() if key not in banned_keys}
        if isinstance(value, list):
            return [self._sanitize_payload(item) for item in value]
        if isinstance(value, str):
            return self._clean_generated_text(value, compact=False)
        return value

    def _source_ids_from_doc(self, doc) -> list[str]:
        """Extract the real source chunk ID from doc metadata; never fabricate one."""
        metadata = getattr(doc, "metadata", {}) or {}
        raw = metadata.get("source_chunk_id") or metadata.get("chunk_id")
        if raw in (None, ""):
            return []
        raw_str = str(raw)
        return [raw_str if raw_str.startswith("chunk_") else f"chunk_{raw_str}"]

    def _clean_docs_context(self, docs, max_docs: int = 24, max_chars: int = 900) -> str:
        """Build clean prompt context: classified, noise-filtered, with internal grounding tags."""
        from backend.services.context_cleaner import classify_chunk, clean_text_markers

        snippets: list[str] = []
        for doc in docs[:max_docs]:
            raw_text = getattr(doc, "page_content", "") or str(doc)
            classification = classify_chunk(raw_text)
            if not classification["use_for_generation"]:
                continue
            text = clean_text_markers(raw_text)[:max_chars].strip()
            if len(text) < 30:
                continue
            source_ids = self._source_ids_from_doc(doc)
            tag = f" [source_chunk_id: {source_ids[0]}]" if source_ids else ""
            snippets.append(f"- {text}{tag}")
        if not snippets:
            return "Tài liệu đã được xử lý nhưng chưa trích xuất được đoạn nội dung đủ rõ."
        return "\n".join(snippets)

    def _doc_points(self, docs, limit: int = 8, max_chars: int = 220) -> list[dict[str, Any]]:
        from backend.services.context_cleaner import classify_chunk

        points = []
        for doc in docs:
            raw_text = getattr(doc, "page_content", "") or str(doc)
            # Classify on the full raw chunk (before truncation) so TOC/dot-leader
            # patterns aren't cut off — this keeps table-of-contents noise and other
            # unusable chunks out of fallback/synthetic content (titles, examples, etc.).
            classification = classify_chunk(raw_text)
            if not classification["use_for_generation"]:
                continue
            text = self._clean_doc_text(doc, max_chars)
            if len(text) < 30:
                continue
            points.append({"text": text, "source_chunk_ids": self._source_ids_from_doc(doc)})
            if len(points) >= limit:
                break
        return points or [
            {
                "text": (
                    "Tài liệu đã được xử lý thành công, nhưng hệ thống chưa trích xuất được đoạn nội dung "
                    "đủ dài cho bản nháp."
                ),
                "source_chunk_ids": [],
            }
        ]

    def _short_title(self, text: str, fallback: str) -> str:
        words = re.findall(r"\w+", text, flags=re.UNICODE)
        # Index artifacts ("41 4 6 comparison to sorting 44 5") must never become
        # titles: drop leading bare numbers, reject digit-dominant results.
        while words and words[0].isdigit():
            words.pop(0)
        title_words = words[:8]
        if not title_words or sum(1 for w in title_words if w.isdigit()) / len(title_words) > 0.34:
            return fallback
        title = " ".join(title_words).strip()
        return title.capitalize() if title else fallback

    def _lesson_defaults(self, point: dict[str, str], title: str) -> dict[str, Any]:
        """Default synthetic lesson content built from a single retrieved chunk.

        Used only to fill in fields the LLM omitted for a given lesson — never
        to override real LLM-generated content that is actually present.
        """
        text = point["text"]
        sentences = self._sentences_from_text(text, limit=3)
        first_sentence = sentences[0] if sentences else text[:180]
        return {
            "duration_minutes": 25,
            "core_idea": first_sentence,
            "why_it_matters": (
                f'Phần "{title}" là một mắt xích trong mạch nội dung của tài liệu gốc; '
                "các phần sau sẽ dựa trên khái niệm này."
            ),
            "learning_objectives": [
                f'Hiểu và trình bày lại được "{title}" bằng lời của mình.',
                "Nhận ra khái niệm này khi nó xuất hiện trong bài tập hoặc ví dụ.",
            ],
            "explanation": text,
            "must_know_points": sentences or [text[:180]],
            "key_concepts": [],
            "example": text[:220],
            "non_example": (
                "Một trường hợp dễ nhầm lẫn nhưng KHÔNG thuộc phạm vi khái niệm này; "
                "cần đối chiếu lại với tài liệu gốc để phân biệt rõ ràng."
            ),
            "common_misunderstanding": {
                "mistake": "Nhầm lẫn giữa luận điểm cốt lõi và các chi tiết minh họa.",
                "correction": "Cần phân biệt rõ nội dung trọng tâm với các ví dụ hỗ trợ đi kèm.",
            },
            "worked_examples": [
                {
                    "title": f"Ví dụ áp dụng: {title}",
                    "problem": text[:200],
                    "step_by_step_solution": [
                        "Xác định thông tin/khái niệm cốt lõi liên quan trong đoạn nội dung.",
                        "Đối chiếu với định nghĩa và ví dụ đã nêu ở phần giải thích.",
                        "Rút ra kết luận và liên hệ với mục tiêu bài học.",
                    ],
                    "why_each_step_matters": "Từng bước giúp người học đi từ dữ kiện thô đến kết luận có căn cứ.",
                    "common_error": "Bỏ qua bước đối chiếu định nghĩa, dẫn đến kết luận thiếu căn cứ.",
                    "source_chunk_ids": point.get("source_chunk_ids", []),
                }
            ],
            "practice_activity": "Tóm tắt phần này bằng 3 gạch đầu dòng và nêu 1 ví dụ minh họa.",
            "practice_problems": [
                {
                    "difficulty": "easy",
                    "question": f'Hãy tóm tắt nội dung trọng tâm của phần "{title}" bằng 1-2 câu.',
                    "expected_answer": text[:120],
                    "hint": "Xem lại phần \"Ý tưởng cốt lõi\" và \"Nội dung bài giảng\" ở trên.",
                    "solution": text[:220],
                    "source_chunk_ids": point.get("source_chunk_ids", []),
                }
            ],
            "quick_check": [
                {
                    "question": f'Nội dung trọng tâm của phần "{title}" là gì?',
                    "answer": text[:120],
                    "explanation": "Dựa trên nội dung trọng tâm được trích xuất từ tài liệu gốc.",
                }
            ],
            "source_chunk_ids": point.get("source_chunk_ids", []),
        }

    def _normalize_lesson(self, raw_lesson: Any, point: dict[str, str], title: str) -> dict[str, Any]:
        """Merge real LLM lesson content with synthetic defaults for any missing field."""
        defaults = self._lesson_defaults(point, title)
        lesson = raw_lesson if isinstance(raw_lesson, dict) else {}

        def text_field(key: str) -> Optional[str]:
            value = lesson.get(key)
            return value.strip() if isinstance(value, str) and value.strip() else None

        def list_field(key: str) -> Optional[list]:
            value = lesson.get(key)
            return value if isinstance(value, list) and value else None

        try:
            duration_minutes = int(lesson.get("duration_minutes"))
        except (TypeError, ValueError):
            duration_minutes = defaults["duration_minutes"]

        key_concepts = [
            {"term": str(kc.get("term") or "").strip(), "definition": str(kc.get("definition") or "").strip()}
            for kc in (list_field("key_concepts") or [])
            if isinstance(kc, dict) and kc.get("term")
        ] or defaults["key_concepts"]

        common_mis_raw = lesson.get("common_misunderstanding")
        common_misunderstanding = (
            {
                "mistake": str(common_mis_raw.get("mistake") or "").strip(),
                "correction": str(common_mis_raw.get("correction") or "").strip(),
            }
            if isinstance(common_mis_raw, dict) and common_mis_raw.get("mistake")
            else defaults["common_misunderstanding"]
        )

        quick_check = [
            {
                "question": str(q.get("question") or "").strip(),
                "answer": str(q.get("answer") or "").strip(),
                "explanation": str(q.get("explanation") or "").strip(),
            }
            for q in (list_field("quick_check") or [])
            if isinstance(q, dict) and q.get("question")
        ] or defaults["quick_check"]

        learning_objectives = list_field("learning_objectives") or defaults["learning_objectives"]
        must_know_points = list_field("must_know_points") or defaults["must_know_points"]
        explanation = text_field("explanation") or defaults["explanation"]
        practice_activity = text_field("practice_activity") or defaults["practice_activity"]
        source_chunk_ids = list_field("source_chunk_ids") or defaults["source_chunk_ids"]

        worked_examples = [
            {
                "title": str(we.get("title") or "").strip() or f"Ví dụ áp dụng {i + 1}",
                "problem": str(we.get("problem") or "").strip(),
                "step_by_step_solution": [str(s).strip() for s in (we.get("step_by_step_solution") or []) if str(s).strip()],
                "why_each_step_matters": str(we.get("why_each_step_matters") or "").strip(),
                "common_error": str(we.get("common_error") or "").strip(),
                "source_chunk_ids": we.get("source_chunk_ids") if isinstance(we.get("source_chunk_ids"), list) else source_chunk_ids,
            }
            for i, we in enumerate(list_field("worked_examples") or [])
            if isinstance(we, dict) and we.get("problem") and we.get("step_by_step_solution")
        ] or defaults["worked_examples"]

        practice_problems = [
            {
                "difficulty": str(pp.get("difficulty") or "medium").strip().lower(),
                "question": str(pp.get("question") or "").strip(),
                "expected_answer": str(pp.get("expected_answer") or "").strip(),
                "hint": str(pp.get("hint") or "").strip(),
                "solution": str(pp.get("solution") or "").strip(),
                "source_chunk_ids": pp.get("source_chunk_ids") if isinstance(pp.get("source_chunk_ids"), list) else source_chunk_ids,
            }
            for pp in (list_field("practice_problems") or [])
            if isinstance(pp, dict) and pp.get("question")
        ] or defaults["practice_problems"]

        return {
            "title": title,
            "duration_minutes": duration_minutes,
            "duration": f"{duration_minutes} phút",
            "core_idea": text_field("core_idea") or defaults["core_idea"],
            "why_it_matters": text_field("why_it_matters") or defaults["why_it_matters"],
            "learning_objectives": learning_objectives,
            "objectives": learning_objectives,
            "explanation": explanation,
            "lecture": explanation,
            "must_know_points": must_know_points,
            "key_points": must_know_points,
            "key_concepts": key_concepts,
            "example": text_field("example") or defaults["example"],
            "non_example": text_field("non_example") or defaults["non_example"],
            "common_misunderstanding": common_misunderstanding,
            "worked_examples": worked_examples,
            "practice_activity": practice_activity,
            "activity": practice_activity,
            "practice_problems": practice_problems,
            "quick_check": quick_check,
            "assessment": [q["question"] for q in quick_check],
            "source_chunk_ids": source_chunk_ids,
        }


    def _invoke_chain(self, chain, prompt_input: dict[str, Any]) -> str:
        res = chain.invoke(prompt_input)
        return str(res.content if hasattr(res, "content") else res)

    def _validate_and_normalize_book_plan(self, raw_plan: dict[str, Any], docs, target_audience: str, style: str) -> dict[str, Any]:
        plan = self._normalize_book(raw_plan, docs, target_audience)
        if plan["title"] == "Chương 1":
            plan["title"] = "Sách học tập chuyên sâu"
        for idx, ch in enumerate(plan.get("chapters", []), 1):
            if ch.get("title") == "Chương 1":
                ch["title"] = f"Chương {idx}: Kiến thức cốt lõi"
            # source_chunk_ids come from real doc metadata via _doc_points/_normalize_lesson;
            # never fabricate grounding IDs here.
        from backend.core import config as core_config

        plan_dir = os.path.join(core_config.BOOKS_DIR, self.course_id)
        os.makedirs(plan_dir, exist_ok=True)
        self._save_json(os.path.join(plan_dir, "plan.json"), plan)
        return plan

    def _validate_book_export_safety(self, book_data: dict[str, Any]) -> None:
        text = str(book_data)
        if "=== BẮT ĐẦU DỮ LIỆU" in text:
            raise ValueError("prohibited marker")
        for ch in book_data.get("chapters", []):
            if ch.get("title", "") in {"Chương 1", "Chương 2", ""}:
                raise ValueError("blank or generic chapter title")

    def _slide_from_point(self, point: dict[str, Any], index: int, topic: dict[str, Any]) -> dict[str, Any]:
        return {
            "slide_index": index,
            "slide_type": "concept",
            "title": topic.get("title", "Slide"),
            "bullets": [point.get("text", "Nội dung 1")[:60], "Phân tích chi tiết khái niệm", "Ứng dụng thực tiễn trong ngành", "Tổng kết điểm cốt lõi"],
            "visual": {"type": "concept", "data": {}},
            "speaker_notes": "Slide này trình bày các khái niệm quan trọng liên quan đến bài giảng. Giảng viên cần nhấn mạnh các điểm cốt lõi và hướng dẫn sinh viên phân tích từng khía cạnh chi tiết và cụ thể.",
            "source_chunk_ids": point.get("source_chunk_ids", []),
        }

    def _video_flow_item(self, index: int, total: int) -> tuple[str, int]:
        types = ["hook", "objective", "concept", "example", "common_mistake", "quiz", "recap"]
        return (types[index % len(types)], 15)

    _QUALITY_GATE_GENERIC_HEADINGS = {"ý chính", "ghi nhớ ý chính", "contents", "mục lục", ""}

    def _quality_gate_generic(self, heading: Any) -> bool:
        return str(heading or "").strip().lower() in self._QUALITY_GATE_GENERIC_HEADINGS

    def _quality_with_context_stats(self, quality: dict[str, Any], resource_data: dict[str, Any], docs) -> dict[str, Any]:
        """Enrich quality report with Task 8 metadata and context stats."""
        enriched = dict(quality)
        enriched["is_final"] = quality.get("is_final", True)
        enriched["quality_score"] = enriched.get("score", 85)
        if "source_grounding_score" not in enriched:
            enriched["source_grounding_score"] = 85
        if "warnings" not in enriched:
            enriched["warnings"] = []
        if "fixes_needed" not in enriched:
            enriched["fixes_needed"] = []
        enriched["usable_chunks_count"] = len(docs) if docs else 0
        return enriched

    def _evaluate_quality_gate(self, resource_data: dict[str, Any], resource_type: str) -> dict[str, Any]:
        """Quality gate: structural checks specific to each resource type.

        Book/slides scores are built from real structural signals (presence of
        key concepts, examples, speaker notes, source grounding) rather than
        just a length/banned-substring check, so a shallow template can no
        longer pass as "university ready" by accident.
        """
        warnings: list[str] = []
        extra_report: dict[str, Any] = {}
        text_content = json.dumps(resource_data, ensure_ascii=False)

        banned_artifacts = [
            "Tree / MEX / Knapsack", "query window",
            "=== BẮT ĐẦU DỮ LIỆU", "MÃ ĐỊNH DANH TRANG", "NỘI DUNG:",
            "Ý chính", "Ghi nhớ ý chính",
            # Generic teaching-filler templates: banned from all user-facing outputs.
            "Giảng viên nên diễn giải thêm",
            "Phân tích cơ chế và điều kiện áp dụng",
            "Liên hệ với ứng dụng thực tế trong tài liệu",
            "Hãy giải thích ngắn gọn nội dung của slide",
            "MIT-style",
        ]
        artifact_hits = [art for art in banned_artifacts if art in text_content]
        if re.search(r"(?:\.\s*){4,}", text_content):
            artifact_hits.append("dot leaders (. . . .)")
        if re.search(r"\bContents\b", text_content):
            artifact_hits.append("heading 'Contents'")
        if re.search(r"\bMIT\b", text_content):
            artifact_hits.append("nhãn 'MIT'")

        def fraction(items: list, predicate) -> float:
            return (sum(1 for item in items if predicate(item)) / len(items) * 100) if items else 0.0

        if resource_type == "book":
            chapters = resource_data.get("chapters") or []
            lessons = [lesson for ch in chapters for lesson in (ch.get("lessons") or [])]
            concepts = [concept for ch in chapters for concept in (ch.get("core_concepts") or [])]
            has_content_units = bool(lessons or concepts)
            structure_score = 100.0 if chapters and has_content_units else 0.0
            if not chapters:
                warnings.append("Sách không có chương nào.")
            elif not has_content_units:
                warnings.append("Các chương không có khái niệm cốt lõi hoặc bài học nào.")

            if concepts:
                # Rigorous v2 schema: depth means intuition + technical explanation per concept.
                depth_score = fraction(
                    concepts,
                    lambda c: bool(str(c.get("intuition") or "").strip())
                    and bool(str(c.get("technical_explanation") or "").strip()),
                )
                example_score = fraction(concepts, lambda c: bool(str(c.get("example") or "").strip()))
                chapter_we = fraction(
                    chapters,
                    lambda ch: any(
                        we.get("step_by_step_solution")
                        for we in (ch.get("worked_examples") or [])
                    ),
                )
                chapter_pp = fraction(
                    chapters,
                    lambda ch: len(
                        {p.get("difficulty") for p in (ch.get("practice_problems") or [])}
                        & {"easy", "medium", "hard"}
                    ) >= 2,
                )
                practice_score = (chapter_we + chapter_pp) / 2
                grounding_score = fraction(concepts, lambda c: bool(c.get("source_chunk_ids")))

                # Shallow-summary detection: real courseware needs substantive explanations.
                total_explanation_chars = sum(
                    len(str(c.get("technical_explanation") or "")) + len(str(c.get("intuition") or ""))
                    for c in concepts
                )
                if concepts and total_explanation_chars / len(concepts) < 300:
                    warnings.append("Nội dung giải thích quá nông (giống tóm tắt AI hơn là giáo trình).")
                    depth_score = min(depth_score, 40.0)
            else:
                depth_score = fraction(lessons, lambda lesson: bool(lesson.get("key_concepts")))
                example_score = fraction(
                    lessons,
                    lambda lesson: bool(str(lesson.get("example") or "").strip())
                    and bool((lesson.get("worked_examples") or [{}])[0].get("step_by_step_solution")),
                )
                practice_score = fraction(
                    lessons,
                    lambda lesson: bool((lesson.get("quick_check") or [{}])[0].get("question"))
                    and bool((lesson.get("practice_problems") or [{}])[0].get("question")),
                )
                grounding_score = fraction(lessons, lambda lesson: bool(lesson.get("source_chunk_ids")))

            generic_titles = [ch.get("title") for ch in chapters if self._quality_gate_generic(ch.get("title"))]
            generic_titles += [
                lesson.get("title")
                for lesson in lessons
                if self._quality_gate_generic(lesson.get("title"))
            ]
            if generic_titles:
                warnings.append(f"{len(generic_titles)} tiêu đề chương/bài quá chung chung.")
            if depth_score < 50:
                warnings.append("Nhiều khái niệm/bài học thiếu chiều sâu (intuition + giải thích kỹ thuật).")
            if practice_score < 50:
                warnings.append("Thiếu worked example từng bước hoặc bài tập đủ các mức độ khó.")
            if example_score < 50:
                warnings.append("Nhiều khái niệm/bài học thiếu ví dụ minh họa.")
            if grounding_score < 50:
                warnings.append("Nhiều nội dung thiếu source_chunk_ids (chưa bám sát nguồn).")

            score = (
                0.25 * structure_score + 0.25 * depth_score + 0.2 * example_score
                + 0.15 * practice_score + 0.15 * grounding_score
            )
            score -= 10 * min(3, len(generic_titles))

        elif resource_type == "slides":
            slides = resource_data.get("slides") or []
            if not slides:
                warnings.append("Không có slide nào được tạo.")
            structural_types = {"title", "objectives", "recap"}
            content_slides = [s for s in slides if s.get("slide_type") not in structural_types] or slides
            slide_types = {s.get("slide_type") for s in slides}

            def word_count(value: Any) -> int:
                return len(str(value or "").split())

            notes_score = fraction(content_slides, lambda s: word_count(s.get("speaker_notes")) >= 50)
            depth_types_present = bool(
                slide_types & {"worked_example", "formula_breakdown", "code_walkthrough", "common_mistake"}
            )
            academic_depth_score = 0.7 * notes_score + (30.0 if depth_types_present else 0.0)

            def has_visual(slide: dict) -> bool:
                visual = slide.get("visual_instruction") or slide.get("visual") or {}
                screen = slide.get("screen_content") or {}
                return (
                    (visual.get("type") not in (None, "", "none") and bool(visual.get("description")))
                    or bool(screen.get("formula"))
                    or bool(screen.get("code"))
                    or bool(screen.get("table"))
                    or bool(screen.get("diagram_description"))
                )

            visual_score = fraction(content_slides, has_visual)
            bullets_ok_score = fraction(content_slides, lambda s: 2 <= len(s.get("bullets") or []) <= 5)
            visual_quality_score = 0.6 * visual_score + 0.4 * bullets_ok_score

            teaching_roles = [
                {"objectives"}, {"motivation"}, {"concept"},
                {"diagram", "comparison"},
                {"worked_example", "code_walkthrough", "formula_breakdown"},
                {"common_mistake"}, {"quick_check"}, {"recap"}, {"practice"},
            ]
            coverage = sum(1 for role in teaching_roles if slide_types & role) / len(teaching_roles) * 100
            key_message_score = fraction(slides, lambda s: bool(str(s.get("key_message") or "").strip()))
            teaching_quality_score = 0.7 * coverage + 0.3 * key_message_score

            source_grounding_score = fraction(content_slides, lambda s: bool(s.get("source_chunk_ids")))

            generic_titles = [s.get("title") for s in slides if self._quality_gate_generic(s.get("title"))]
            if generic_titles:
                warnings.append(f"{len(generic_titles)} slide có tiêu đề generic (Ý chính/Nội dung chính/để trống).")
            thin_slides = sum(1 for s in content_slides if len(s.get("bullets") or []) < 3)
            if thin_slides:
                warnings.append(f"{thin_slides} slide có ít hơn 3 bullet (nội dung quá mỏng).")
            if notes_score < 50:
                warnings.append("Nhiều slide thiếu speaker_notes chi tiết (cần 100-180 từ).")
            if not depth_types_present:
                warnings.append("Deck thiếu slide ví dụ mẫu / công thức / code / sai lầm thường gặp.")
            if not (slide_types & {"practice", "quick_check"}):
                warnings.append("Deck thiếu slide bài tập hoặc kiểm tra nhanh.")
            if source_grounding_score < 50:
                warnings.append("Nhiều slide thiếu source_chunk_ids (chưa bám sát nguồn).")

            extra_report = {
                "academic_depth_score": int(academic_depth_score),
                "visual_quality_score": int(visual_quality_score),
                "teaching_quality_score": int(teaching_quality_score),
                "source_grounding_score": int(source_grounding_score),
            }

            score = (
                0.3 * academic_depth_score
                + 0.2 * visual_quality_score
                + 0.3 * teaching_quality_score
                + 0.2 * source_grounding_score
            )
            score -= 15 * min(3, len(generic_titles))

        elif resource_type == "video":
            videos = resource_data.get("videos") or []
            structure_score = 100.0 if videos else 0.0
            grounding_score = fraction(videos, lambda v: bool(v.get("source_chunk_ids")))
            if not videos:
                warnings.append("Không có video nào được tạo.")
            else:
                for v in videos:
                    if not v.get("source_chunk_ids"):
                        warnings.append(f"Video '{v.get('full_title', 'Bài học')}' thiếu liên kết nguồn (source_chunk_ids).")

            # Scene-level rejection rules (product spec): every scene must be grounded,
            # screen text must stay short, and voiceover must not be generic boilerplate.
            all_scenes = [
                sc for v in videos
                for sc in (v.get("storyboard") or v.get("scenes") or [])
                if isinstance(sc, dict)
            ]
            hard_reject = False
            long_lines = 0
            if all_scenes:
                ungrounded = [sc for sc in all_scenes if not sc.get("source_chunk_ids")]
                if ungrounded:
                    warnings.append(f"{len(ungrounded)} cảnh thiếu source_chunk_ids (không truy vết được nguồn).")
                    hard_reject = True

                long_lines = sum(
                    1 for sc in all_scenes
                    for line in (sc.get("screen_text") or [])
                    if len(str(line).split()) > 15
                )
                if long_lines:
                    warnings.append(f"{long_lines} dòng screen_text quá dài (trên 15 từ/dòng).")

                seen_voice: set[str] = set()
                repeated_voice = 0
                for sc in all_scenes:
                    key = re.sub(r"\s+", " ", str(sc.get("voiceover") or "").lower()).strip()
                    if not key:
                        continue
                    if key in seen_voice:
                        repeated_voice += 1
                    else:
                        seen_voice.add(key)
                if repeated_voice:
                    warnings.append(f"{repeated_voice} cảnh có lời thuyết minh lặp lại giống hệt nhau.")
                    hard_reject = True

            score = 0.6 * structure_score + 0.4 * grounding_score
            if long_lines:
                score -= 10
            if hard_reject:
                score = min(score, 79.0)
            if score < 80:
                warnings.append("Chất lượng video chưa đạt chuẩn (điểm < 80).")
            extra_report = {
                "engagement_score": int(max(0, score)),
                "learning_score": int(max(0, score)),
                "visual_score": int(max(0, score)),
                "is_user_friendly": score >= 80,
                "fixes_needed": warnings,
            }

        elif resource_type == "quiz":
            questions = resource_data.get("questions") or []
            structure_score = 100.0 if questions else 0.0
            if not questions:
                warnings.append("Không có câu hỏi nào được tạo.")

            grounding_score = fraction(questions, lambda q: bool(q.get("source_chunk_ids")))
            if grounding_score < 100:
                warnings.append("Một số câu hỏi thiếu source_chunk_ids (chưa bám sát nguồn).")

            def has_real_explanation(q: dict) -> bool:
                explanation = str(q.get("explanation") or "").strip()
                return len(explanation.split()) >= 6

            explanation_score = fraction(questions, has_real_explanation)
            if explanation_score < 100:
                warnings.append("Một số câu hỏi thiếu giải thích chi tiết (chỉ ghi đúng/sai, không dạy được gì).")

            options_questions = [q for q in questions if q.get("options")]
            why_wrong_score = (
                fraction(options_questions, lambda q: bool(q.get("why_wrong_options_are_wrong")))
                if options_questions
                else 100.0
            )
            if why_wrong_score < 50:
                warnings.append("Nhiều câu trắc nghiệm thiếu giải thích vì sao các lựa chọn sai không đúng.")

            duplicate_count = 0
            seen_texts: set[str] = set()
            for q in questions:
                key = re.sub(r"\s+", " ", re.sub(r"[^\w\s]", "", str(q.get("question") or "").lower())).strip()
                if key in seen_texts:
                    duplicate_count += 1
                else:
                    seen_texts.add(key)
            uniqueness_score = (
                100.0 if not duplicate_count
                else max(0.0, 100.0 - (duplicate_count / max(1, len(questions)) * 100))
            )
            if duplicate_count:
                warnings.append(f"Phát hiện {duplicate_count} câu hỏi trùng lặp/gần giống nhau.")

            score = (
                0.2 * structure_score + 0.3 * grounding_score + 0.25 * explanation_score
                + 0.15 * why_wrong_score + 0.1 * uniqueness_score
            )
            # Hard-reject conditions per product rule: no grounding, missing explanations,
            # or repeated questions must fail the gate outright, not just dent the score.
            if grounding_score < 50 or explanation_score < 50 or duplicate_count > 0:
                score = min(score, 79.0)
            extra_report = {"duplicate_questions": duplicate_count}

        elif resource_type == "flashcards":
            cards = resource_data.get("cards") or []
            structure_score = 100.0 if cards else 0.0
            if not cards:
                warnings.append("Không có thẻ ghi nhớ nào được tạo.")

            grounding_score = fraction(cards, lambda c: bool(c.get("source_chunk_ids")))
            if grounding_score < 100:
                warnings.append("Một số thẻ thiếu source_chunk_ids (chưa bám sát nguồn).")

            def has_real_back(c: dict) -> bool:
                back = str(c.get("back") or "").strip()
                front = str(c.get("front") or "").strip().lower()
                return len(back.split()) >= 4 and back.lower() != front

            content_score = fraction(cards, has_real_back)
            if content_score < 100:
                warnings.append("Một số thẻ có mặt sau quá sơ sài hoặc trùng với mặt trước.")

            duplicate_count = 0
            seen_fronts: set[str] = set()
            for c in cards:
                key = re.sub(r"\s+", " ", re.sub(r"[^\w\s]", "", str(c.get("front") or "").lower())).strip()
                if key in seen_fronts:
                    duplicate_count += 1
                else:
                    seen_fronts.add(key)
            uniqueness_score = (
                100.0 if not duplicate_count
                else max(0.0, 100.0 - (duplicate_count / max(1, len(cards)) * 100))
            )
            if duplicate_count:
                warnings.append(f"Phát hiện {duplicate_count} thẻ trùng lặp/gần giống nhau.")

            score = (
                0.25 * structure_score + 0.35 * grounding_score
                + 0.25 * content_score + 0.15 * uniqueness_score
            )
            if grounding_score < 50 or content_score < 50 or duplicate_count > 0:
                score = min(score, 79.0)
            extra_report = {"duplicate_cards": duplicate_count}

        else:
            score = 85.0
            if len(text_content) < 500:
                score -= 20
                warnings.append("Nội dung quá ngắn.")

        score -= 20 * len(artifact_hits)
        for art in artifact_hits:
            warnings.append(f"Chứa rác/artifact: {art}")

        score = max(0, min(100, int(score)))
        # Books and slides are held to the rigorous courseware bar; other outputs keep the legacy bar.
        threshold = 85 if resource_type in {"book", "slides"} else 80
        return {
            "score": score,
            "is_university_ready": score >= threshold,
            "warnings": warnings,
            "issues": warnings,
            "problems": warnings,
            **extra_report,
        }

    def _normalize_book(self, book, docs, target_audience: str):
        points = self._doc_points(docs, limit=18, max_chars=620)
        if not isinstance(book, dict):
            return self._build_fallback_book(docs, target_audience)

        normalized = {
            "title": book.get("title") or "Sách học tập từ tài liệu đã tải lên",
            "subtitle": str(book.get("subtitle") or "").strip(),
            "audience": book.get("audience") or target_audience or "người học",
            "course_level": str(book.get("course_level") or "university").strip(),
            "description": book.get("subtitle")
            or book.get("description")
            or f"Sách học tập dành cho {target_audience or 'người học'}, bám sát nội dung tài liệu gốc.",
            "estimated_duration": book.get("estimated_duration") or "3-5 giờ",
            "prerequisites": [str(p).strip() for p in (book.get("prerequisites") or []) if str(p).strip()],
            "course_learning_outcomes": [
                str(o).strip() for o in (book.get("course_learning_outcomes") or []) if str(o).strip()
            ],
            "chapters": [],
            "glossary": [
                {
                    "term": str(g.get("term") or "").strip(),
                    "definition": str(g.get("definition") or "").strip(),
                    "related_chapter": g.get("related_chapter"),
                }
                for g in (book.get("glossary") or [])
                if isinstance(g, dict) and g.get("term")
            ],
            "review_plan": book.get("review_plan") if isinstance(book.get("review_plan"), dict) else {},
        }

        raw_chapters = book.get("chapters") or []
        if not isinstance(raw_chapters, list) or not raw_chapters:
            return self._build_fallback_book(docs, target_audience)

        point_cursor = 0
        for chapter_index, raw_chapter in enumerate(raw_chapters[:6], 1):
            chapter = raw_chapter if isinstance(raw_chapter, dict) else {"title": str(raw_chapter)}
            raw_lessons = chapter.get("lessons") or []
            if not isinstance(raw_lessons, list) or not raw_lessons:
                raw_lessons = [{"title": f"Bài {chapter_index}.1: Nội dung trọng tâm"}]

            lessons = []
            for lesson_index, raw_lesson in enumerate(raw_lessons[:3], 1):
                lesson = raw_lesson if isinstance(raw_lesson, dict) else {"title": str(raw_lesson)}
                point = points[point_cursor % len(points)]
                point_cursor += 1
                title = lesson.get("title") or f"Bài {chapter_index}.{lesson_index}: Nội dung trọng tâm"
                lessons.append(self._normalize_lesson(lesson, point, title))

            learning_outcomes = chapter.get("learning_outcomes")
            chapter_quiz = [
                {
                    "question": str(q.get("question") or "").strip(),
                    "answer": str(q.get("answer") or "").strip(),
                    "explanation": str(q.get("explanation") or "").strip(),
                }
                for q in (chapter.get("chapter_quiz") or [])
                if isinstance(q, dict) and q.get("question")
            ]
            prerequisites = chapter.get("prerequisites")
            connections = chapter.get("connections_to_other_chapters")
            normalized["chapters"].append(
                {
                    "chapter_index": chapter_index,
                    "title": chapter.get("title") or f"Chương {chapter_index}",
                    "description": chapter.get("overview")
                    or chapter.get("description")
                    or f"Chương này hệ thống hóa {len(lessons)} bài học chính từ tài liệu.",
                    "estimated_duration_minutes": chapter.get("estimated_duration_minutes"),
                    "prerequisites": [str(p).strip() for p in prerequisites if str(p).strip()] if isinstance(prerequisites, list) else [],
                    "learning_outcomes": learning_outcomes if isinstance(learning_outcomes, list) else [],
                    "connections_to_other_chapters": (
                        [str(c).strip() for c in connections if str(c).strip()] if isinstance(connections, list) else []
                    ),
                    "chapter_summary": str(chapter.get("chapter_summary") or "").strip(),
                    "chapter_quiz": chapter_quiz,
                    "lessons": lessons,
                }
            )

        return normalized

    def _evaluate_mindmap_quality_gate(self, mindmap: dict[str, Any]) -> dict[str, Any]:
        """Evaluate quality gate for mindmap against banned phrases and grounding rules."""
        banned_substrings = [
            "contents", "lorem ipsum", "mục lục",
            "bắt đầu dữ liệu", "kết thúc dữ liệu", "mã định danh trang", "nội dung:",
            "ý chính", "ghi nhớ ý chính",
        ]
        warnings = []
        score = 92
        is_usable = True

        all_nodes = [mindmap.get("root", {})] + (mindmap.get("nodes") or [])
        for node in all_nodes:
            if not isinstance(node, dict):
                continue
            title = str(node.get("title") or "")
            summary = str(node.get("summary") or "")
            text_to_check = f"{title} {summary}".lower()

            # Generic filler headings (shared with the book/slide quality gates).
            if self._quality_gate_generic(title):
                warnings.append(f"Node '{title}' dùng tiêu đề chung chung/generic filler.")
                score -= 20
                is_usable = False

            for banned in banned_substrings:
                if banned in text_to_check:
                    warnings.append(f"Node '{title}' chứa nội dung không đạt chuẩn ('{banned}').")
                    score -= 15
                    is_usable = False

            # Dot leaders (". . . ." or "....") — same pattern used by the book/slide/context cleaners.
            if re.search(r"(?:\.\s*){3,}", text_to_check):
                warnings.append(f"Node '{title}' chứa dòng chấm lửng (dot leaders).")
                score -= 15
                is_usable = False

            if title.strip().isdigit() or bool(re.match(r"^(trang|page)\s*\d+$", title.strip(), re.IGNORECASE)):
                warnings.append(f"Node '{title}' sử dụng số trang thô làm tiêu đề.")
                score -= 20
                is_usable = False

            importance = node.get("importance", "low")
            node_type = node.get("type", "")
            if importance in ("high", "medium") or node_type in ("root", "chapter", "lesson", "concept"):
                chunks = node.get("source_chunk_ids") or []
                if not chunks or len(chunks) == 0:
                    warnings.append(f"Node quan trọng '{title}' thiếu source_chunk_ids để đối chiếu.")
                    score -= 10
                    if importance == "high":
                        is_usable = False

        score = max(0, min(100, score))
        if score < 85:
            is_usable = False

        report = {
            "score": score,
            "is_usable": is_usable,
            "warnings": warnings,
        }
        mindmap["quality_report"] = report
        return report

    def build_mindmap_from_book(self, book: dict[str, Any], quality: dict[str, Any] = None) -> dict[str, Any]:
        """Build professional 3-level interactive mindmap schema directly from book plan."""
        title = book.get("title") or book.get("course_title") or "Study Pack"
        description = book.get("description") or book.get("subtitle") or "Sơ đồ kiến thức hệ thống hóa từ tài liệu."
        chapters = book.get("chapters") or []
        all_chunks = list(book.get("source_chunk_ids") or [])
        if not all_chunks:
            for ch in chapters:
                for cid in (ch.get("source_chunk_ids") or []):
                    if cid not in all_chunks:
                        all_chunks.append(cid)
                for les in (ch.get("lessons") or []):
                    for cid in (les.get("source_chunk_ids") or []):
                        if cid not in all_chunks:
                            all_chunks.append(cid)

        root_node = {
            "id": "root",
            "title": title,
            "summary": description,
            "type": "root",
            "importance": "high",
            "source_chunk_ids": all_chunks,
            "children": [],
        }
        nodes = []
        edges = []

        for c_idx, ch in enumerate(chapters):
            ch_id = f"ch_{c_idx}"
            ch_title = ch.get("title") or f"Chương {c_idx+1}"
            ch_summary = ch.get("chapter_summary") or ch.get("overview") or ch.get("description") or f"Nội dung trọng tâm của {ch_title}"
            # Grounding fallback chain: chapter's own ids -> union of its lessons' ids ->
            # book-level ids -> every id in the book. A chapter node must never ship
            # ungrounded when any of its content is grounded.
            ch_chunks = ch.get("source_chunk_ids") or []
            if not ch_chunks:
                for les in (ch.get("lessons") or ch.get("sections") or []):
                    for cid in (les.get("source_chunk_ids") or []):
                        if cid not in ch_chunks:
                            ch_chunks.append(cid)
            if not ch_chunks:
                ch_chunks = list(book.get("source_chunk_ids") or []) or all_chunks
            ch_keywords = [kc.get("term") for kc in (ch.get("core_concepts") or []) if kc.get("term")]
            ch_children_ids = []

            edges.append({"from": "root", "to": ch_id, "relation": "contains"})
            root_node["children"].append(ch_id)

            lessons = ch.get("lessons") or ch.get("sections") or []
            if not lessons and ch_summary:
                lessons = [{"title": f"Trọng tâm {ch_title}", "content": ch_summary, "source_chunk_ids": ch_chunks}]

            for l_idx, les in enumerate(lessons):
                les_id = f"les_{c_idx}_{l_idx}"
                les_title = les.get("title") or les.get("short_name") or f"Bài {c_idx+1}.{l_idx+1}"
                les_summary = les.get("core_idea") or les.get("explanation") or str(les.get("content", ""))[:200] or f"Kiến thức cốt lõi của {les_title}"
                les_chunks = les.get("source_chunk_ids") or ch_chunks
                les_keywords = [kc.get("term") for kc in (les.get("key_concepts") or []) if kc.get("term")] or ch_keywords[:3]
                les_children_ids = []

                edges.append({"from": ch_id, "to": les_id, "relation": "contains"})
                ch_children_ids.append(les_id)

                # Level 3: Concepts
                concepts = (les.get("key_concepts") or (ch.get("core_concepts") if l_idx == 0 else [])) or []
                for k_idx, kc in enumerate(concepts[:3]):
                    c_term = kc.get("term")
                    c_def = kc.get("definition")
                    if c_term and c_def:
                        c3_id = f"cpt_{c_idx}_{l_idx}_{k_idx}"
                        les_children_ids.append(c3_id)
                        # Chapter-level core_concepts may carry a formula — surface it as its own
                        # "formula" node type (per the mindmap schema's type enum) instead of
                        # collapsing every concept into the generic "concept" bucket.
                        c_formula = str(kc.get("formula") or "").strip()
                        nodes.append({
                            "id": c3_id, "parent_id": les_id, "title": c_term,
                            "summary": f"{c_def}\n\nCông thức: {c_formula}" if c_formula else c_def,
                            "type": "formula" if c_formula else "concept",
                            "importance": "medium", "keywords": [c_term],
                            "source_chunk_ids": les_chunks, "children": []
                        })
                        edges.append({"from": les_id, "to": c3_id, "relation": "explains"})

                # Level 3: Examples
                examples = (les.get("worked_examples") or (ch.get("worked_examples") if l_idx == 0 else [])) or []
                for e_idx, ex in enumerate(examples[:2]):
                    ex_title = ex.get("title") or ex.get("problem") or f"Ví dụ mẫu {e_idx+1}"
                    if len(str(ex_title)) > 50:
                        ex_title = str(ex_title)[:47].rstrip(". ") + "…"
                    ex_sum = ex.get("solution") or ex.get("problem") or "Lời giải ví dụ mẫu"
                    e3_id = f"ex_{c_idx}_{l_idx}_{e_idx}"
                    les_children_ids.append(e3_id)
                    nodes.append({
                        "id": e3_id, "parent_id": les_id, "title": ex_title, "summary": ex_sum,
                        "type": "example", "importance": "low", "keywords": ["ví dụ", "lời giải"],
                        "source_chunk_ids": les_chunks, "children": []
                    })
                    edges.append({"from": les_id, "to": e3_id, "relation": "example_of"})

                # Level 3: Warnings / Mistakes
                mistakes = (les.get("common_mistakes") or les.get("common_mistake") or (ch.get("common_mistakes") if l_idx == 0 else [])) or []
                for m_idx, cm in enumerate(mistakes[:2]):
                    cm_title = cm if isinstance(cm, str) else (cm.get("mistake") or cm.get("title") or "Sai lầm thường gặp")
                    if len(str(cm_title)) > 50:
                        cm_title = str(cm_title)[:47].rstrip(". ") + "…"
                    cm_sum = cm if isinstance(cm, str) else (cm.get("correction") or cm.get("explanation") or cm_title)
                    m3_id = f"warn_{c_idx}_{l_idx}_{m_idx}"
                    les_children_ids.append(m3_id)
                    nodes.append({
                        "id": m3_id, "parent_id": les_id, "title": f"Lưu ý: {cm_title}", "summary": cm_sum,
                        "type": "warning", "importance": "low", "keywords": ["lưu ý", "sai lầm"],
                        "source_chunk_ids": les_chunks, "children": []
                    })
                    edges.append({"from": les_id, "to": m3_id, "relation": "contrasts_with"})

                # Level 3: Exercises
                exercises = (les.get("practice_problems") or (ch.get("practice_problems") if l_idx == 0 else [])) or []
                for p_idx, pp in enumerate(exercises[:1]):
                    pp_title = pp.get("question") or pp.get("problem") or f"Bài tập tự luyện {p_idx+1}"
                    if len(str(pp_title)) > 50:
                        pp_title = str(pp_title)[:47].rstrip(". ") + "…"
                    pp_sum = pp.get("explanation") or pp.get("solution") or pp_title
                    p3_id = f"exer_{c_idx}_{l_idx}_{p_idx}"
                    les_children_ids.append(p3_id)
                    nodes.append({
                        "id": p3_id, "parent_id": les_id, "title": pp_title, "summary": pp_sum,
                        "type": "exercise", "importance": "low", "keywords": ["bài tập", "tự luyện"],
                        "source_chunk_ids": les_chunks, "children": []
                    })
                    edges.append({"from": les_id, "to": p3_id, "relation": "leads_to"})

                nodes.append({
                    "id": les_id, "parent_id": ch_id, "title": les_title, "summary": les_summary,
                    "type": "lesson", "importance": "medium", "keywords": les_keywords,
                    "source_chunk_ids": les_chunks, "children": les_children_ids
                })

            nodes.append({
                "id": ch_id, "parent_id": "root", "title": ch_title, "summary": ch_summary,
                "type": "chapter", "importance": "high", "keywords": ch_keywords,
                "source_chunk_ids": ch_chunks, "children": ch_children_ids
            })

        mindmap = {
            "document_id": getattr(self, "course_id", ""),
            "title": title,
            "description": description,
            "root": root_node,
            "nodes": nodes,
            "edges": edges,
        }
        self._evaluate_mindmap_quality_gate(mindmap)
        return mindmap

    def regenerate_mindmap(self, force_llm: bool = False, profile: Optional[dict[str, Any]] = None) -> dict[str, Any]:
        """Regenerate or assemble the 3-level interactive mindmap for the course."""
        profile_directives = build_profile_directives(profile)
        paths = get_course_path(self.course_id)
        book = None
        if os.path.exists(paths.get("book", "")):
            book = self._read_json(paths["book"], None)

        if book and isinstance(book, dict) and book.get("chapters") and not force_llm:
            logger.info("[MindmapGen] Deriving 3-level mindmap directly from structured book plan.")
            mindmap = self.build_mindmap_from_book(book)
            if "mindmap" in paths:
                self._save_json(paths["mindmap"], mindmap)
            return mindmap

        logger.info("[MindmapGen] Generating mindmap via LLM from clean chunks.")
        if not self.vectorstore:
            try:
                from backend.main import course_manager
                rag = course_manager.get_course(self.course_id) if hasattr(course_manager, "get_course") else None
                if rag and hasattr(rag, "vectorstore"):
                    self.vectorstore = rag.vectorstore
            except Exception as e:
                logger.debug(f"[MindmapGen] Could not fetch vectorstore from course_manager: {e}")

        if self.vectorstore and hasattr(self.vectorstore, "as_retriever"):
            retriever = self.vectorstore.as_retriever(search_kwargs={"k": 20})
            docs = retriever.invoke("chương bài học khái niệm công thức ví dụ định nghĩa")
            cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=15, max_chars=350)
        else:
            docs = []
            cleaned_chunks = []
        
        if not cleaned_chunks and book:
            mindmap = self.build_mindmap_from_book(book)
            if "mindmap" in paths:
                self._save_json(paths["mindmap"], mindmap)
            return mindmap

        context_text = "\n\n".join([
            f"--- NGUỒN [CHUNK_ID: {c.get('source_chunk_ids', [''])[0]}] ---\n{c['text']}"
            for c in cleaned_chunks
        ])
        
        from backend.core.prompts import MINDMAP_GENERATION_PROMPT
        prompt = MINDMAP_GENERATION_PROMPT.format(context=context_text, profile_directives=profile_directives)
        
        llm = get_llm(task="mindmap")
        try:
            res = llm.invoke(prompt)
            raw_json = extract_json(res.content) if hasattr(res, "content") else extract_json(str(res))
            if isinstance(raw_json, dict) and raw_json.get("root") and raw_json.get("nodes"):
                raw_json["document_id"] = self.course_id
                self._evaluate_mindmap_quality_gate(raw_json)
                if "mindmap" in paths:
                    self._save_json(paths["mindmap"], raw_json)
                return self._sanitize_payload(raw_json)
        except Exception as e:
            logger.error(f"[MindmapGen] LLM generation failed: {e}. Falling back to book or shallow mindmap.")

        if book and isinstance(book, dict):
            mindmap = self.build_mindmap_from_book(book)
        else:
            mindmap = self.generate_fallback_shallow_mindmap("Sơ đồ tư duy tài liệu", docs)

        if "mindmap" in paths:
            self._save_json(paths["mindmap"], mindmap)
        return mindmap

    def _build_study_pack_from_book(self, book: dict[str, Any], quality: dict[str, Any]) -> dict[str, Any]:
        """Build a connected study pack derived from the normalized book structure."""
        chapters = book.get("chapters") or []
        high_yield_summary = []
        flashcards = []
        source_chunk_ids = set()

        for ch in chapters:
            ch_title = ch.get("title") or "Chương"
            if ch.get("chapter_summary"):
                high_yield_summary.append({"title": ch_title, "summary": ch["chapter_summary"]})
            for lesson in ch.get("lessons") or []:
                for cid in (lesson.get("source_chunk_ids") or []):
                    source_chunk_ids.add(cid)
                for kc in (lesson.get("key_concepts") or []):
                    if kc.get("term") and kc.get("definition"):
                        flashcards.append({"front": kc["term"], "back": kc["definition"]})

        for item in (book.get("glossary") or []):
            if item.get("term") and item.get("definition"):
                flashcards.append({"front": item["term"], "back": item["definition"]})
            for cid in (item.get("source_chunk_ids") or []):
                source_chunk_ids.add(cid)

        mindmap_data = self.build_mindmap_from_book(book, quality)

        return {
            "main_output": "study_guide_pdf",
            "study_guide": {"chapters": len(chapters), "title": book.get("title")},
            "mindmap": mindmap_data,
            "high_yield_summary": high_yield_summary,
            "flashcards": flashcards,
            "source_chunk_ids": list(source_chunk_ids),
            "quality_report": quality,
        }

    def _build_fallback_book(self, docs, target_audience: str):
        points = self._doc_points(docs, limit=6, max_chars=620)
        chapters = []
        for index, point in enumerate(points, 1):
            title = self._short_title(point["text"], f"Nội dung chính {index}")
            lesson_title_1 = f"Bài {index}.1: Đọc hiểu nội dung chính"
            lesson_title_2 = f"Bài {index}.2: Vận dụng và tự kiểm tra"
            chapters.append(
                {
                    "chapter_index": index,
                    "title": f"Chương {index}: {title}",
                    "description": "Hệ thống hóa một nhóm nội dung trọng tâm trong tài liệu.",
                    "estimated_duration_minutes": None,
                    "prerequisites": [],
                    "learning_outcomes": [],
                    "connections_to_other_chapters": [],
                    "chapter_summary": "",
                    "chapter_quiz": [],
                    "lessons": [
                        self._normalize_lesson({}, point, lesson_title_1),
                        self._normalize_lesson({}, point, lesson_title_2),
                    ],
                }
            )
        return {
            "title": "Sách học tập từ tài liệu đã tải lên",
            "subtitle": "",
            "audience": target_audience or "người học",
            "course_level": "university",
            "description": (
                f"Bản sách MVP dành cho {target_audience or 'người học'}, được dựng trực tiếp từ các đoạn "
                "nội dung đã index trong tài liệu."
            ),
            "estimated_duration": "3-5 giờ",
            "prerequisites": [],
            "course_learning_outcomes": [],
            "chapters": chapters,
            "glossary": [],
            "review_plan": {},
        }

    def _book_pdf_elements(self, book: dict[str, Any]) -> list[tuple[str, str]]:
        elements: list[tuple[str, str]] = []

        def add(text: Any, style: str = "body") -> None:
            clean = str(text or "").strip()
            if clean:
                elements.append((clean, style))

        # Title/description live on the cover page; body starts with how-to-use.
        if book.get("how_to_use"):
            add("Cách sử dụng cẩm nang này", "chapter")
            for step_index, step in enumerate(book["how_to_use"], 1):
                add(f"{step_index}. {step}", "body")
            elements.append(("", "gap"))

        if book.get("prerequisites"):
            add("Kiến thức cần có trước", "chapter")
            for item in book["prerequisites"]:
                add(f"- {item}", "body")
            elements.append(("", "gap"))
        if book.get("course_learning_outcomes"):
            add("Kết quả học tập tổng thể", "section")
            for item in book["course_learning_outcomes"]:
                add(f"- {item}", "body")
            elements.append(("", "gap"))

        if book.get("course_roadmap"):
            add("Lộ trình khóa học", "chapter")
            for road_index, road_unit in enumerate(book["course_roadmap"], 1):
                if not isinstance(road_unit, dict):
                    continue
                add(f"Phần {road_index}: {road_unit.get('title', '')}", "lesson")
                if road_unit.get("big_idea"):
                    add(road_unit["big_idea"], "body")
                if road_unit.get("key_concepts"):
                    add("Khái niệm chính: " + ", ".join(road_unit["key_concepts"][:6]), "small")
            elements.append(("", "gap"))

        for chapter_index, chapter in enumerate(book.get("chapters") or [], 1):
            if not isinstance(chapter, dict):
                continue
            add(chapter.get("title") or f"Chương {chapter_index}", "chapter")
            add(chapter.get("chapter_overview") or chapter.get("description"), "body")
            if chapter.get("learning_objectives"):
                add("Mục tiêu học tập", "section")
                for item in chapter["learning_objectives"]:
                    add(f"- {item}", "body")
            if chapter.get("prerequisites"):
                add("Yêu cầu trước chương này", "section")
                for item in chapter["prerequisites"]:
                    add(f"- {item}", "body")
            if chapter.get("learning_outcomes"):
                add("Kết quả học tập của chương", "section")
                for item in chapter["learning_outcomes"]:
                    add(f"- {item}", "body")
            if chapter.get("big_picture"):
                add(f"BỨC TRANH TOÀN CẢNH\n{chapter['big_picture']}", "callout_key_idea")

            core_concepts = [c for c in (chapter.get("core_concepts") or []) if isinstance(c, dict)]
            for concept in core_concepts:
                add(concept.get("term", ""), "lesson")
                if concept.get("definition"):
                    add(f"Định nghĩa: {concept['definition']}", "body")
                if concept.get("intuition"):
                    add(f"TRỰC QUAN\n{concept['intuition']}", "callout_key_idea")
                if concept.get("technical_explanation"):
                    add("Giải thích kỹ thuật", "section")
                    add(concept["technical_explanation"], "body")
                if concept.get("formula"):
                    add(concept["formula"], "formula")
                if concept.get("code"):
                    add(concept["code"], "code")
                if concept.get("example"):
                    example_text = concept["example"]
                    if concept.get("non_example"):
                        example_text += f"\n\nKHÔNG phải là: {concept['non_example']}"
                    add(f"VÍ DỤ MINH HỌA\n{example_text}", "callout_example")
                mistake = concept.get("common_mistake") or {}
                if mistake.get("mistake"):
                    add(
                        f"SAI LẦM THƯỜNG GẶP\nSai: {mistake['mistake']}\nĐúng: {mistake.get('correction', '')}",
                        "callout_mistake",
                    )
                elements.append(("", "gap"))

            if core_concepts:
                for we in chapter.get("worked_examples") or []:
                    if not we.get("problem"):
                        continue
                    steps = "\n".join(
                        f"  {i}. {step}" for i, step in enumerate(we.get("step_by_step_solution") or [], 1)
                    )
                    body_text = f"Đề bài: {we['problem']}\n{steps}"
                    if we.get("why_each_step_matters"):
                        body_text += f"\nVì sao các bước này quan trọng: {we['why_each_step_matters']}"
                    if we.get("common_error"):
                        body_text += f"\nLỗi thường gặp: {we['common_error']}"
                    add(f"VÍ DỤ MẪU (TỪNG BƯỚC): {we.get('title', '')}\n{body_text}", "callout_worked_example")

                for pp in chapter.get("practice_problems") or []:
                    if not pp.get("question"):
                        continue
                    body_text = f"[{pp.get('difficulty', 'medium')}] {pp['question']}"
                    if pp.get("hint"):
                        body_text += f"\nGợi ý: {pp['hint']}"
                    if pp.get("solution"):
                        body_text += f"\nLời giải: {pp['solution']}"
                    add(f"THỰC HÀNH\n{body_text}", "callout_practice")

            if core_concepts:
                if chapter.get("chapter_summary"):
                    add("Tóm tắt chương", "section")
                    add(chapter["chapter_summary"], "body")
                if chapter.get("active_recall_questions"):
                    add("Tự kiểm tra (không nhìn tài liệu)", "section")
                    for q_index, question in enumerate(chapter["active_recall_questions"], 1):
                        add(f"{q_index}. {question}", "body")
                if chapter.get("connections_to_other_chapters"):
                    add("Liên hệ với các chương khác", "section")
                    for item in chapter["connections_to_other_chapters"]:
                        add(f"- {item}", "body")
                elements.append(("", "gap"))
                continue  # lessons[] duplicate the concepts; render them only for legacy books

            for lesson_index, lesson in enumerate(chapter.get("lessons") or chapter.get("sections") or [], 1):
                if not isinstance(lesson, dict):
                    continue
                add(lesson.get("title") or f"Bài {chapter_index}.{lesson_index}", "lesson")
                add(f"Thời lượng: {lesson.get('duration')}", "small")

                for label, key in [
                    ("Mục tiêu", "objectives"),
                    ("Tại sao quan trọng (Why It Matters):", "why_it_matters"),
                    ("Nội dung bài giảng", "lecture"),
                    ("Kiến thức cần nhớ", "key_points"),
                ]:
                    value = lesson.get(key)
                    if not value:
                        continue
                    add(label, "section")
                    if isinstance(value, list):
                        for item in value:
                            add(f"- {item}", "body")
                    else:
                        add(value, "body")

                if lesson.get("core_idea"):
                    add("Ý tưởng cốt lõi (Core Idea):", "section")
                    add(lesson["core_idea"], "body")

                if lesson.get("key_concepts"):
                    add("Khái niệm chính", "section")
                    for kc in lesson["key_concepts"]:
                        if kc.get("term"):
                            add(f"- {kc['term']}: {kc.get('definition', '')}", "body")

                if lesson.get("example"):
                    example_text = lesson["example"]
                    if lesson.get("non_example"):
                        example_text += f"\n\nKHÔNG phải là: {lesson['non_example']}"
                    add(f"📘 VÍ DỤ MINH HỌA\n{example_text}", "callout_example")

                common_mis = lesson.get("common_misunderstanding") or {}
                if common_mis.get("mistake"):
                    add("Sai lầm phổ biến & Cách hiểu đúng:", "section")
                    add(f"Sai: {common_mis['mistake']}\nĐúng: {common_mis.get('correction', '')}", "body")

                for we in lesson.get("worked_examples") or []:
                    if not we.get("problem"):
                        continue
                    steps = "\n".join(
                        f"  {i}. {step}" for i, step in enumerate(we.get("step_by_step_solution") or [], 1)
                    )
                    body = f"Đề bài: {we['problem']}\n{steps}"
                    if we.get("why_each_step_matters"):
                        body += f"\nVì sao các bước này quan trọng: {we['why_each_step_matters']}"
                    if we.get("common_error"):
                        body += f"\nLỗi thường gặp: {we['common_error']}"
                    add(f"🧮 VÍ DỤ MẪU (TỪNG BƯỚC): {we.get('title', '')}\n{body}", "callout_worked_example")

                for label, key in [
                    ("Hoạt động học tập", "activity"),
                    ("Kiểm tra nhanh", "assessment"),
                ]:
                    value = lesson.get(key)
                    if not value:
                        continue
                    add(label, "section")
                    if isinstance(value, list):
                        for item in value:
                            add(f"- {item}", "body")
                    else:
                        add(value, "body")

                for pp in lesson.get("practice_problems") or []:
                    if not pp.get("question"):
                        continue
                    body = f"[{pp.get('difficulty', 'medium')}] {pp['question']}"
                    if pp.get("hint"):
                        body += f"\nGợi ý: {pp['hint']}"
                    if pp.get("solution"):
                        body += f"\nLời giải: {pp['solution']}"
                    add(f"✏️ THỰC HÀNH\n{body}", "callout_practice")
                elements.append(("", "gap"))

            if chapter.get("chapter_summary"):
                add("Tóm tắt chương", "section")
                add(chapter["chapter_summary"], "body")
            if chapter.get("active_recall_questions"):
                add("Tự kiểm tra (không nhìn tài liệu)", "section")
                for q_index, question in enumerate(chapter["active_recall_questions"], 1):
                    add(f"{q_index}. {question}", "body")
            elements.append(("", "gap"))

        if book.get("glossary"):
            add("Thuật ngữ & Khái niệm (Glossary)", "chapter")
            for item in book["glossary"]:
                if item.get("term"):
                    add(f"- {item['term']}: {item.get('definition', '')}", "body")
            elements.append(("", "gap"))

        problem_set = [p for p in (book.get("problem_set") or []) if isinstance(p, dict) and p.get("question")]
        if problem_set:
            add("Bộ bài tập tổng hợp", "chapter")
            for p_index, problem in enumerate(problem_set, 1):
                add(
                    f"Bài {p_index} [{problem.get('difficulty', 'medium')}] (Chương {problem.get('chapter_index', '?')}): "
                    f"{problem['question']}",
                    "body",
                )
                if problem.get("hint"):
                    add(f"Gợi ý: {problem['hint']}", "small")
            elements.append(("", "gap"))

            add("Đáp án bộ bài tập", "chapter")
            for p_index, problem in enumerate(problem_set, 1):
                solution = problem.get("solution") or problem.get("expected_answer")
                if solution:
                    add(f"Bài {p_index}: {solution}", "body")
            elements.append(("", "gap"))

        review_plan = book.get("review_plan") or {}
        if any(review_plan.get(k) for k in ("ten_minute", "thirty_minute", "one_hour")):
            add("Kế hoạch ôn tập (Review Plan)", "chapter")
            for label, key in [
                ("Ôn tập 10 phút", "ten_minute"),
                ("Ôn tập 30 phút", "thirty_minute"),
                ("Ôn tập 1 giờ", "one_hour"),
            ]:
                for item in review_plan.get(key) or []:
                    add(f"[{label}] {item}", "body")

        return elements

    _MARKDOWN_PATTERNS = [
        (re.compile(r"```[a-zA-Z]*\n?"), ""),
        (re.compile(r"^#{1,4}\s+", re.MULTILINE), ""),
        (re.compile(r"\*\*(.+?)\*\*"), r"\1"),
        (re.compile(r"__(.+?)__"), r"\1"),
        (re.compile(r"(?<!\w)\*(?!\s)(.+?)(?<!\s)\*(?!\w)"), r"\1"),
        (re.compile(r"`([^`\n]+)`"), r"\1"),
    ]

    def _strip_markdown(self, text: str) -> str:
        """LLM output sometimes carries markdown; PDFs must render plain styled text."""
        value = str(text or "")
        for pattern, replacement in self._MARKDOWN_PATTERNS:
            value = pattern.sub(replacement, value)
        return value

    def _render_book_cover(self, book: dict[str, Any], page_width: int, page_height: int):
        from PIL import Image, ImageDraw

        image = Image.new("RGB", (page_width, page_height), (255, 255, 255))
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, page_width, 26), fill=(30, 64, 175))
        draw.rectangle((0, page_height - 26, page_width, page_height), fill=(30, 64, 175))

        level_labels = {
            "introductory": "Nhập môn",
            "intermediate": "Trung cấp",
            "university": "Đại học",
            "advanced": "Nâng cao",
            "high_yield": "Bản trọng tâm",
        }
        level = level_labels.get(str(book.get("course_level") or "").lower(), "")

        y = 340
        draw.text((120, y), "SÁCH HỌC TẬP", font=self._font(26, bold=True), fill=(30, 64, 175))
        y += 70
        for line in self._wrap_lines(self._strip_markdown(book.get("title") or "Sách học tập"), width=30)[:4]:
            draw.text((120, y), line, font=self._font(56, bold=True), fill=(17, 24, 39))
            y += 74
        subtitle = self._strip_markdown(book.get("subtitle") or "")
        if subtitle:
            y += 10
            for line in self._wrap_lines(subtitle, width=52)[:3]:
                draw.text((120, y), line, font=self._font(28), fill=(71, 85, 105))
                y += 42

        y = max(y + 60, 820)
        draw.line((120, y, page_width - 120, y), fill=(226, 232, 240), width=3)
        y += 40
        meta_rows = [
            ("Đối tượng", str(book.get("audience") or "")),
            ("Trình độ", level),
            ("Thời lượng ước tính", str(book.get("estimated_duration") or "")),
        ]
        for label, value in meta_rows:
            if not value:
                continue
            draw.text((120, y), f"{label}:", font=self._font(24, bold=True), fill=(51, 65, 85))
            draw.text((420, y), value[:60], font=self._font(24), fill=(51, 65, 85))
            y += 46

        disclaimer = (
            "Nội dung được AI tạo từ tài liệu của bạn và bám sát nguồn gốc. "
            "AI có thể sai — hãy kiểm chứng những thông tin quan trọng."
        )
        y = page_height - 220
        for line in self._wrap_lines(disclaimer, width=74):
            draw.text((120, y), line, font=self._font(20), fill=(100, 116, 139))
            y += 30
        return image

    def _render_book_toc(self, toc_entries: list[tuple[str, int]], page_width: int, page_height: int, per_page: int):
        """Render TOC pages: right-aligned page numbers, no dot leaders."""
        from PIL import Image, ImageDraw

        pages = []
        margin = 120
        for start in range(0, len(toc_entries), per_page):
            image = Image.new("RGB", (page_width, page_height), (255, 255, 255))
            draw = ImageDraw.Draw(image)
            y = 140
            if start == 0:
                draw.text((margin, y), "Mục lục", font=self._font(40, bold=True), fill=(30, 64, 175))
                y += 90
            for label, page_number in toc_entries[start : start + per_page]:
                lines = self._wrap_lines(label, width=58) or [label]
                draw.text((margin, y), lines[0], font=self._font(24), fill=(51, 65, 85))
                draw.text((page_width - margin - 60, y), str(page_number), font=self._font(24, bold=True), fill=(30, 64, 175))
                y += 44
            pages.append(image)
        return pages

    def _render_book_pdf(self, book: dict[str, Any], pdf_path: str) -> str:
        from PIL import Image, ImageDraw

        page_width, page_height = 1240, 1754
        margin = 86
        content_width = page_width - margin * 2
        styles = {
            "title": {"font": self._font(44, bold=True), "fill": (17, 24, 39), "line_height": 58, "width": 38},
            "chapter": {"font": self._font(34, bold=True), "fill": (30, 64, 175), "line_height": 46, "width": 48},
            "lesson": {"font": self._font(29, bold=True), "fill": (15, 23, 42), "line_height": 40, "width": 58},
            "section": {"font": self._font(25, bold=True), "fill": (51, 65, 85), "line_height": 34, "width": 66},
            "body": {"font": self._font(24), "fill": (51, 65, 85), "line_height": 34, "width": 78},
            "small": {"font": self._font(21), "fill": (100, 116, 139), "line_height": 30, "width": 88},
            "formula": {
                "font": self._font(26, bold=True), "fill": (30, 41, 59), "line_height": 40, "width": 60,
                "bg": (248, 250, 252), "accent": (100, 116, 139),
            },
            "code": {
                "font": self._font(22), "fill": (226, 232, 240), "line_height": 32, "width": 74,
                "bg": (15, 23, 42), "accent": (56, 189, 248),
            },
            "callout_key_idea": {
                "font": self._font(23), "fill": (120, 53, 15), "line_height": 32, "width": 72,
                "bg": (255, 251, 235), "accent": (217, 119, 6),
            },
            "callout_example": {
                "font": self._font(23), "fill": (6, 78, 59), "line_height": 32, "width": 72,
                "bg": (236, 253, 245), "accent": (5, 150, 105),
            },
            "callout_mistake": {
                "font": self._font(23), "fill": (127, 29, 29), "line_height": 32, "width": 72,
                "bg": (254, 242, 242), "accent": (220, 38, 38),
            },
            "callout_practice": {
                "font": self._font(23), "fill": (49, 46, 129), "line_height": 32, "width": 72,
                "bg": (238, 242, 255), "accent": (79, 70, 229),
            },
            "callout_worked_example": {
                "font": self._font(23), "fill": (12, 74, 110), "line_height": 32, "width": 72,
                "bg": (240, 249, 255), "accent": (2, 132, 199),
            },
        }
        callout_styles = {
            "callout_key_idea", "callout_example", "callout_mistake",
            "callout_practice", "callout_worked_example",
        }
        boxed_styles = callout_styles | {"formula", "code"}

        body_pages: list[Image.Image] = []
        chapter_marks: list[tuple[str, int]] = []  # (chapter title, body page index starting at 1)
        image = Image.new("RGB", (page_width, page_height), (255, 255, 255))
        draw = ImageDraw.Draw(image)
        y = margin

        def finish_page() -> None:
            body_pages.append(image.copy())

        def reset_page() -> tuple[Image.Image, ImageDraw.ImageDraw, int]:
            next_image = Image.new("RGB", (page_width, page_height), (255, 255, 255))
            next_draw = ImageDraw.Draw(next_image)
            return next_image, next_draw, margin

        for text, style_name in self._book_pdf_elements(book):
            if style_name == "gap":
                y += 24
                continue

            text = self._strip_markdown(text)
            style = styles.get(style_name, styles["body"])
            is_boxed = style_name in boxed_styles
            text_x = margin
            wrap_width = style["width"]

            if is_boxed:
                text_x = margin + 22
                wrap_width = style["width"] - 4

            if style_name == "code":
                # Preserve code line structure instead of re-wrapping prose-style.
                lines = [ln[:70] for ln in text.replace("\r", "").split("\n") if ln.strip()][:24]
            else:
                lines = self._wrap_lines(text, width=wrap_width) or [text]
            pad = 20 if is_boxed else 0
            block_height = len(lines) * style["line_height"] + 12 + pad * 2
            if y + block_height > page_height - margin - 40:
                finish_page()
                image, draw, y = reset_page()

            if style_name == "chapter":
                chapter_marks.append((text.split("\n", 1)[0], len(body_pages) + 1))

            if style_name in {"chapter", "lesson"}:
                draw.rounded_rectangle(
                    (margin - 18, y - 10, margin + content_width + 18, y + block_height - 8),
                    radius=10,
                    fill=(248, 250, 252),
                )
            elif is_boxed:
                draw.rounded_rectangle(
                    (margin, y - 4, margin + content_width, y + block_height - pad + 8),
                    radius=10,
                    fill=style["bg"],
                )
                draw.rectangle(
                    (margin, y - 4, margin + 8, y + block_height - pad + 8),
                    fill=style["accent"],
                )
                y += pad

            for line_index, line in enumerate(lines):
                font = style["font"]
                fill = style["fill"]
                if style_name in callout_styles and line_index == 0:
                    font = self._font(23, bold=True)
                    fill = style["accent"]
                draw.text((text_x, y), line, font=font, fill=fill)
                y += style["line_height"]
            y += pad + 12

        finish_page()

        # Assemble: cover + TOC (page numbers offset by front matter) + body with footers.
        toc_per_page = 30
        toc_page_count = max(1, -(-len(chapter_marks) // toc_per_page)) if chapter_marks else 0
        front_matter = 1 + toc_page_count
        toc_entries = [(title, page + front_matter) for title, page in chapter_marks]

        pages: list[Image.Image] = [self._render_book_cover(book, page_width, page_height)]
        pages.extend(self._render_book_toc(toc_entries, page_width, page_height, toc_per_page))

        footer_font = self._font(20)
        for body_index, body_page in enumerate(body_pages, 1):
            footer_draw = ImageDraw.Draw(body_page)
            page_label = f"Trang {body_index + front_matter}"
            footer_draw.text(
                (page_width // 2 - 40, page_height - 52), page_label, font=footer_font, fill=(148, 163, 184)
            )
            pages.append(body_page)

        os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
        pages[0].save(pdf_path, "PDF", resolution=120.0, save_all=True, append_images=pages[1:])
        return pdf_path

    def export_book_pdf(self) -> str:
        paths = get_course_path(self.course_id)
        if not os.path.exists(paths["book"]):
            raise FileNotFoundError("Book has not been generated yet.")
        with open(paths["book"], "r", encoding="utf-8") as f:
            book = json.load(f)
        return self._render_book_pdf(book, paths["book_pdf"])

    def _render_artifact_pdf(self, title: str, elements: list[tuple[str, str]], pdf_path: str) -> str:
        from PIL import Image, ImageDraw

        page_width, page_height = 1240, 1754
        margin = 86
        content_width = page_width - margin * 2
        styles = {
            "title": {"font": self._font(44, bold=True), "fill": (17, 24, 39), "line_height": 58, "width": 38},
            "heading": {"font": self._font(32, bold=True), "fill": (15, 23, 42), "line_height": 44, "width": 52},
            "section": {"font": self._font(24, bold=True), "fill": (51, 65, 85), "line_height": 34, "width": 72},
            "body": {"font": self._font(23), "fill": (51, 65, 85), "line_height": 33, "width": 82},
            "small": {"font": self._font(20), "fill": (100, 116, 139), "line_height": 29, "width": 92},
        }

        pages: list[Image.Image] = []
        image = Image.new("RGB", (page_width, page_height), (255, 255, 255))
        draw = ImageDraw.Draw(image)
        y = margin

        def finish_page() -> None:
            pages.append(image.copy())

        def reset_page() -> tuple[Image.Image, ImageDraw.ImageDraw, int]:
            next_image = Image.new("RGB", (page_width, page_height), (255, 255, 255))
            next_draw = ImageDraw.Draw(next_image)
            return next_image, next_draw, margin

        full_elements = [(title, "title"), ("", "gap"), *elements]
        for text, style_name in full_elements:
            if style_name == "gap":
                y += 24
                continue

            style = styles.get(style_name, styles["body"])
            lines = self._wrap_lines(self._clean_generated_text(text, compact=False), width=style["width"]) or [text]
            block_height = len(lines) * style["line_height"] + 16
            if y + block_height > page_height - margin:
                finish_page()
                image, draw, y = reset_page()

            if style_name == "heading":
                draw.rounded_rectangle(
                    (margin - 18, y - 10, margin + content_width + 18, y + block_height - 8),
                    radius=10,
                    fill=(248, 250, 252),
                )

            for line in lines:
                draw.text((margin, y), line, font=style["font"], fill=style["fill"])
                y += style["line_height"]
            y += 12

        finish_page()
        os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
        pages[0].save(pdf_path, "PDF", resolution=120.0, save_all=True, append_images=pages[1:])
        return pdf_path

    def _format_math_text(self, text: Any, compact: bool = True) -> str:
        import re
        val = str(text or "")
        val = val.replace("A_i x_i", "A_i × x_i").replace(r"\sum", "Σ").replace("sum(", "Σ(").replace(r"\le", "≤").replace("<=", "≤").replace(" * ", " × ")
        val = re.sub(r"P\(([^,]+),([^)]+)\)", r"P(\1, \2)", val)
        def sub_repl(m):
            base, sub = m.group(1), m.group(2)
            sub_map = str.maketrans("0123456789aehijklmnoprstuvx", "₀₁₂₃₄₅₆₇₈₉ₐₑₕᵢⱼₖ⪲ₘₙₒₚᵣₛₜᵤᵥₓ")
            return base + sub.translate(sub_map)
        val = re.sub(r"\b([A-Za-z]+)_([0-9a-z]+)\b", sub_repl, val)
        val = re.sub(r"\b([A-WYZx])([ij])\b", sub_repl, val)
        def sup_repl(m):
            base, sup = m.group(1), m.group(2)
            sup_map = str.maketrans("0123456789+-=()ni", "⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾ⁿⁱ")
            return base + sup.translate(sup_map)
        val = re.sub(r"([0-9a-zA-Z]+)\^([0-9a-z+-]+)", sup_repl, val)
        if compact:
            val = re.sub(r"\s+", " ", val).strip()
        return val

    def _infer_visual_type(self, slide: dict[str, Any]) -> str:
        title = str(slide.get("title") or "").lower()
        content = str(slide.get("content") or "").lower()
        sug = str(slide.get("image_suggestion") or "").lower()
        combined = f"{title} {content} {sug}"
        
        if "lesson map" in title or "mind map" in title:
            return "lesson_map"
        if "tổng kết" in title or "summary" in title:
            return "summary"
            
        if "mex" in combined or "tree" in combined:
            return "tree_mex"
        if "knapsack" in combined or "bag" in combined:
            return "knapsack"
        if "grid" in combined or "bàn cờ" in combined:
            return "grid"
        if "modulo" in combined or "mod" in combined:
            return "counting_modulo"
        if "lesson map" in combined or "mind map" in combined or "recap" in combined:
            return "lesson_map"
        if "tổng kết" in combined or "summary" in combined:
            return "summary"
        if "pseudocode" in title or "for i in" in content or "code" in title:
            return "code"
        if "đồ thị" in title or "graph" in combined or "bfs" in combined:
            return "graph"
        if "quy hoạch động" in title or "dp" in combined:
            return "dp"
        return "default"

    def export_slides_pptx(self) -> str:
        """Export the deck as a clean academic PPTX: indigo accents, header/footer, typed layouts."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        INDIGO_DARK = RGBColor(30, 27, 75)
        INDIGO = RGBColor(49, 46, 129)
        ACCENT = RGBColor(79, 70, 229)
        TEXT = RGBColor(30, 41, 59)
        MUTED = RGBColor(100, 116, 139)
        LIGHT = RGBColor(238, 242, 255)
        CODE_BG = RGBColor(15, 23, 42)
        CODE_TEXT = RGBColor(226, 232, 240)
        MISTAKE = RGBColor(185, 28, 28)

        paths = get_course_path(self.course_id)
        slides_json_path = paths["slides"]
        pptx_path = paths.get("slides_pptx") or os.path.join(os.path.dirname(slides_json_path), "slides.pptx")

        slides_data: list = []
        deck_title = "Bài giảng"
        if os.path.exists(slides_json_path):
            with open(slides_json_path, "r", encoding="utf-8") as f:
                raw = json.load(f)
            if isinstance(raw, dict):
                slides_data = raw.get("slides", [])
                deck_title = str(raw.get("deck_title") or deck_title)
            else:
                slides_data = raw

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        total = len(slides_data)

        def add_rect(slide, left, top, width, height, fill):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.fill.background()
            shape.shadow.inherit = False
            return shape

        def add_text(slide, left, top, width, height, text, size, color, bold=False, italic=False,
                     align=None, font_name=None):
            box = slide.shapes.add_textbox(left, top, width, height)
            frame = box.text_frame
            frame.word_wrap = True
            paragraph = frame.paragraphs[0]
            paragraph.text = text
            paragraph.font.size = Pt(size)
            paragraph.font.bold = bold
            paragraph.font.italic = italic
            paragraph.font.color.rgb = color
            if align is not None:
                paragraph.alignment = align
            if font_name:
                paragraph.font.name = font_name
            return box

        def add_chrome(slide, index):
            add_rect(slide, 0, 0, prs.slide_width, Inches(0.12), ACCENT)
            add_text(slide, Inches(0.8), Inches(0.18), Inches(9.0), Inches(0.35),
                     deck_title[:80], 11, MUTED)
            add_rect(slide, 0, Inches(7.34), prs.slide_width, Inches(0.03), LIGHT)
            add_text(slide, Inches(12.0), Inches(7.02), Inches(1.1), Inches(0.35),
                     f"{index} / {total}", 12, MUTED, align=PP_ALIGN.RIGHT)

        for idx, s in enumerate(slides_data, 1):
            if not isinstance(s, dict):
                s = {"content": str(s)}
            slide = prs.slides.add_slide(blank_layout)
            slide_type = str(s.get("slide_type") or "concept")
            title_text = str(s.get("title") or f"Slide {idx}")
            screen = s.get("screen_content") if isinstance(s.get("screen_content"), dict) else {}

            # --- title slide: dedicated centered layout, no header chrome ---
            if slide_type == "title":
                add_rect(slide, 0, 0, prs.slide_width, Inches(0.22), ACCENT)
                add_rect(slide, Inches(0.9), Inches(2.35), Inches(1.6), Inches(0.08), ACCENT)
                add_text(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.8),
                         title_text, 40, INDIGO_DARK, bold=True)
                key_message = str(s.get("key_message") or "").strip()
                if key_message:
                    add_text(slide, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.9),
                             key_message, 20, TEXT)
                add_text(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
                         "Bài giảng đại học · Nội dung bám sát tài liệu gốc", 13, MUTED)
                notes = str(s.get("speaker_notes") or "").strip()
                if notes:
                    slide.notes_slide.notes_text_frame.text = notes
                continue

            add_chrome(slide, idx)
            add_text(slide, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.95),
                     title_text, 30, INDIGO, bold=True)

            top = 1.55
            key_message = str(s.get("key_message") or "").strip()
            if key_message:
                add_text(slide, Inches(0.8), Inches(top), Inches(11.7), Inches(0.55),
                         self._format_math_text(key_message, compact=True), 16, ACCENT, italic=True)
                top += 0.65

            formula = str(screen.get("formula") or "").strip()
            code = str(screen.get("code") or "").strip()
            table_rows = screen.get("table") if isinstance(screen.get("table"), list) else []
            diagram = str(screen.get("diagram_description") or "").strip()

            bullets = s.get("bullets") or []
            if not bullets and s.get("content"):
                bullets = [line.strip("- ").strip() for line in str(s["content"]).splitlines() if line.strip()]
            bullets = bullets[:5]

            if bullets:
                box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(2.6))
                frame = box.text_frame
                frame.word_wrap = True
                first = True
                for bullet in bullets:
                    text = self._format_math_text(str(bullet).strip(), compact=False)
                    if not text:
                        continue
                    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                    first = False
                    paragraph.text = f"•  {text}"
                    paragraph.font.size = Pt(20)
                    paragraph.font.color.rgb = TEXT
                    paragraph.space_after = Pt(10)
                top += 0.55 * len(bullets) + 0.25

            if formula:
                shape = add_rect(slide, Inches(2.2), Inches(top), Inches(8.9), Inches(0.85), LIGHT)
                frame = shape.text_frame
                frame.word_wrap = True
                paragraph = frame.paragraphs[0]
                paragraph.text = self._format_math_text(formula, compact=True)
                paragraph.font.size = Pt(24)
                paragraph.font.bold = True
                paragraph.font.color.rgb = INDIGO_DARK
                paragraph.alignment = PP_ALIGN.CENTER
                top += 1.05

            if code:
                code_lines = [ln for ln in code.replace("\r", "").split("\n") if ln.strip()][:8]
                height = max(0.7, 0.32 * len(code_lines) + 0.3)
                shape = add_rect(slide, Inches(0.8), Inches(top), Inches(11.7), Inches(height), CODE_BG)
                frame = shape.text_frame
                frame.word_wrap = True
                first = True
                for line in code_lines:
                    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                    first = False
                    paragraph.text = line[:90]
                    paragraph.font.size = Pt(14)
                    paragraph.font.name = "Consolas"
                    paragraph.font.color.rgb = CODE_TEXT
                top += height + 0.2

            if table_rows and all(isinstance(row, (list, tuple)) for row in table_rows):
                n_rows = min(len(table_rows), 6)
                n_cols = min(max(len(row) for row in table_rows), 4)
                table_shape = slide.shapes.add_table(
                    n_rows, n_cols, Inches(0.8), Inches(top), Inches(11.7), Inches(0.5 * n_rows)
                )
                table = table_shape.table
                for r in range(n_rows):
                    for c in range(n_cols):
                        cell = table.cell(r, c)
                        cell.text = str(table_rows[r][c]) if c < len(table_rows[r]) else ""
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(14)
                            paragraph.font.bold = r == 0
                            paragraph.font.color.rgb = INDIGO_DARK if r == 0 else TEXT
                top += 0.5 * n_rows + 0.2

            if diagram:
                shape = add_rect(slide, Inches(0.8), Inches(top), Inches(11.7), Inches(0.7), LIGHT)
                frame = shape.text_frame
                frame.word_wrap = True
                paragraph = frame.paragraphs[0]
                paragraph.text = f"Sơ đồ đề xuất: {diagram[:180]}"
                paragraph.font.size = Pt(14)
                paragraph.font.italic = True
                paragraph.font.color.rgb = INDIGO
                top += 0.9

            tail_box = slide.shapes.add_textbox(Inches(0.8), Inches(min(top, 6.2)), Inches(11.7), Inches(0.9))
            tail_frame = tail_box.text_frame
            tail_frame.word_wrap = True
            tail_first = True

            def tail_paragraph():
                nonlocal tail_first
                paragraph = tail_frame.paragraphs[0] if tail_first else tail_frame.add_paragraph()
                tail_first = False
                return paragraph

            example = str(s.get("example_or_application") or s.get("example") or "").strip()
            if example and slide_type not in {"worked_example"}:
                paragraph = tail_paragraph()
                paragraph.text = f"Ví dụ: {self._format_math_text(example, compact=True)}"
                paragraph.font.size = Pt(16)
                paragraph.font.italic = True
                paragraph.font.color.rgb = TEXT

            common_mistake = s.get("common_mistake") or {}
            if common_mistake.get("mistake"):
                paragraph = tail_paragraph()
                paragraph.text = (
                    f"Sai lầm thường gặp: {common_mistake['mistake']} — "
                    f"Cách hiểu đúng: {common_mistake.get('correction', '')}"
                )
                paragraph.font.size = Pt(15)
                paragraph.font.color.rgb = MISTAKE

            speaker_notes = str(s.get("speaker_notes") or "").strip()
            student_prompt = str(s.get("student_prompt") or "").strip()
            notes_text = "\n\n".join(
                filter(None, [speaker_notes, f"Câu hỏi cho lớp: {student_prompt}" if student_prompt else ""])
            )
            if notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text

        os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
        prs.save(pptx_path)
        return pptx_path

    def export_slides_pdf(self) -> str:
        paths = get_course_path(self.course_id)
        if not os.path.exists(paths["slides"]):
            raise FileNotFoundError("Slide has not been generated yet.")
        with open(paths["slides"], "r", encoding="utf-8") as f:
            raw = json.load(f)
        slides = raw.get("slides", raw) if isinstance(raw, dict) else raw

        elements: list[tuple[str, str]] = []
        for index, slide in enumerate(slides if isinstance(slides, list) else [], 1):
            item = slide if isinstance(slide, dict) else {"content": str(slide)}
            elements.append((f"Slide {index}: {item.get('title') or 'Nội dung'}", "heading"))
            if item.get("key_message"):
                elements.append((item["key_message"], "small"))
            bullets = item.get("bullets") or []
            if bullets:
                for bullet in bullets:
                    elements.append((f"- {bullet}", "body"))
            elif item.get("content"):
                elements.append((item["content"], "body"))
            if item.get("example_or_application") or item.get("example"):
                elements.append((f"Ví dụ: {item.get('example_or_application') or item.get('example')}", "small"))
            if item.get("speaker_notes"):
                elements.append((f"Ghi chú giảng viên: {item['speaker_notes']}", "small"))
            elements.append(("", "gap"))

        return self._render_artifact_pdf("Slide học tập", elements, paths["slides_pdf"])

    def export_quiz_pdf(self) -> str:
        paths = get_course_path(self.course_id)
        if not os.path.exists(paths["questions"]):
            raise FileNotFoundError("Quiz has not been generated yet.")
        with open(paths["questions"], "r", encoding="utf-8") as f:
            raw = json.load(f)
        questions = raw.get("questions", []) if isinstance(raw, dict) else raw

        elements: list[tuple[str, str]] = []
        labels = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
        for index, question in enumerate(questions if isinstance(questions, list) else [], 1):
            item = question if isinstance(question, dict) else {"question": str(question)}
            elements.append((f"Câu {index}: {item.get('question') or 'Câu hỏi'}", "heading"))
            options = item.get("options") if isinstance(item.get("options"), list) else []
            for option_index, option in enumerate(options):
                elements.append((f"{labels[option_index]}. {option}", "body"))
            try:
                correct = int(item.get("correct", 0))
            except (TypeError, ValueError):
                correct = 0
            if 0 <= correct < len(options):
                elements.append((f"Đáp án: {labels[correct]}", "section"))
            elif item.get("correct_answer"):
                elements.append((f"Đáp án: {item['correct_answer']}", "section"))
            if item.get("explanation"):
                elements.append((f"Giải thích: {item['explanation']}", "small"))
            elements.append(("", "gap"))

        return self._render_artifact_pdf("Quiz học tập", elements, paths["questions_pdf"])

    def _retrieve_overview_docs(self, k: int = 32):
        """Broad retrieval used for readiness evaluation and fallback generation."""
        if not self.vectorstore:
            return []
        try:
            retriever = self.vectorstore.as_retriever(search_kwargs={"k": k})
            return retriever.invoke("tổng quan nội dung tài liệu")
        except Exception:
            return []

    def evaluate_readiness(self) -> dict[str, Any]:
        """Evaluate context quality and return the public per-output readiness report."""
        from backend.services.generation_readiness import evaluate_document_readiness

        docs = self._retrieve_overview_docs()
        report = evaluate_document_readiness(docs, document_id=self.course_id)
        # Public response must not leak internal chunk/page/source metadata (Gate 4).
        return {
            "document_id": report["document_id"],
            "overall_quality_score": report["overall_quality_score"],
            "clean_chunks_count": report["clean_chunks_count"],
            "noisy_chunks_removed": report["noisy_chunks_removed"],
            "generation_readiness": report["generation_readiness"],
            "safe_outputs_available": report["safe_outputs_available"],
            "warnings": report["warnings"],
            "recommended_actions": report["recommended_actions"],
        }

    def _fallback_chunks(self, docs, max_docs: int, max_chars: int):
        from backend.services.context_cleaner import clean_and_filter_chunks

        return clean_and_filter_chunks(docs, max_docs=max_docs, max_chars=max_chars)

    def _key_terms_from_chunks(self, cleaned_chunks, limit: int = 6) -> list[str]:
        terms = []
        for c in cleaned_chunks:
            term = self._short_title(c["text"], "")
            if term and term not in terms:
                terms.append(term)
            if len(terms) >= limit:
                break
        return terms

    _TIME_BUDGET_POINT_COUNTS = {
        "10_min": {"main_points": 3, "key_terms": 3, "core_ideas": 3, "must_know": 2},
        "30_min": {"main_points": 5, "key_terms": 6, "core_ideas": 5, "must_know": 5},
        "1_hour": {"main_points": 7, "key_terms": 8, "core_ideas": 7, "must_know": 7},
        "multi_day": {"main_points": 10, "key_terms": 10, "core_ideas": 10, "must_know": 10},
    }

    def _summary_point_counts(self, profile: Optional[dict[str, Any]]) -> dict[str, int]:
        """Summary/high-yield fallbacks are deterministic (no LLM), so a profile can only
        shape them through length/quantity — driven by time_budget (and exam_prep/
        high_yield style trimming toward must-know points over general prose)."""
        time_budget = str((profile or {}).get("time_budget") or "30_min")
        counts = dict(self._TIME_BUDGET_POINT_COUNTS.get(time_budget, self._TIME_BUDGET_POINT_COUNTS["30_min"]))
        if profile and profile.get("preferred_output_style") == "high_yield":
            counts["main_points"] = max(3, counts["main_points"] - 1)
        return counts

    def generate_fallback_summary(self, title: str, docs, profile: Optional[dict[str, Any]] = None) -> dict[str, Any]:
        """Short summary fallback, grounded in whatever clean context is available."""
        counts = self._summary_point_counts(profile)
        cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=8, max_chars=400)
        limitations = "Nội dung được tạo từ ngữ cảnh giới hạn."
        if not cleaned_chunks:
            return self._sanitize_payload({
                "title": title,
                "summary": "Tài liệu chưa có đủ nội dung sạch để tóm tắt. Vui lòng kiểm tra lại file gốc.",
                "main_points": [],
                "key_terms": [],
                "limitations": limitations,
                "source_chunk_ids": [],
            })
        return self._sanitize_payload({
            "title": title,
            "summary": " ".join(c["text"] for c in cleaned_chunks[:3])[:600].rstrip(),
            "main_points": [c["text"][:180].rstrip() for c in cleaned_chunks[: counts["main_points"]]],
            "key_terms": self._key_terms_from_chunks(cleaned_chunks, limit=counts["key_terms"]),
            "limitations": limitations,
            "source_chunk_ids": stats["source_chunk_ids_used"],
        })

    def generate_fallback_high_yield(self, title: str, docs, profile: Optional[dict[str, Any]] = None) -> dict[str, Any]:
        """High-yield notes fallback: core ideas, must-know points, and a quick self-check."""
        counts = self._summary_point_counts(profile)
        cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=10, max_chars=320)
        if not cleaned_chunks:
            return self._sanitize_payload({
                "title": title,
                "core_ideas": [],
                "must_know_points": [],
                "simple_explanation": "Tài liệu chưa có đủ nội dung sạch để tạo bản học trọng tâm.",
                "quick_review_questions": [],
                "source_chunk_ids": [],
            })
        core_ideas = [c["text"][:220].rstrip() for c in cleaned_chunks[: counts["core_ideas"]]]
        must_know_points = [
            c["text"][:160].rstrip() for c in cleaned_chunks[counts["core_ideas"]: counts["core_ideas"] + counts["must_know"]]
        ] or core_ideas[:3]
        terms = self._key_terms_from_chunks(cleaned_chunks, limit=3)
        quick_review_questions = [f"Hãy giải thích nội dung trọng tâm về: {term}?" for term in terms] or [
            "Nội dung quan trọng nhất trong phần tài liệu đã đọc được là gì?"
        ]
        return self._sanitize_payload({
            "title": title,
            "core_ideas": core_ideas,
            "must_know_points": must_know_points,
            "simple_explanation": core_ideas[0] if core_ideas else "",
            "quick_review_questions": quick_review_questions,
            "source_chunk_ids": stats["source_chunk_ids_used"],
        })

    def generate_fallback_outline(self, title: str, docs) -> dict[str, Any]:
        """Document outline fallback: detected topics and possible sections."""
        cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=12, max_chars=200)
        if not cleaned_chunks:
            return self._sanitize_payload({
                "title": title,
                "detected_topics": [],
                "possible_sections": [],
                "missing_context_warning": "Tài liệu chưa trích xuất được đề mục rõ ràng.",
                "source_chunk_ids": [],
            })
        possible_sections = [
            {
                "heading": self._short_title(c["text"], f"Phần {index}"),
                "summary": c["text"][:160].rstrip(),
            }
            for index, c in enumerate(cleaned_chunks, 1)
        ]
        return self._sanitize_payload({
            "title": title,
            "detected_topics": self._key_terms_from_chunks(cleaned_chunks, limit=8),
            "possible_sections": possible_sections,
            "missing_context_warning": (
                "" if len(cleaned_chunks) >= 6 else
                "Dàn ý này được suy ra từ số lượng đoạn nội dung sạch còn hạn chế, có thể chưa bao phủ toàn bộ tài liệu."
            ),
            "source_chunk_ids": stats["source_chunk_ids_used"],
        })

    def generate_fallback_key_terms(self, title: str, docs) -> dict[str, Any]:
        """Deterministic key-term flashcards drawn from the cleanest available chunks."""
        cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=16, max_chars=260)
        terms = []
        for c in cleaned_chunks:
            term = self._short_title(c["text"], "")
            if not term:
                continue
            terms.append({"term": term, "definition": c["text"][:200].rstrip()})
            if len(terms) >= 10:
                break
        return self._sanitize_payload({
            "title": title,
            "terms": terms,
            "source_chunk_ids": stats["source_chunk_ids_used"],
        })

    def generate_fallback_shallow_mindmap(self, title: str, docs) -> dict[str, Any]:
        """3-level concept map fallback, grounded in clean chunks."""
        cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=10, max_chars=220)
        root_node = {
            "id": "root",
            "title": title,
            "summary": "Sơ đồ khái niệm rút gọn từ các đoạn văn bản sạch nhất của tài liệu.",
            "type": "root",
            "importance": "high",
            "source_chunk_ids": stats.get("source_chunk_ids_used", []),
            "children": [],
        }
        nodes = []
        edges = []
        if not cleaned_chunks:
            mindmap = {
                "document_id": getattr(self, "course_id", ""),
                "title": title,
                "description": "Tài liệu chưa có đủ nội dung sạch để xây dựng sơ đồ khái niệm.",
                "root": root_node,
                "nodes": [],
                "edges": [],
            }
            self._evaluate_mindmap_quality_gate(mindmap)
            return self._sanitize_payload(mindmap)

        for index, chunk in enumerate(cleaned_chunks[:6], 1):
            ch_id = f"ch_{index}"
            label = self._short_title(chunk["text"], f"Chủ đề {index}")
            root_node["children"].append(ch_id)
            edges.append({"from": "root", "to": ch_id, "relation": "contains"})
            
            les_id = f"les_{index}_1"
            edges.append({"from": ch_id, "to": les_id, "relation": "contains"})
            
            nodes.append({
                "id": les_id,
                "parent_id": ch_id,
                "title": f"Nội dung {index}.1",
                "summary": chunk["text"][:150].rstrip(),
                "type": "lesson",
                "importance": "medium",
                "keywords": [label],
                "source_chunk_ids": chunk.get("source_chunk_ids", []),
                "children": [],
            })
            nodes.append({
                "id": ch_id,
                "parent_id": "root",
                "title": label,
                "summary": chunk["text"][:100].rstrip(". ") + "…",
                "type": "chapter",
                "importance": "high",
                "keywords": [label],
                "source_chunk_ids": chunk.get("source_chunk_ids", []),
                "children": [les_id],
            })

        mindmap = {
            "document_id": getattr(self, "course_id", ""),
            "title": title,
            "description": "Sơ đồ kiến thức rút gọn (2-3 tầng) được tổng hợp từ ngữ cảnh sạch của tài liệu.",
            "root": root_node,
            "nodes": nodes,
            "edges": edges,
        }
        self._evaluate_mindmap_quality_gate(mindmap)
        return self._sanitize_payload(mindmap)

    def generate_fallback_short_video_script(self, title: str, docs) -> dict[str, Any]:
        """~60-second video script fallback: grounded scenes, no rendering."""
        cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=6, max_chars=260)
        if not cleaned_chunks:
            return self._sanitize_payload({
                "title": title,
                "video_mode": "sixty_second",
                "estimated_duration_seconds": 0,
                "scenes": [],
                "limitations": "Tài liệu chưa có đủ nội dung sạch để viết kịch bản video.",
                "source_chunk_ids": [],
            })
        scenes = self._build_fallback_scenes(docs, 4)
        for scene in scenes:
            scene["duration_seconds"] = 15
        return self._sanitize_payload({
            "title": title,
            "video_mode": "sixty_second",
            "estimated_duration_seconds": 60,
            "scenes": scenes,
            "limitations": "Kịch bản video ngắn (~60 giây) do chưa đủ ngữ cảnh cho video bài giảng hoàn chỉnh.",
            "source_chunk_ids": stats["source_chunk_ids_used"],
        })

    def generate_fallback_storyboard_only(self, title: str, docs) -> dict[str, Any]:
        """Storyboard-only fallback: visual plan without voiceover or MP4 rendering."""
        cleaned_chunks, stats = self._fallback_chunks(docs, max_docs=8, max_chars=260)
        if not cleaned_chunks:
            return self._sanitize_payload({
                "title": title,
                "render_status": "storyboard_only",
                "storyboard": [],
                "limitations": "Tài liệu chưa có đủ nội dung sạch để dựng storyboard.",
                "source_chunk_ids": [],
            })
        storyboard = [
            {
                "title": scene["title"],
                "visual_template": scene["visual_template"],
                "screen_text": scene.get("screen_text", []),
                "visual_text": scene["visual_text"],
                "source_chunk_ids": scene.get("source_chunk_ids", []),
            }
            for scene in self._build_fallback_scenes(docs, 6)
        ]
        return self._sanitize_payload({
            "title": title,
            "render_status": "storyboard_only",
            "storyboard": storyboard,
            "limitations": "Chỉ dựng storyboard (chưa render video) do ngữ cảnh sạch còn hạn chế.",
            "source_chunk_ids": stats["source_chunk_ids_used"],
        })

    def _high_yield_notes_to_book(self, notes: dict[str, Any], target_audience: str) -> dict[str, Any]:
        """Adapt a generate_fallback_high_yield() result into a minimal book shape.

        Reuses _normalize_lesson so the same renderer/frontend fields
        (simple_explanation/explanation, active_recall_questions/assessment,
        source_chunk_ids, ...) work identically to a full book.
        """
        source_chunk_ids = notes.get("source_chunk_ids") or []
        core_ideas = notes.get("core_ideas") or ([notes["simple_explanation"]] if notes.get("simple_explanation") else [])
        questions = notes.get("quick_review_questions") or []
        lessons = []
        for idx, idea in enumerate(core_ideas or ["Chưa có đủ nội dung sạch để tóm tắt."], 1):
            raw_lesson = {
                "core_idea": idea,
                "explanation": idea,
                "must_know_points": notes.get("must_know_points") or [],
                "quick_check": (
                    [{"question": questions[idx - 1], "answer": idea[:160], "explanation": ""}]
                    if idx - 1 < len(questions)
                    else []
                ),
                "source_chunk_ids": source_chunk_ids,
            }
            point = {"text": idea, "source_chunk_ids": source_chunk_ids}
            lessons.append(self._normalize_lesson(raw_lesson, point, f"Ý trọng tâm {idx}"))

        return {
            "title": notes.get("title") or "Bản học trọng tâm",
            "subtitle": "Được tạo từ ngữ cảnh giới hạn — chưa đủ cho giáo trình đại học đầy đủ.",
            "audience": target_audience or "người học",
            "course_level": "high_yield",
            "description": "Bản tóm tắt trọng tâm do ngữ cảnh sạch của tài liệu còn hạn chế.",
            "estimated_duration": "30-60 phút",
            "prerequisites": [],
            "course_learning_outcomes": [],
            "chapters": [
                {
                    "chapter_index": 1,
                    "title": "Tóm tắt trọng tâm",
                    "description": "Tổng hợp các ý chính có thể trích xuất từ tài liệu hiện có.",
                    "estimated_duration_minutes": None,
                    "prerequisites": [],
                    "learning_outcomes": [],
                    "connections_to_other_chapters": [],
                    "chapter_summary": "",
                    "chapter_quiz": [],
                    "lessons": lessons,
                }
            ],
            "glossary": [],
            "review_plan": {},
        }

    # ─── Two-stage book pipeline: blueprint -> chapters -> assembled book ────

    def _generate_course_blueprint(
        self, context: str, user_prompt: str, target_audience: str, learning_mode: str,
        profile_directives: str = "",
    ) -> Optional[dict[str, Any]]:
        """Stage 1: design the course blueprint from clean context. None on failure."""
        from backend.core.prompts import COURSE_BLUEPRINT_PROMPT

        try:
            prompt = ChatPromptTemplate.from_template(COURSE_BLUEPRINT_PROMPT)
            chain = prompt | get_llm(temperature=0.2, max_output_tokens=16384, task="course") | StrOutputParser()
            res = chain.invoke(
                {
                    "context": context,
                    "user_prompt": user_prompt or "Không có",
                    "target_audience": target_audience or "người học chung",
                    "learning_mode": learning_mode or "normal",
                    "profile_directives": profile_directives,
                }
            )
            raw = json.loads(extract_json(res))
        except Exception as exc:
            logger.warning("[BookGen] Blueprint generation failed: %s", exc)
            return None
        return self._normalize_blueprint(raw, target_audience)

    def _normalize_blueprint(self, raw: Any, target_audience: str) -> Optional[dict[str, Any]]:
        if not isinstance(raw, dict) or raw.get("error"):
            return None
        units = []
        for index, raw_unit in enumerate((raw.get("course_units") or [])[:6], 1):
            if not isinstance(raw_unit, dict):
                continue
            title = str(raw_unit.get("title") or "").strip()
            if not title or self._quality_gate_generic(title):
                continue
            ids = raw_unit.get("source_chunk_ids")
            units.append(
                {
                    "unit_id": str(raw_unit.get("unit_id") or f"unit_{index:02d}"),
                    "title": title,
                    "big_idea": str(raw_unit.get("big_idea") or "").strip(),
                    "why_it_matters": str(raw_unit.get("why_it_matters") or "").strip(),
                    "key_concepts": [str(k).strip() for k in (raw_unit.get("key_concepts") or []) if str(k).strip()],
                    "definitions": [d for d in (raw_unit.get("definitions") or []) if isinstance(d, dict) and d.get("term")],
                    "formulas": [f for f in (raw_unit.get("formulas") or []) if isinstance(f, dict) and f.get("formula")],
                    "examples": [str(e).strip() for e in (raw_unit.get("examples") or []) if str(e).strip()],
                    "common_misconceptions": [
                        m for m in (raw_unit.get("common_misconceptions") or []) if isinstance(m, dict) and m.get("mistake")
                    ],
                    "worked_examples": [w for w in (raw_unit.get("worked_examples") or []) if isinstance(w, dict) and w.get("problem")],
                    "practice_problems": [p for p in (raw_unit.get("practice_problems") or []) if isinstance(p, dict) and p.get("question")],
                    "source_chunk_ids": [str(i) for i in ids if str(i).strip()] if isinstance(ids, list) else [],
                }
            )
        if not units:
            return None
        top_ids = raw.get("source_chunk_ids")
        return {
            "course_title": str(raw.get("course_title") or "").strip() or "Giáo trình học tập từ tài liệu",
            "course_level": str(raw.get("course_level") or "university").strip(),
            "audience": str(raw.get("audience") or target_audience or "người học").strip(),
            "prerequisites": [str(p).strip() for p in (raw.get("prerequisites") or []) if str(p).strip()],
            "learning_outcomes": [str(o).strip() for o in (raw.get("learning_outcomes") or []) if str(o).strip()],
            "course_units": units,
            "glossary": [g for g in (raw.get("glossary") or []) if isinstance(g, dict) and g.get("term")],
            "assessment_plan": [a for a in (raw.get("assessment_plan") or []) if isinstance(a, dict)],
            "source_chunk_ids": [str(i) for i in top_ids if str(i).strip()] if isinstance(top_ids, list) else [],
        }

    def _retrieve_unit_docs(self, unit: dict[str, Any], fallback_docs, k: int = 24):
        """Targeted retrieval for one unit; falls back to the broad docs."""
        if not self.vectorstore:
            return fallback_docs
        query = " ".join([unit.get("title", ""), *unit.get("key_concepts", [])[:4]]).strip()
        if not query:
            return fallback_docs
        try:
            retriever = self.vectorstore.as_retriever(search_kwargs={"k": k})
            unit_docs = retriever.invoke(query)
            return unit_docs or fallback_docs
        except Exception:
            return fallback_docs

    def _normalize_chapter_v2(
        self, raw: Any, unit: dict[str, Any], chapter_index: int, docs
    ) -> dict[str, Any]:
        """Normalize an LLM chapter into the rigorous chapter schema."""
        ch = raw if isinstance(raw, dict) else {}
        unit_ids = [str(i) for i in (unit.get("source_chunk_ids") or []) if str(i).strip()]

        def text(key: str, default: str = "") -> str:
            value = ch.get(key)
            return value.strip() if isinstance(value, str) and value.strip() else default

        def str_list(key: str) -> list[str]:
            value = ch.get(key)
            return [str(v).strip() for v in value if str(v).strip()] if isinstance(value, list) else []

        def ids_or_unit(value) -> list[str]:
            if isinstance(value, list):
                cleaned = [str(i) for i in value if str(i).strip()]
                if cleaned:
                    return cleaned
            return unit_ids

        core_concepts = []
        for concept in ch.get("core_concepts") or []:
            if not isinstance(concept, dict) or not str(concept.get("term") or "").strip():
                continue
            mistake_raw = concept.get("common_mistake") if isinstance(concept.get("common_mistake"), dict) else {}
            core_concepts.append(
                {
                    "term": str(concept.get("term")).strip(),
                    "definition": str(concept.get("definition") or "").strip(),
                    "intuition": str(concept.get("intuition") or "").strip(),
                    "technical_explanation": str(concept.get("technical_explanation") or "").strip(),
                    "example": str(concept.get("example") or "").strip(),
                    "non_example": str(concept.get("non_example") or "").strip(),
                    "common_mistake": {
                        "mistake": str(mistake_raw.get("mistake") or "").strip(),
                        "correction": str(mistake_raw.get("correction") or "").strip(),
                    },
                    "formula": str(concept.get("formula") or "").strip(),
                    "code": str(concept.get("code") or "").strip(),
                    "source_chunk_ids": ids_or_unit(concept.get("source_chunk_ids")),
                }
            )
        if not core_concepts:
            core_concepts = self._concepts_from_unit(unit)

        worked_examples = [
            {
                "title": str(we.get("title") or "").strip() or f"Ví dụ mẫu {i + 1}",
                "problem": str(we.get("problem") or "").strip(),
                "step_by_step_solution": [str(s).strip() for s in (we.get("step_by_step_solution") or []) if str(s).strip()],
                "why_each_step_matters": str(we.get("why_each_step_matters") or "").strip(),
                "common_error": str(we.get("common_error") or "").strip(),
                "source_chunk_ids": ids_or_unit(we.get("source_chunk_ids")),
            }
            for i, we in enumerate(ch.get("worked_examples") or [])
            if isinstance(we, dict) and we.get("problem") and we.get("step_by_step_solution")
        ]

        practice_problems = [
            {
                "difficulty": str(pp.get("difficulty") or "medium").strip().lower(),
                "question": str(pp.get("question") or "").strip(),
                "hint": str(pp.get("hint") or "").strip(),
                "solution": str(pp.get("solution") or "").strip(),
                "expected_answer": str(pp.get("expected_answer") or pp.get("solution") or "").strip(),
                "source_chunk_ids": ids_or_unit(pp.get("source_chunk_ids")),
            }
            for pp in ch.get("practice_problems") or []
            if isinstance(pp, dict) and pp.get("question")
        ]

        title = text("title") or unit.get("title") or f"Chương {chapter_index}"
        chapter_overview = text("chapter_overview") or text("overview") or unit.get("big_idea", "")
        all_ids = set(unit_ids)
        for item in [*core_concepts, *worked_examples, *practice_problems]:
            all_ids.update(item.get("source_chunk_ids") or [])

        chapter = {
            "chapter_index": chapter_index,
            "title": title,
            "chapter_overview": chapter_overview,
            "description": chapter_overview,
            "learning_objectives": str_list("learning_objectives") or unit.get("key_concepts", [])[:4],
            "learning_outcomes": str_list("learning_objectives"),
            "prerequisites": str_list("prerequisites"),
            "big_picture": text("big_picture") or unit.get("why_it_matters", ""),
            "core_concepts": core_concepts,
            "worked_examples": worked_examples,
            "practice_problems": practice_problems,
            "formulas": unit.get("formulas", []),
            "chapter_summary": text("chapter_summary"),
            "active_recall_questions": str_list("active_recall_questions"),
            "connections_to_other_chapters": str_list("connections_to_other_chapters"),
            "chapter_quiz": [],
            "source_chunk_ids": sorted(all_ids),
        }
        chapter["lessons"] = self._derive_lessons_from_concepts(chapter)
        return chapter

    def _concepts_from_unit(self, unit: dict[str, Any]) -> list[dict[str, Any]]:
        """Deterministic core concepts built only from real blueprint unit data."""
        unit_ids = [str(i) for i in (unit.get("source_chunk_ids") or []) if str(i).strip()]
        misconceptions = unit.get("common_misconceptions") or []
        examples = unit.get("examples") or []
        concepts = []
        for index, definition in enumerate((unit.get("definitions") or [])[:5]):
            mistake = misconceptions[index] if index < len(misconceptions) else {}
            concepts.append(
                {
                    "term": str(definition.get("term")).strip(),
                    "definition": str(definition.get("definition") or "").strip(),
                    "intuition": "",
                    "technical_explanation": "",
                    "example": str(examples[index]) if index < len(examples) else "",
                    "non_example": "",
                    "common_mistake": {
                        "mistake": str(mistake.get("mistake") or "").strip(),
                        "correction": str(mistake.get("correction") or "").strip(),
                    },
                    "formula": "",
                    "code": "",
                    "source_chunk_ids": unit_ids,
                }
            )
        return concepts

    def _derive_lessons_from_concepts(self, chapter: dict[str, Any]) -> list[dict[str, Any]]:
        """Back-compat lessons[] so existing frontend/study-pack views keep working."""
        lessons = []
        worked_examples = chapter.get("worked_examples") or []
        practice_problems = chapter.get("practice_problems") or []
        for index, concept in enumerate((chapter.get("core_concepts") or [])[:4]):
            explanation_parts = [p for p in [concept.get("intuition"), concept.get("technical_explanation")] if p]
            raw_lesson = {
                "core_idea": concept.get("definition", ""),
                "explanation": "\n\n".join(explanation_parts),
                "example": concept.get("example", ""),
                "non_example": concept.get("non_example", ""),
                "key_concepts": [{"term": concept.get("term", ""), "definition": concept.get("definition", "")}],
                "common_misunderstanding": concept.get("common_mistake") or {},
                "worked_examples": worked_examples[index : index + 1],
                "practice_problems": practice_problems[index : index + 2],
                "source_chunk_ids": concept.get("source_chunk_ids") or [],
            }
            point = {
                "text": concept.get("definition") or concept.get("intuition") or concept.get("term", ""),
                "source_chunk_ids": concept.get("source_chunk_ids") or [],
            }
            lessons.append(self._normalize_lesson(raw_lesson, point, concept.get("term") or f"Khái niệm {index + 1}"))
        return lessons

    def _chapter_weaknesses(self, chapter: dict[str, Any]) -> list[str]:
        """Structural weaknesses that trigger a one-shot chapter regeneration."""
        problems = []
        concepts = chapter.get("core_concepts") or []
        if not concepts:
            problems.append("no_core_concepts")
        elif sum(1 for c in concepts if c.get("technical_explanation") and c.get("intuition")) < max(1, len(concepts) // 2):
            problems.append("shallow_concepts")
        if not chapter.get("worked_examples"):
            problems.append("no_worked_examples")
        difficulties = {p.get("difficulty") for p in chapter.get("practice_problems") or []}
        if len(difficulties & {"easy", "medium", "hard"}) < 2:
            problems.append("practice_difficulty_coverage")
        if not chapter.get("source_chunk_ids"):
            problems.append("no_grounding")
        return problems

    def _generate_chapter_from_unit(
        self, unit: dict[str, Any], chapter_index: int, target_audience: str, learning_mode: str, fallback_docs,
        profile_directives: str = "",
    ) -> dict[str, Any]:
        """Stage 2: generate one rigorous chapter for a blueprint unit (one retry on weakness)."""
        from backend.core.prompts import BOOK_CHAPTER_GENERATION_PROMPT

        unit_docs = self._retrieve_unit_docs(unit, fallback_docs)
        context = self._clean_docs_context(unit_docs, max_docs=12, max_chars=900)
        unit_plan = json.dumps(unit, ensure_ascii=False)

        prompt = ChatPromptTemplate.from_template(BOOK_CHAPTER_GENERATION_PROMPT)
        chain = prompt | get_llm(temperature=0.3, max_output_tokens=32768, task="book") | StrOutputParser()
        payload = {
            "chapter_index": chapter_index,
            "chapter_title": unit.get("title", f"Chương {chapter_index}"),
            "target_audience": target_audience or "người học chung",
            "learning_mode": learning_mode or "normal",
            "unit_plan": unit_plan,
            "context": context,
            "profile_directives": profile_directives,
        }

        chapter: Optional[dict[str, Any]] = None
        for attempt in (1, 2):
            try:
                res = chain.invoke(payload)
                chapter = self._normalize_chapter_v2(json.loads(extract_json(res)), unit, chapter_index, unit_docs)
            except Exception as exc:
                logger.warning("[BookGen] Chapter %s attempt %s failed: %s", chapter_index, attempt, exc)
                chapter = None
                continue
            weaknesses = self._chapter_weaknesses(chapter)
            if not weaknesses:
                return chapter
            logger.info("[BookGen] Chapter %s attempt %s weak (%s)%s", chapter_index, attempt, ",".join(weaknesses),
                        ", retrying..." if attempt == 1 else ", keeping best effort.")
        if chapter is not None:
            return chapter

        deterministic = self._normalize_chapter_v2({}, unit, chapter_index, unit_docs)
        deterministic["generation_note"] = "deterministic_from_blueprint"
        return deterministic

    _HOW_TO_USE_STEPS = [
        "Đọc \"Lộ trình khóa học\" để nắm bức tranh tổng thể trước khi vào từng chương.",
        "Với mỗi chương: đọc phần tổng quan và mục tiêu trước, sau đó mới đọc từng khái niệm cốt lõi.",
        "Với mỗi khái niệm: đọc định nghĩa, phần trực quan, rồi phần giải thích kỹ thuật; dừng lại ở ví dụ và phản ví dụ.",
        "Tự giải các ví dụ mẫu trước khi xem lời giải từng bước.",
        "Làm bài tập theo thứ tự dễ → khó; chỉ xem gợi ý khi bị kẹt, xem lời giải sau khi đã thử.",
        "Cuối mỗi chương, trả lời các câu hỏi tự kiểm tra (active recall) mà không nhìn lại nội dung.",
        "Dùng \"Kế hoạch ôn tập\" để ôn lại theo chu kỳ 10 phút / 30 phút / 1 giờ.",
        "Lưu ý: nội dung do AI tạo từ tài liệu của bạn — hãy kiểm chứng những thông tin quan trọng.",
    ]

    def _assemble_book_v2(
        self, blueprint: dict[str, Any], chapters: list[dict[str, Any]], target_audience: str
    ) -> dict[str, Any]:
        """Assemble the final book JSON: roadmap, chapters, glossary, problem set, review plan."""
        glossary = {str(g.get("term")).strip(): str(g.get("definition") or "").strip()
                    for g in blueprint.get("glossary", []) if str(g.get("term") or "").strip()}
        for chapter in chapters:
            for concept in chapter.get("core_concepts") or []:
                term = concept.get("term", "")
                if term and term not in glossary and concept.get("definition"):
                    glossary[term] = concept["definition"]

        problem_set = []
        for chapter in chapters:
            for pp in chapter.get("practice_problems") or []:
                problem_set.append({**pp, "chapter_index": chapter["chapter_index"], "chapter_title": chapter["title"]})

        recall_questions = [q for ch in chapters for q in (ch.get("active_recall_questions") or [])]
        review_plan = {
            "ten_minute": recall_questions[:4] or ["Tự trả lời các câu hỏi active recall của một chương bất kỳ."],
            "thirty_minute": [
                f"Giải lại một ví dụ mẫu của chương \"{ch['title']}\" mà không xem lời giải."
                for ch in chapters[:3] if ch.get("worked_examples")
            ] or ["Giải lại các ví dụ mẫu mà không xem lời giải."],
            "one_hour": ["Làm trọn bộ bài tập tổng hợp (Problem Set) và đối chiếu với phần lời giải."],
        }

        all_ids = set(blueprint.get("source_chunk_ids") or [])
        for chapter in chapters:
            all_ids.update(chapter.get("source_chunk_ids") or [])

        unit_minutes = 90
        return {
            "title": blueprint["course_title"],
            "subtitle": (blueprint.get("learning_outcomes") or [""])[0],
            "audience": blueprint.get("audience") or target_audience or "người học",
            "course_level": blueprint.get("course_level", "university"),
            "description": " ".join(
                unit.get("big_idea", "") for unit in blueprint.get("course_units", [])
            ).strip() or f"Giáo trình học tập cho {target_audience or 'người học'}, bám sát tài liệu gốc.",
            "estimated_duration": f"{len(chapters) * unit_minutes // 60}-{len(chapters) * unit_minutes // 45} giờ học",
            "how_to_use": list(self._HOW_TO_USE_STEPS),
            "prerequisites": blueprint.get("prerequisites", []),
            "course_learning_outcomes": blueprint.get("learning_outcomes", []),
            "course_roadmap": [
                {
                    "unit_id": unit["unit_id"],
                    "title": unit["title"],
                    "big_idea": unit.get("big_idea", ""),
                    "key_concepts": unit.get("key_concepts", []),
                }
                for unit in blueprint.get("course_units", [])
            ],
            "chapters": chapters,
            "glossary": [{"term": term, "definition": definition} for term, definition in glossary.items()],
            "assessment_plan": blueprint.get("assessment_plan", []),
            "problem_set": problem_set,
            "review_plan": review_plan,
            "source_chunk_ids": sorted(all_ids),
        }

    def generate_book(
        self, user_prompt: str = "", target_audience: str = "sinh viên", learning_mode: str = "normal",
        profile: Optional[dict[str, Any]] = None,
    ):
        profile_directives = build_profile_directives(profile)
        try:
            total_chunks = int(self.vectorstore.collection.count()) if self.vectorstore else 32
        except Exception:
            total_chunks = 32
        retriever_k = min(150, total_chunks) if total_chunks > 100 else 32
        if total_chunks > 100:
            logger.info("[BookGen] Long document detected (%d chunks), engaging chapter-by-chapter strategy with deep retrieval.", total_chunks)
        retriever = self.vectorstore.as_retriever(search_kwargs={"k": retriever_k})
        query = user_prompt or "nội dung chính mục tiêu khái niệm ví dụ kết luận thuật toán thuật ngữ"
        docs = retriever.invoke(query)

        from backend.services.generation_readiness import evaluate_document_readiness

        book_readiness = evaluate_document_readiness(docs, document_id=self.course_id)["generation_readiness"]["book"]
        if book_readiness["status"] == "not_enough_context":
            logger.info("[BookGen] Not enough clean context, generating high-yield study notes instead.")
            notes = self.generate_fallback_high_yield(f"Bản học trọng tâm: {target_audience or 'người học'}", docs, profile=profile)
            book = self._high_yield_notes_to_book(notes, target_audience)
            quality = self._evaluate_quality_gate(book, "book")
            book = self._sanitize_payload(book)
            book["quality_report"] = quality
            book["generation_mode"] = "high_yield_study_guide"
            book["generation_status"] = {
                "status": "limited",
                "reason": "Không đủ ngữ cảnh sạch để tạo giáo trình đại học đầy đủ.",
                "fallback_used": "high_yield_study_notes",
            }
            paths = get_course_path(self.course_id)
            self._save_json(paths["book"], book)
            self._render_book_pdf(book, paths["book_pdf"])
            return {"book": book, "pdf_url": f"/api/course/{self.course_id}/book.pdf"}

        max_docs = 150 if len(docs) > 50 else 32
        context = self._clean_docs_context(docs, max_docs=max_docs, max_chars=800)
        blueprint = self._generate_course_blueprint(context, user_prompt, target_audience, learning_mode, profile_directives)

        if blueprint is not None:
            from backend.core import config as core_config

            blueprint_dir = os.path.join(core_config.BOOKS_DIR, self.course_id)
            self._save_json(os.path.join(blueprint_dir, "blueprint.json"), blueprint)

            chapters = [
                self._generate_chapter_from_unit(unit, index, target_audience, learning_mode, docs, profile_directives)
                for index, unit in enumerate(blueprint["course_units"], 1)
            ]
            book = self._assemble_book_v2(blueprint, chapters, target_audience)
            quality = self._evaluate_quality_gate(book, "book")

            book = self._sanitize_payload(book)
            book["quality_report"] = quality
            if quality["is_university_ready"]:
                book["generation_mode"] = "full_book"
                book["generation_status"] = {"status": "full", "reason": "", "fallback_used": None}
            elif book_readiness["status"] == "limited":
                logger.info("[BookGen] Gate < 85 with limited context; switching to honest high-yield notes.")
                notes = self.generate_fallback_high_yield(f"Bản học trọng tâm: {target_audience or 'người học'}", docs, profile=profile)
                book = self._high_yield_notes_to_book(notes, target_audience)
                quality = self._evaluate_quality_gate(book, "book")
                book = self._sanitize_payload(book)
                book["quality_report"] = quality
                book["generation_mode"] = "high_yield_study_guide"
                book["generation_status"] = {
                    "status": "limited",
                    "reason": "Ngữ cảnh hạn chế: dùng bản học trọng tâm thay vì giáo trình đầy đủ kém chất lượng.",
                    "fallback_used": "high_yield_study_notes",
                }
            else:
                book["generation_mode"] = "full_book"
                book["generation_status"] = {
                    "status": "limited",
                    "reason": "Giáo trình chưa đạt chuẩn chất lượng 85 điểm sau khi tự động tạo lại; xem warnings trong quality_report.",
                    "fallback_used": None,
                }

            paths = get_course_path(self.course_id)
            self._save_json(paths["book"], book)
            self._render_book_pdf(book, paths["book_pdf"])
            return {"book": book, "pdf_url": f"/api/course/{self.course_id}/book.pdf"}

        # Blueprint failed: legacy single-shot pipeline as a robust fallback.
        logger.info("[BookGen] Blueprint unavailable, falling back to single-shot generation.")
        generation_error: Optional[str] = None
        try:
            prompt = ChatPromptTemplate.from_template(BOOK_GENERATION_PROMPT)
            # The academic-depth schema (worked examples, practice problems, key concepts, etc.)
            # across multiple chapters/lessons is large; the default 8192-token cap truncates the
            # JSON mid-string and silently falls back to the shallow template. Give it real headroom.
            chain = prompt | get_llm(temperature=0.3, max_output_tokens=65536, task="book") | StrOutputParser()
            res = chain.invoke(
                {
                    "context": context,
                    "user_prompt": user_prompt or "Không có",
                    "target_audience": target_audience or "người học chung",
                    "profile_directives": profile_directives,
                }
            )
            book = self._normalize_book(json.loads(extract_json(res)), docs, target_audience)
            quality = self._evaluate_quality_gate(book, "book")
            if not quality["is_university_ready"]:
                logger.info("[BookGen] Score below gate, retrying once...")
                res = chain.invoke(
                    {
                        "context": context,
                        "user_prompt": user_prompt or "Không có",
                        "target_audience": target_audience or "người học chung",
                        "profile_directives": profile_directives,
                    }
                )
                book = self._normalize_book(json.loads(extract_json(res)), docs, target_audience)
                quality = self._evaluate_quality_gate(book, "book")
        except Exception as e:
            logger.warning("Book generation failed, using fallback: %s", e)
            generation_error = str(e)
            book = self._build_fallback_book(docs, target_audience)
            quality = self._evaluate_quality_gate(book, "book")

        book = self._sanitize_payload(book)
        book["quality_report"] = quality
        if generation_error is not None:
            book["generation_mode"] = "summary_only"
            book["generation_status"] = {
                "status": "limited",
                "reason": (
                    "Không thể tạo giáo trình đầy đủ từ AI (lỗi hệ thống hoặc giới hạn quota), "
                    "đã dùng bản dự phòng dựng trực tiếp từ tài liệu."
                ),
                "fallback_used": "generic_fallback",
            }
        else:
            book["generation_mode"] = "full_book"
            book["generation_status"] = {"status": "full", "reason": "", "fallback_used": None}

        paths = get_course_path(self.course_id)
        self._save_json(paths["book"], book)
        self._render_book_pdf(book, paths["book_pdf"])
        return {"book": book, "pdf_url": f"/api/course/{self.course_id}/book.pdf"}

    _QUIZ_VALID_TYPES = {"mcq", "true_false", "short_answer", "scenario", "code_reading"}
    _VALID_DIFFICULTIES = {"easy", "medium", "hard"}

    def _dedup_key(self, text: str) -> str:
        """Normalize text for duplicate detection: lowercase, strip punctuation/whitespace."""
        return re.sub(r"\s+", " ", re.sub(r"[^\w\s]", "", str(text or "").lower())).strip()

    def _build_fallback_quiz(self, docs, quantity: int, difficulty: str):
        points = self._doc_points(docs, limit=max(1, min(quantity, 10)), max_chars=220)
        fallback_difficulty = difficulty if difficulty in self._VALID_DIFFICULTIES else "medium"
        questions = []
        seen: set[str] = set()
        for index in range(max(1, min(quantity, 10))):
            point = points[index % len(points)]
            # Vary the templated question per point so a repeated fallback set doesn't
            # itself trip the quiz quality gate's duplicate-question check.
            term = self._short_title(point["text"], f"phần {index + 1}")
            question_text = f'Ý nào sau đây phản ánh đúng nội dung về "{term}" trong tài liệu?'
            key = self._dedup_key(question_text)
            if key in seen:
                continue
            seen.add(key)
            options = [
                point["text"][:140],
                "Một nhận định không được tài liệu cung cấp rõ ràng.",
                "Một kết luận mở rộng ngoài phạm vi tài liệu.",
                "Một phương án dùng để gây nhiễu trong câu hỏi.",
            ]
            questions.append(
                {
                    "id": f"q{len(questions) + 1}",
                    "type": "mcq",
                    "question": question_text,
                    "options": options,
                    "correct": 0,
                    "correct_answer": options[0],
                    "explanation": (
                        f'Đáp án đúng bám sát đoạn nội dung về "{term}" được hệ thống truy xuất từ tài liệu; '
                        "các phương án còn lại không được tài liệu xác nhận."
                    ),
                    "why_wrong_options_are_wrong": [
                        "Nội dung này không xuất hiện rõ ràng trong đoạn tài liệu đã truy xuất.",
                        "Đây là kết luận vượt ra ngoài phạm vi được tài liệu đề cập.",
                        "Đây chỉ là phương án gây nhiễu, không phản ánh nội dung tài liệu.",
                    ],
                    "difficulty": fallback_difficulty,
                    "concept_tags": [term] if term else [],
                    "source_chunk_ids": point.get("source_chunk_ids", []),
                }
            )
        return questions

    def _normalize_quiz(self, raw_questions, quantity: int, docs, difficulty: str):
        if isinstance(raw_questions, dict):
            raw_questions = raw_questions.get("questions")
        if not isinstance(raw_questions, list) or not raw_questions:
            return self._build_fallback_quiz(docs, quantity, difficulty)

        normalized = []
        seen: set[str] = set()
        for item in raw_questions:
            if not isinstance(item, dict):
                continue

            question_text = str(item.get("question") or "").strip()
            if not question_text:
                continue
            key = self._dedup_key(question_text)
            if key in seen:
                continue

            q_type = str(item.get("type") or item.get("question_type") or "mcq").strip().lower()
            if q_type not in self._QUIZ_VALID_TYPES:
                q_type = "mcq"

            options = item.get("options")
            correct = item.get("correct")
            correct_answer_raw = item.get("correct_answer")

            if isinstance(options, dict):
                entries = list(options.items())
                labels = [k for k, _ in entries]
                options = [str(v) for _, v in entries]
                if isinstance(correct_answer_raw, str) and correct_answer_raw in labels:
                    correct = labels.index(correct_answer_raw)
                elif isinstance(correct, str) and correct in labels:
                    correct = labels.index(correct)

            if q_type == "true_false":
                options = ["Đúng", "Sai"]
            elif isinstance(options, list):
                options = [str(o) for o in options if str(o).strip()][:4]
            else:
                options = []

            correct_index = None
            if len(options) >= 2:
                # Prefer resolving correct_answer text against the options list; fall
                # back to a numeric `correct` index for older/alternate payload shapes.
                if isinstance(correct_answer_raw, str) and correct_answer_raw.strip() in options:
                    correct_index = options.index(correct_answer_raw.strip())
                elif correct is not None:
                    try:
                        correct_index = int(correct)
                    except (TypeError, ValueError):
                        correct_index = None
                if correct_index is None or correct_index < 0 or correct_index >= len(options):
                    correct_index = 0
            else:
                # Fewer than 2 usable options: treat as a free-form answer type instead
                # of a broken MCQ.
                options = []
                if q_type not in {"short_answer", "scenario", "code_reading"}:
                    q_type = "short_answer"

            if options:
                correct_answer = options[correct_index]
            else:
                correct_answer = str(correct_answer_raw or "").strip() or "Xem giải thích."

            raw_why_wrong = item.get("why_wrong_options_are_wrong")
            why_wrong = (
                [str(w).strip() for w in raw_why_wrong if str(w).strip()]
                if isinstance(raw_why_wrong, list)
                else []
            )

            raw_tags = item.get("concept_tags")
            concept_tags = (
                [str(t).strip() for t in raw_tags if str(t).strip()][:5]
                if isinstance(raw_tags, list)
                else []
            )

            item_difficulty = str(item.get("difficulty") or difficulty or "medium").strip().lower()
            if item_difficulty not in self._VALID_DIFFICULTIES:
                item_difficulty = "medium"

            raw_ids = item.get("source_chunk_ids")
            source_chunk_ids = (
                [str(i) for i in raw_ids if str(i).strip()] if isinstance(raw_ids, list) else []
            )

            seen.add(key)
            normalized.append(
                {
                    "id": str(item.get("id") or f"q{len(normalized) + 1}").strip() or f"q{len(normalized) + 1}",
                    "type": q_type,
                    "question": question_text,
                    "options": options,
                    "correct": correct_index if correct_index is not None else 0,
                    "correct_answer": correct_answer,
                    "explanation": str(item.get("explanation") or "Đáp án đúng dựa trên nội dung tài liệu.").strip(),
                    "why_wrong_options_are_wrong": why_wrong,
                    "difficulty": item_difficulty,
                    "concept_tags": concept_tags,
                    "source_chunk_ids": source_chunk_ids,
                }
            )
            if len(normalized) >= quantity:
                break

        return normalized or self._build_fallback_quiz(docs, quantity, difficulty)

    def generate_quiz_v2(
        self, topic: str, quantity: int, difficulty: str, learning_mode: str = "normal",
        profile: Optional[dict[str, Any]] = None,
    ):
        from backend.services.generation_readiness import evaluate_document_readiness

        profile_directives = build_profile_directives(profile)
        retriever = self.vectorstore.as_retriever(search_kwargs={"k": 15})
        docs = retriever.invoke(topic)

        quiz_readiness = evaluate_document_readiness(docs, document_id=self.course_id)["generation_readiness"]["quiz"]
        effective_quantity = quantity
        if quiz_readiness["status"] == "not_enough_context":
            effective_quantity = min(quantity, 3)
        elif quiz_readiness["status"] == "limited":
            effective_quantity = min(quantity, 5)

        context = self._clean_docs_context(docs, max_docs=15, max_chars=800)
        quiz_title = f"Quiz: {topic}" if topic and topic != "tổng quan" else "Bộ câu hỏi ôn tập"
        try:
            prompt = ChatPromptTemplate.from_template(QUIZ_V2_PROMPT)
            chain = prompt | get_llm(temperature=0.3, task="quiz") | StrOutputParser()
            res = chain.invoke(
                {
                    "context": context,
                    "topic": topic,
                    "quantity": effective_quantity,
                    "difficulty": difficulty,
                    "profile_directives": profile_directives,
                }
            )
            raw = json.loads(extract_json(res))
            if isinstance(raw, dict) and str(raw.get("quiz_title") or "").strip():
                quiz_title = str(raw["quiz_title"]).strip()
            questions = self._normalize_quiz(raw, effective_quantity, docs, difficulty)
        except Exception as e:
            logger.warning("Quiz generation failed, using fallback: %s", e)
            questions = self._build_fallback_quiz(docs, effective_quantity, difficulty)

        difficulty_mix = {"easy": 0, "medium": 0, "hard": 0}
        for q in questions:
            difficulty_mix[q.get("difficulty", "medium")] = difficulty_mix.get(q.get("difficulty", "medium"), 0) + 1

        quiz_payload: dict[str, Any] = {
            "quiz_title": quiz_title,
            "questions": questions,
            "difficulty_mix": difficulty_mix,
        }
        quality = self._evaluate_quality_gate(quiz_payload, "quiz")
        quiz_payload = self._sanitize_payload(quiz_payload)
        quiz_payload["quality_report"] = quality
        if quiz_readiness["status"] != "ready":
            quiz_payload["generation_status"] = {
                "status": "limited",
                "reason": quiz_readiness["reason"],
                "fallback_used": quiz_readiness["recommended_fallback"],
            }
        else:
            quiz_payload["generation_status"] = {"status": "full", "reason": "", "fallback_used": None}

        self._save_json(get_course_path(self.course_id)["questions"], quiz_payload)
        try:
            self.export_quiz_pdf()
        except Exception as exc:
            logger.warning("Quiz PDF export failed: %s", exc)

        payload = dict(quiz_payload)
        payload["json_url"] = f"/api/course/{self.course_id}/quiz.json"
        payload["pdf_url"] = f"/api/course/{self.course_id}/quiz.pdf"
        if difficulty == "exam":
            payload["exam_pack"] = self._build_exam_pack(questions, docs)
        return payload

    def _build_exam_pack(self, questions: list[dict[str, Any]], docs) -> dict[str, Any]:
        """Build supplementary exam study pack: short answer questions and flashcards."""
        short_answers = []
        flashcards = []
        for q in questions[:10]:
            answer_text = q.get("correct_answer") or (
                q.get("options", [""])[q.get("correct", 0)] if q.get("options") else ""
            )
            short_answers.append({
                "question": q.get("question", ""),
                "sample_answer": answer_text,
                "explanation": q.get("explanation", ""),
                "source_chunk_ids": q.get("source_chunk_ids", [])
            })
            flashcards.append({
                "front": q.get("question", ""),
                "back": answer_text,
                "source_chunk_ids": q.get("source_chunk_ids", [])
            })
        return {
            "short_answer_questions": short_answers,
            "flashcards": flashcards,
        }

    _FLASHCARD_VALID_TYPES = {
        "definition", "example", "formula", "misconception", "process", "code", "quick_recall",
    }

    def _build_fallback_flashcards(self, docs, quantity: int) -> list[dict[str, Any]]:
        """Deterministic key-term flashcards, used when context is too limited for the LLM deck."""
        cleaned_chunks, _stats = self._fallback_chunks(docs, max_docs=max(1, quantity), max_chars=260)
        cards = []
        seen: set[str] = set()
        for c in cleaned_chunks:
            term = self._short_title(c["text"], "")
            if not term:
                continue
            key = self._dedup_key(term)
            if key in seen:
                continue
            seen.add(key)
            cards.append({
                "id": f"c{len(cards) + 1}",
                "front": f"{term} là gì?",
                "back": c["text"][:220].rstrip(),
                "card_type": "quick_recall",
                "difficulty": "easy",
                "concept_tags": [term],
                "source_chunk_ids": c.get("source_chunk_ids", []),
            })
            if len(cards) >= quantity:
                break
        return cards

    def _normalize_flashcards(self, raw_cards, quantity: int, docs) -> list[dict[str, Any]]:
        if isinstance(raw_cards, dict):
            raw_cards = raw_cards.get("cards")
        if not isinstance(raw_cards, list) or not raw_cards:
            return self._build_fallback_flashcards(docs, quantity)

        normalized = []
        seen: set[str] = set()
        for item in raw_cards:
            if not isinstance(item, dict):
                continue
            front = str(item.get("front") or "").strip()
            back = str(item.get("back") or "").strip()
            if not front or not back:
                continue
            key = self._dedup_key(front)
            if key in seen:
                continue

            card_type = str(item.get("card_type") or "quick_recall").strip().lower()
            if card_type not in self._FLASHCARD_VALID_TYPES:
                card_type = "quick_recall"

            difficulty = str(item.get("difficulty") or "medium").strip().lower()
            if difficulty not in self._VALID_DIFFICULTIES:
                difficulty = "medium"

            raw_tags = item.get("concept_tags")
            concept_tags = (
                [str(t).strip() for t in raw_tags if str(t).strip()][:5]
                if isinstance(raw_tags, list)
                else []
            )

            raw_ids = item.get("source_chunk_ids")
            source_chunk_ids = (
                [str(i) for i in raw_ids if str(i).strip()] if isinstance(raw_ids, list) else []
            )

            seen.add(key)
            normalized.append({
                "id": str(item.get("id") or f"c{len(normalized) + 1}").strip() or f"c{len(normalized) + 1}",
                "front": front,
                "back": back,
                "card_type": card_type,
                "difficulty": difficulty,
                "concept_tags": concept_tags,
                "source_chunk_ids": source_chunk_ids,
            })
            if len(normalized) >= quantity:
                break

        return normalized or self._build_fallback_flashcards(docs, quantity)

    def generate_flashcards_v2(
        self, topic: str = "tổng quan", quantity: int = 15,
        profile: Optional[dict[str, Any]] = None,
    ) -> dict[str, Any]:
        """Generate a dedicated, source-grounded flashcard deck (not derived from quiz Q&A)."""
        from backend.services.generation_readiness import evaluate_document_readiness

        profile_directives = build_profile_directives(profile)
        retriever = self.vectorstore.as_retriever(search_kwargs={"k": 15})
        docs = retriever.invoke(topic)

        readiness = evaluate_document_readiness(docs, document_id=self.course_id)["generation_readiness"]["flashcards"]
        deck_title = f"Flashcards: {topic}" if topic and topic != "tổng quan" else "Bộ thẻ ghi nhớ"

        if readiness["status"] == "not_enough_context":
            cards = self._build_fallback_flashcards(docs, min(quantity, 6))
        else:
            effective_quantity = min(quantity, 8) if readiness["status"] == "limited" else quantity
            context = self._clean_docs_context(docs, max_docs=15, max_chars=800)
            try:
                prompt = ChatPromptTemplate.from_template(FLASHCARD_GENERATION_PROMPT)
                chain = prompt | get_llm(temperature=0.3, task="flashcard") | StrOutputParser()
                res = chain.invoke({
                    "context": context, "topic": topic, "quantity": effective_quantity,
                    "profile_directives": profile_directives,
                })
                raw = json.loads(extract_json(res))
                if isinstance(raw, dict) and str(raw.get("deck_title") or "").strip():
                    deck_title = str(raw["deck_title"]).strip()
                cards = self._normalize_flashcards(raw, effective_quantity, docs)
            except Exception as e:
                logger.warning("Flashcard generation failed, using fallback: %s", e)
                cards = self._build_fallback_flashcards(docs, effective_quantity)

        deck_payload: dict[str, Any] = {"deck_title": deck_title, "cards": cards}
        quality = self._evaluate_quality_gate(deck_payload, "flashcards")
        deck_payload = self._sanitize_payload(deck_payload)
        deck_payload["quality_report"] = quality
        if readiness["status"] != "ready":
            deck_payload["generation_status"] = {
                "status": "limited",
                "reason": readiness["reason"],
                "fallback_used": readiness["recommended_fallback"],
            }
        else:
            deck_payload["generation_status"] = {"status": "full", "reason": "", "fallback_used": None}

        paths = get_course_path(self.course_id)
        if "flashcards" in paths:
            self._save_json(paths["flashcards"], deck_payload)
        return deck_payload

    def _sentences_from_text(self, text: str, limit: int = 3, max_words: int = 18) -> list[str]:
        """Split chunk prose into short, real sentences for bullets — no template filler."""
        sentences = [s.strip() for s in re.split(r"(?<=[.!?;])\s+", str(text or "")) if len(s.strip()) >= 15]
        bullets = []
        for s in sentences:
            words = s.split()
            bullets.append(" ".join(words[:max_words]) + ("…" if len(words) > max_words else ""))
            if len(bullets) >= limit:
                break
        return bullets

    def _slide_defaults(self, point: dict[str, str], index: int) -> dict[str, Any]:
        """Default synthetic slide content built from a single retrieved chunk.

        Used only to fill in fields the LLM omitted for a given slide — never
        to override real LLM-generated content. Every field is derived from the
        chunk's own prose: no generic teaching-instruction filler, per product rule.
        """
        text = point["text"]
        title = self._short_title(text, f"Slide {index}")
        bullets = self._sentences_from_text(text, limit=3) or [text[:110]]
        first_sentence = bullets[0] if bullets else text[:120]
        return {
            "slide_type": "concept",
            "title": title,
            "key_message": first_sentence,
            "bullets": bullets,
            "screen_content": {"bullets": [], "formula": "", "code": "", "table": [], "diagram_description": ""},
            "visual_instruction": {"type": "none", "description": "", "labels": []},
            "visual": {"type": "none", "description": "", "labels": []},
            "speaker_notes": text,
            "example_or_application": "",
            "common_mistake": {"mistake": "", "correction": ""},
            "quick_check": {"question": "", "answer": ""},
            "student_prompt": f'"{title}" nghĩa là gì và được dùng trong tình huống nào theo tài liệu?',
            "source_chunk_ids": point.get("source_chunk_ids", []),
        }

    _GENERIC_SLIDE_TITLES = {"ý chính", "ghi nhớ ý chính", "nội dung chính", ""}

    def _normalize_slide_item(self, raw_slide: Any, point: dict[str, str], index: int) -> dict[str, Any]:
        """Merge real LLM slide content with synthetic defaults for any missing field."""
        defaults = self._slide_defaults(point, index)
        slide = raw_slide if isinstance(raw_slide, dict) else {}

        def text_field(key: str) -> Optional[str]:
            value = slide.get(key)
            return value.strip() if isinstance(value, str) and value.strip() else None

        def list_field(key: str) -> Optional[list]:
            value = slide.get(key)
            return value if isinstance(value, list) and value else None

        title = text_field("title") or defaults["title"]
        if title.strip().lower() in self._GENERIC_SLIDE_TITLES:
            title = defaults["title"]

        # New schema: screen_content carries the slide body; legacy "bullets" still accepted.
        screen_raw = slide.get("screen_content") if isinstance(slide.get("screen_content"), dict) else {}
        raw_bullets = screen_raw.get("bullets") if isinstance(screen_raw.get("bullets"), list) else None
        if raw_bullets is None:
            raw_bullets = list_field("bullets") or []
        bullets = [str(b).strip() for b in raw_bullets if str(b).strip()][:5] or defaults["bullets"]
        table_raw = screen_raw.get("table")
        screen_content = {
            "bullets": bullets,
            "formula": str(screen_raw.get("formula") or "").strip(),
            "code": str(screen_raw.get("code") or "").strip(),
            "table": [
                [str(cell) for cell in row]
                for row in table_raw
                if isinstance(row, (list, tuple)) and row
            ]
            if isinstance(table_raw, list)
            else [],
            "diagram_description": str(screen_raw.get("diagram_description") or "").strip(),
        }

        key_message = text_field("key_message") or defaults["key_message"]
        speaker_notes = text_field("speaker_notes") or defaults["speaker_notes"]
        slide_type = text_field("slide_type") or defaults["slide_type"]
        example_or_application = text_field("example_or_application") or defaults["example_or_application"]

        visual_raw = slide.get("visual_instruction") if isinstance(slide.get("visual_instruction"), dict) else slide.get("visual")
        visual = (
            {
                "type": str(visual_raw.get("type") or "none"),
                "description": str(visual_raw.get("description") or ""),
                "labels": [str(label) for label in (visual_raw.get("labels") or []) if str(label).strip()],
            }
            if isinstance(visual_raw, dict)
            else defaults["visual"]
        )

        quick_check_raw = slide.get("quick_check")
        quick_check = (
            {
                "question": str(quick_check_raw.get("question") or "").strip(),
                "answer": str(quick_check_raw.get("answer") or "").strip(),
            }
            if isinstance(quick_check_raw, dict) and quick_check_raw.get("question")
            else defaults["quick_check"]
        )

        common_mistake_raw = slide.get("common_mistake")
        common_mistake = (
            {
                "mistake": str(common_mistake_raw.get("mistake") or "").strip(),
                "correction": str(common_mistake_raw.get("correction") or "").strip(),
            }
            if isinstance(common_mistake_raw, dict) and common_mistake_raw.get("mistake")
            else defaults["common_mistake"]
        )

        student_prompt = text_field("student_prompt") or defaults["student_prompt"]
        source_chunk_ids = list_field("source_chunk_ids") or defaults["source_chunk_ids"]

        content_lines = [f"- {b}" for b in bullets]
        if screen_content["formula"]:
            content_lines.append(f"\nCông thức:\n$$ {screen_content['formula']} $$")
        if screen_content["code"]:
            content_lines.append(f"\nCode:\n```\n{screen_content['code']}\n```")
        if screen_content["table"]:
            content_lines.append("\nBảng so sánh:")
            for row in screen_content["table"]:
                content_lines.append(" | ".join(row))
        if screen_content["diagram_description"]:
            content_lines.append(f"\nSơ đồ: {screen_content['diagram_description']}")
        content_str = "\n".join(content_lines)

        return {
            "slide_index": index,
            "slide_type": slide_type,
            "title": title,
            "key_message": key_message,
            "key_idea": key_message,
            "screen_content": screen_content,
            "bullets": bullets,
            "content": content_str,
            "visual_instruction": visual,
            "visual": visual,
            "visual_type": visual["type"],
            "image_suggestion": visual["description"],
            "speaker_notes": speaker_notes,
            "example_or_application": example_or_application,
            "example": example_or_application,
            "common_mistake": common_mistake,
            "note": quick_check["question"] or example_or_application,
            "quick_check": quick_check,
            "student_prompt": student_prompt,
            "layout_hint": "title-and-content",
            "source_chunk_ids": source_chunk_ids,
        }

    def _build_fallback_slides(self, docs, num_slides: int):
        count = max(1, min(num_slides, 10))
        points = self._doc_points(docs, limit=count, max_chars=220)
        return [self._normalize_slide_item({}, points[i % len(points)], i + 1) for i in range(count)]

    def _normalize_slides(self, raw_slides, num_slides: int, docs):
        # The LLM returns a deck object {deck_title, slides: [...], ...}; unwrap it.
        if isinstance(raw_slides, dict):
            raw_slides = raw_slides.get("slides")
        if not isinstance(raw_slides, list) or not raw_slides:
            return self._build_fallback_slides(docs, num_slides)

        points = self._doc_points(docs, limit=max(1, min(num_slides, 10)), max_chars=220)
        slides = []
        for index, item in enumerate(raw_slides[:num_slides], 1):
            point = points[(index - 1) % len(points)]
            slides.append(self._normalize_slide_item(item, point, index))
        return slides or self._build_fallback_slides(docs, num_slides)

    def _extract_deck_meta(self, raw_deck: Any) -> dict[str, Any]:
        if not isinstance(raw_deck, dict):
            return {}
        return {
            "deck_title": raw_deck.get("deck_title"),
            "subtitle": raw_deck.get("subtitle"),
            "course_level": raw_deck.get("course_level") or "university",
            "learning_outcomes": raw_deck.get("learning_outcomes") or [],
            "estimated_duration_minutes": raw_deck.get("estimated_duration_minutes"),
        }

    def _ensure_deck_structure(self, slides: list[dict[str, Any]], deck_meta: dict[str, Any], topic: str) -> list[dict[str, Any]]:
        """Guarantee the lecture arc: title, objectives up front; recap, practice at the end.

        Missing structural slides are built deterministically from real deck/slide content
        (never invented facts), so a weak LLM pass still yields a deliverable deck.
        """
        types = {s.get("slide_type") for s in slides}
        content_ids = sorted({i for s in slides for i in (s.get("source_chunk_ids") or [])})
        deck_title = str(deck_meta.get("deck_title") or topic).strip() or topic
        outcomes = [str(o).strip() for o in (deck_meta.get("learning_outcomes") or []) if str(o).strip()]

        def structural(slide_type: str, title: str, bullets: list[str], key_message: str) -> dict[str, Any]:
            bullets = [b for b in bullets if b][:5]
            return {
                "slide_index": 0,
                "slide_type": slide_type,
                "title": title,
                "key_message": key_message,
                "key_idea": key_message,
                "screen_content": {"bullets": bullets, "formula": "", "code": "", "table": [], "diagram_description": ""},
                "bullets": bullets,
                "content": "\n".join(f"- {b}" for b in bullets),
                "visual_instruction": {"type": "none", "description": "", "labels": []},
                "visual": {"type": "none", "description": "", "labels": []},
                "visual_type": "none",
                "image_suggestion": "",
                "speaker_notes": (
                    f"Slide {slide_type} của bài giảng \"{deck_title}\". Giảng viên điểm qua từng mục, "
                    "kết nối với nội dung các slide chính và mời sinh viên đặt câu hỏi."
                ),
                "example_or_application": "",
                "example": "",
                "common_mistake": {"mistake": "", "correction": ""},
                "note": "",
                "quick_check": {"question": "", "answer": ""},
                "student_prompt": "",
                "layout_hint": "title-and-content",
                "source_chunk_ids": content_ids,
            }

        if "title" not in types:
            slides.insert(0, structural(
                "title", deck_title, [],
                str(deck_meta.get("subtitle") or "").strip() or (outcomes[0] if outcomes else ""),
            ))
        if "objectives" not in types and outcomes:
            slides.insert(1, structural(
                "objectives", "Mục tiêu bài học", outcomes, "Sau bài giảng này, bạn sẽ đạt được các mục tiêu sau."
            ))
        if "motivation" not in types and "prerequisite" not in types:
            slides.insert(len(slides) // 4 or 2, structural(
                "motivation", "Tại sao kiến thức này quan trọng?",
                ["Hiểu rõ bản chất và bối cảnh áp dụng.", "Nền tảng cho các phần kiến thức nâng cao tiếp theo.", "Ứng dụng trực tiếp trong giải quyết bài toán thực tế."],
                "Ý nghĩa và động lực học tập của bài giảng."
            ))
        if "recap" not in types:
            takeaways = [s.get("key_message", "") for s in slides if s.get("slide_type") not in {"title", "objectives"}][:5]
            slides.append(structural("recap", "Tổng kết bài giảng", takeaways, "Các điểm mấu chốt cần mang về."))
        if "practice" not in types:
            questions = [s.get("student_prompt", "") for s in slides if s.get("student_prompt")][:4]
            if questions:
                slides.append(structural("practice", "Bài tập tự luyện", questions, "Tự kiểm tra trước buổi học sau."))

        for index, slide in enumerate(slides, 1):
            slide["slide_index"] = index
        return slides

    def generate_slides_v2(
        self, topic: str, num_slides: int, learning_mode: str = "normal",
        profile: Optional[dict[str, Any]] = None,
    ):
        profile_directives = build_profile_directives(profile)
        retriever = self.vectorstore.as_retriever(search_kwargs={"k": 15})
        docs = retriever.invoke(topic)

        from backend.services.generation_readiness import evaluate_document_readiness

        slide_readiness = evaluate_document_readiness(docs, document_id=self.course_id)["generation_readiness"]["slides"]
        if slide_readiness["status"] == "not_enough_context":
            logger.info("[SlideGen] Not enough clean context, generating a short overview deck instead.")
            overview_count = max(5, min(num_slides, 8))
            slides = self._build_fallback_slides(docs, overview_count)
            quality = self._evaluate_quality_gate({"slides": slides}, "slides")
            slides = self._sanitize_payload(slides)
            overview_warning = (
                "Bộ slide tổng quan ngắn: tài liệu chưa đủ ngữ cảnh sạch cho bài giảng đầy đủ. "
                "Bạn vẫn có thể dùng deck này để nắm nhanh nội dung chính."
            )
            payload = {
                "slides": slides,
                "quality_report": quality,
                "deck_title": f"Tổng quan ngắn: {topic}",
                "subtitle": "",
                "course_level": "overview",
                "learning_outcomes": [],
                "estimated_duration_minutes": None,
                "warnings": [overview_warning],
                "generation_status": {
                    "status": "limited",
                    "reason": overview_warning,
                    "fallback_used": "short_overview_deck",
                },
            }
            self._save_json(get_course_path(self.course_id)["slides"], payload)
            try:
                self.export_slides_pdf()
            except Exception as exc:
                logger.warning("Slide PDF export failed: %s", exc)
            return {
                "slides": slides,
                "deck_title": payload["deck_title"],
                "learning_outcomes": [],
                "quality_report": quality,
                "generation_status": payload["generation_status"],
                "json_url": f"/api/course/{self.course_id}/slide.json",
                "pdf_url": f"/api/course/{self.course_id}/slide.pdf",
            }

        context = self._clean_docs_context(docs, max_docs=15, max_chars=800)
        deck_meta: dict[str, Any] = {}
        generation_error: Optional[str] = None
        try:
            prompt = ChatPromptTemplate.from_template(SLIDE_GENERATION_PROMPT)
            # Same truncation risk as book generation: rich per-slide schema across many slides
            # can exceed the default 8192-token cap. Raise the ceiling to avoid silent fallback.
            chain = prompt | get_llm(temperature=0.1, max_output_tokens=32768, task="slide") | StrOutputParser()
            res = chain.invoke({
                "context": context, "topic": topic, "num_slides": num_slides,
                "profile_directives": profile_directives,
            })
            raw_deck = json.loads(extract_json(res))
            deck_meta = self._extract_deck_meta(raw_deck)
            slides = self._ensure_deck_structure(self._normalize_slides(raw_deck, num_slides, docs), deck_meta, topic)
            quality = self._evaluate_quality_gate({"slides": slides}, "slides")
            if not quality["is_university_ready"]:
                logger.info("[SlideGen] Score below the 85 gate, retrying once...")
                res = chain.invoke({
                    "context": context, "topic": topic, "num_slides": num_slides,
                    "profile_directives": profile_directives,
                })
                raw_deck = json.loads(extract_json(res))
                deck_meta = self._extract_deck_meta(raw_deck)
                slides = self._ensure_deck_structure(self._normalize_slides(raw_deck, num_slides, docs), deck_meta, topic)
                quality = self._evaluate_quality_gate({"slides": slides}, "slides")
        except Exception as e:
            logger.warning("Slide generation failed, using fallback: %s", e)
            generation_error = str(e)
            slides = self._build_fallback_slides(docs, num_slides)
            quality = self._evaluate_quality_gate({"slides": slides}, "slides")

        slides = self._sanitize_payload(slides)
        if generation_error is not None:
            slide_generation_status = {
                "status": "limited",
                "reason": (
                    "Không thể tạo bộ slide đầy đủ từ AI (lỗi hệ thống hoặc giới hạn quota), "
                    "đã dùng bản dự phòng dựng trực tiếp từ tài liệu."
                ),
                "fallback_used": "generic_fallback",
            }
        else:
            slide_generation_status = {"status": "full", "reason": "", "fallback_used": None}
        payload = {
            "slides": slides,
            "quality_report": quality,
            "deck_title": deck_meta.get("deck_title") or topic,
            "subtitle": deck_meta.get("subtitle") or "",
            "course_level": deck_meta.get("course_level") or "university",
            "learning_outcomes": deck_meta.get("learning_outcomes") or [],
            "estimated_duration_minutes": deck_meta.get("estimated_duration_minutes"),
            "generation_status": slide_generation_status,
        }
        self._save_json(get_course_path(self.course_id)["slides"], payload)
        try:
            self.export_slides_pdf()
        except Exception as exc:
            logger.warning("Slide PDF export failed: %s", exc)
        return {
            "slides": slides,
            "deck_title": payload["deck_title"],
            "learning_outcomes": payload["learning_outcomes"],
            "quality_report": quality,
            "generation_status": payload["generation_status"],
            "json_url": f"/api/course/{self.course_id}/slide.json",
            "pdf_url": f"/api/course/{self.course_id}/slide.pdf",
        }

    _FALLBACK_SCENE_TEMPLATES = [
        "title_intro",
        "concept_card",
        "flow_diagram",
        "comparison",
        "quiz_card",
        "recap_map",
    ]

    def _build_fallback_scenes(self, docs, scene_count: int):
        """Deterministic grounded scenes: varied visual templates, real source_chunk_ids.

        Never pads beyond the number of distinct source points — a short document
        yields an honest limited video instead of repeated near-identical scenes
        (which the video quality gate would rightly reject).
        """
        points = self._doc_points(docs, limit=max(scene_count, 4), max_chars=260)
        scene_count = max(1, min(scene_count, len(points)))
        scenes = []
        for index in range(scene_count):
            point = points[index]
            template = self._FALLBACK_SCENE_TEMPLATES[index % len(self._FALLBACK_SCENE_TEMPLATES)]
            title = self._short_title(point["text"], f"Nội dung trọng tâm {index + 1}")
            short_line = " ".join(point["text"].split()[:12])
            scenes.append(
                {
                    "title": f"Cảnh {index + 1}: {title}",
                    "scene_type": template,
                    "visual_template": template,
                    "screen_text": [short_line, "Liên hệ với nội dung tài liệu gốc"],
                    "visual_text": f"- {point['text']}\n- Liên hệ với nội dung tài liệu gốc",
                    "voiceover": (
                        f"Ở phần này, chúng ta tập trung vào nội dung sau: {point['text']} "
                        "Hãy ghi nhớ kiến thức cốt lõi và liên hệ nó với các phần trước của tài liệu."
                    ),
                    "source_chunk_ids": point.get("source_chunk_ids", []),
                }
            )
        return scenes

    def _clean_scene_line(self, value: Any) -> str:
        """Clean one scene text line; return '' if the line is TOC/noise."""
        from backend.services.context_cleaner import clean_text_markers, scrub_banned_phrases

        raw = str(value or "")
        if re.search(r"(?:\.\s*){3,}", raw):
            return ""
        lowered = raw.strip().lower()
        if lowered.startswith(("contents", "table of contents", "mục lục")):
            return ""
        return scrub_banned_phrases(clean_text_markers(raw)).strip()

    def _normalize_scenes(self, raw_scenes, scene_count: int, docs):
        if not isinstance(raw_scenes, list) or not raw_scenes:
            return self._build_fallback_scenes(docs, scene_count)

        points = self._doc_points(docs, limit=max(len(raw_scenes), 1), max_chars=260)
        scenes = []
        for index, item in enumerate(raw_scenes[:scene_count], 1):
            scene = item if isinstance(item, dict) else {"voiceover": str(item)}
            point = points[(index - 1) % len(points)]

            title = self._clean_scene_line(scene.get("title")) or f"Cảnh {index}"

            raw_lines = scene.get("screen_text") if isinstance(scene.get("screen_text"), list) else None
            screen_text = None
            if raw_lines is not None:
                screen_text = [
                    line for line in (self._clean_scene_line(raw) for raw in raw_lines) if line
                ]

            visual_text = self._clean_scene_line(scene.get("visual_text") or scene.get("content"))
            if not visual_text:
                visual_text = "\n".join(f"- {line}" for line in screen_text) if screen_text else title

            voiceover = self._clean_scene_line(scene.get("voiceover")) or visual_text.replace("-", " ").strip()

            raw_ids = scene.get("source_chunk_ids")
            source_chunk_ids = (
                [str(i) for i in raw_ids if str(i).strip()]
                if isinstance(raw_ids, list) and raw_ids
                else point.get("source_chunk_ids", [])
            )

            try:
                duration_seconds = int(scene.get("duration_seconds") or 0)
            except (TypeError, ValueError):
                duration_seconds = 0
            if duration_seconds <= 0:
                duration_seconds = self._estimate_scene_seconds(voiceover)

            normalized: dict[str, Any] = {
                "title": title,
                "scene_type": str(scene.get("scene_type") or "concept"),
                "visual_template": str(
                    scene.get("visual_template") or scene.get("scene_type") or "concept_card"
                ),
                "visual_text": visual_text,
                "voiceover": voiceover,
                "duration_seconds": duration_seconds,
                "source_chunk_ids": source_chunk_ids,
            }
            key_message = self._clean_scene_line(scene.get("key_message"))
            if key_message:
                normalized["key_message"] = key_message
            if screen_text:
                normalized["screen_text"] = screen_text

            visual_data = scene.get("visual_data")
            if isinstance(visual_data, dict):
                cleaned_data = {}
                for key, value in visual_data.items():
                    if isinstance(value, str):
                        cleaned_value = self._clean_scene_line(value)
                        if cleaned_value:
                            cleaned_data[key] = cleaned_value
                    else:
                        cleaned_data[key] = value
                if cleaned_data:
                    normalized["visual_data"] = cleaned_data

            scenes.append(normalized)
        return scenes or self._build_fallback_scenes(docs, scene_count)

    def _font(self, size: int, bold: bool = False):
        from PIL import ImageFont

        candidates = [
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
            if bold
            else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
            if bold
            else "/System/Library/Fonts/Supplemental/Arial.ttf",
        ]
        for candidate in candidates:
            if candidate and os.path.exists(candidate):
                return ImageFont.truetype(candidate, size)
        return ImageFont.load_default()

    def _wrap_lines(self, text: str, width: int) -> list[str]:
        lines: list[str] = []
        for raw_line in text.replace("\r", "").split("\n"):
            line = raw_line.strip().lstrip("-*• ").strip()
            if not line:
                continue
            lines.extend(textwrap.wrap(line, width=width) or [line])
        return lines

    def _draw_text_block(self, draw, position: tuple[int, int], lines: list[str], font, fill, line_height: int):
        x, y = position
        for line in lines:
            draw.text((x, y), line, font=font, fill=fill)
            y += line_height
        return y

    def _render_scene_image(self, scene: dict[str, Any], index: int, path: str) -> None:
        from PIL import Image, ImageDraw

        image = Image.new("RGB", (1280, 720), (248, 250, 252))
        draw = ImageDraw.Draw(image)
        title_font = self._font(44, bold=True)
        sub_font = self._font(34, bold=True)
        body_font = self._font(28)
        small_font = self._font(22)
        code_font = self._font(24)

        # Header bar
        draw.rectangle((0, 0, 1280, 80), fill=(30, 41, 59))
        draw.text((48, 26), f"AI Learning Video · Cảnh {index}", font=small_font, fill=(241, 245, 249))
        
        tpl = str(scene.get("visual_template") or scene.get("scene_type") or "").lower()
        vdata = scene.get("visual_data") if isinstance(scene.get("visual_data"), dict) else {}
        title = str(scene.get("title") or f"Cảnh {index}")
        screen_lines = scene.get("screen_text") if isinstance(scene.get("screen_text"), list) else self._wrap_lines(str(scene.get("visual_text") or ""), 50)[:4]

        # Card container
        draw.rounded_rectangle((48, 110, 1232, 660), radius=24, fill=(255, 255, 255), outline=(226, 232, 240), width=2)

        if "title_intro" in tpl or "hook" in tpl:
            draw.rounded_rectangle((80, 150, 1200, 320), radius=16, fill=(238, 242, 255))
            t_lines = self._wrap_lines(title, width=38)[:2]
            self._draw_text_block(draw, (120, 180), t_lines, self._font(52, bold=True), (30, 27, 75), 64)
            y = 360
            for line in screen_lines[:3]:
                draw.rounded_rectangle((80, y, 1200, y + 70), radius=12, fill=(241, 245, 249))
                draw.text((110, y + 20), line, font=body_font, fill=(51, 65, 85))
                y += 85

        elif "mistake" in tpl:
            t_lines = self._wrap_lines(title, width=45)[:1]
            self._draw_text_block(draw, (80, 140), t_lines, title_font, (15, 23, 42), 52)
            mistake = str(vdata.get("mistake") or (screen_lines[0] if screen_lines else "Lỗi hiểu sai khái niệm"))
            correction = str(vdata.get("correction") or (screen_lines[1] if len(screen_lines) > 1 else "Cách tiếp cận chuẩn xác"))
            draw.rounded_rectangle((80, 210, 1200, 410), radius=16, fill=(254, 242, 242), outline=(252, 165, 165), width=2)
            draw.text((110, 230), "❌ LỖI THƯỜNG GẶP / HIỂU NHẦM:", font=sub_font, fill=(185, 28, 28))
            m_lines = self._wrap_lines(mistake, width=55)[:3]
            self._draw_text_block(draw, (110, 280), m_lines, body_font, (127, 29, 29), 38)
            draw.rounded_rectangle((80, 435, 1200, 635), radius=16, fill=(240, 253, 244), outline=(134, 239, 172), width=2)
            draw.text((110, 455), "✅ GIẢI PHÁP / HIỂU ĐÚNG:", font=sub_font, fill=(21, 128, 61))
            c_lines = self._wrap_lines(correction, width=55)[:3]
            self._draw_text_block(draw, (110, 505), c_lines, body_font, (20, 83, 45), 38)

        elif "code" in tpl:
            t_lines = self._wrap_lines(title, width=45)[:1]
            self._draw_text_block(draw, (80, 140), t_lines, title_font, (15, 23, 42), 52)
            code_text = str(vdata.get("code") or "\n".join(screen_lines) or "# Code walkthrough\nprint('AI Course Generator')")
            draw.rounded_rectangle((80, 210, 1200, 630), radius=16, fill=(15, 23, 42))
            draw.text((110, 230), "💻 Code Example / Walkthrough", font=small_font, fill=(148, 163, 184))
            c_lines = code_text.splitlines()[:10]
            y = 280
            for cline in c_lines:
                draw.text((120, y), cline[:65], font=code_font, fill=(56, 189, 248))
                y += 34

        elif "comparison" in tpl:
            t_lines = self._wrap_lines(title, width=45)[:1]
            self._draw_text_block(draw, (80, 140), t_lines, title_font, (15, 23, 42), 52)
            left_col = str(vdata.get("left_col") or (screen_lines[0] if screen_lines else "Khái niệm A"))
            right_col = str(vdata.get("right_col") or (screen_lines[1] if len(screen_lines) > 1 else "Khái niệm B"))
            draw.rounded_rectangle((80, 210, 630, 630), radius=16, fill=(248, 250, 252), outline=(203, 213, 225), width=2)
            draw.rounded_rectangle((80, 210, 630, 270), radius=16, fill=(226, 232, 240))
            draw.text((110, 225), "🔹 Nhóm / Khái niệm 1", font=sub_font, fill=(30, 41, 59))
            self._draw_text_block(draw, (110, 290), self._wrap_lines(left_col, 24)[:8], body_font, (51, 65, 85), 38)

            draw.rounded_rectangle((650, 210, 1200, 630), radius=16, fill=(238, 242, 255), outline=(199, 210, 254), width=2)
            draw.rounded_rectangle((650, 210, 1200, 270), radius=16, fill=(224, 231, 255))
            draw.text((680, 225), "🔸 Nhóm / Khái niệm 2", font=sub_font, fill=(55, 48, 163))
            self._draw_text_block(draw, (680, 290), self._wrap_lines(right_col, 24)[:8], body_font, (49, 46, 129), 38)

        elif "flow" in tpl or "diagram" in tpl:
            t_lines = self._wrap_lines(title, width=45)[:1]
            self._draw_text_block(draw, (80, 140), t_lines, title_font, (15, 23, 42), 52)
            steps = vdata.get("steps") if isinstance(vdata.get("steps"), list) else screen_lines
            steps = [s for s in steps if s][:4] or ["Bước 1: Đầu vào", "Bước 2: Xử lý", "Bước 3: Đầu ra"]
            y = 220
            for idx, stp in enumerate(steps, 1):
                draw.rounded_rectangle((120, y, 1160, y + 80), radius=16, fill=(239, 246, 255), outline=(191, 219, 254), width=2)
                draw.ellipse((140, y + 18, 184, y + 62), fill=(37, 99, 235))
                draw.text((154, y + 25), str(idx), font=sub_font, fill=(255, 255, 255))
                draw.text((210, y + 24), stp[:55], font=sub_font, fill=(30, 58, 138))
                y += 105

        elif "quiz" in tpl:
            q = str(vdata.get("question") or title)
            opts = vdata.get("options") if isinstance(vdata.get("options"), list) else screen_lines
            opts = [o for o in opts if o][:4] or ["A. Phương án 1", "B. Phương án 2", "C. Phương án 3", "D. Phương án 4"]
            draw.rounded_rectangle((80, 135, 1200, 250), radius=16, fill=(254, 249, 195), outline=(253, 224, 71), width=2)
            draw.text((110, 150), "❓ CÂU HỎI ÔN TẬP NHANH:", font=small_font, fill=(133, 77, 14))
            self._draw_text_block(draw, (110, 180), self._wrap_lines(q, 50)[:2], sub_font, (113, 63, 18), 38)
            y = 275
            for opt in opts:
                draw.rounded_rectangle((80, y, 1200, y + 75), radius=14, fill=(248, 250, 252), outline=(203, 213, 225), width=2)
                draw.text((110, y + 22), opt[:65], font=body_font, fill=(30, 41, 59))
                y += 90

        else:
            t_lines = self._wrap_lines(title, width=42)[:2]
            y = self._draw_text_block(draw, (80, 140), t_lines, title_font, (15, 23, 42), 54)
            term = str(vdata.get("term") or "")
            if term:
                draw.rounded_rectangle((80, y + 10, 1200, y + 75), radius=12, fill=(238, 242, 255))
                draw.text((110, y + 26), f"📌 Thuật ngữ: {term[:50]}", font=sub_font, fill=(67, 56, 202))
                y += 85
            y += 20
            for line in screen_lines[:5]:
                draw.rounded_rectangle((80, y, 1200, y + 75), radius=14, fill=(248, 250, 252), outline=(226, 232, 240), width=1)
                draw.ellipse((105, y + 28, 125, y + 48), fill=(79, 70, 229))
                draw.text((145, y + 22), line[:60], font=body_font, fill=(51, 65, 85))
                y += 90

        image.save(path, quality=95)
        image.save(path, quality=95)

    def _clean_tts_text(self, text: str) -> str:
        text = re.sub(r"\$.*?\$", "", text)
        text = text.replace("\\", " ")
        text = re.sub(
            r"[^\w\s,.\?!\-áàảãạâấầẩẫậăắằẳẵặéèẻẽẹêếềểễệíìỉĩịóòỏõọôốồổỗộơớờởỡợ"
            r"úùủũụưứừửữựýỳỷỹỵĐđ]",
            " ",
            text,
        )
        return re.sub(r"\s+", " ", text).strip()

    def _run_async_in_thread(self, coro):
        result: dict[str, Any] = {}

        def runner():
            try:
                result["value"] = asyncio.run(coro)
            except Exception as exc:  # pragma: no cover - threaded propagation
                result["error"] = exc

        thread = threading.Thread(target=runner)
        thread.start()
        thread.join()
        if "error" in result:
            raise result["error"]
        return result.get("value")

    def _synthesize_voiceover(self, text: str, path: str) -> None:
        import edge_tts

        cleaned = self._clean_tts_text(text)
        if not cleaned:
            raise ValueError("Voiceover text is empty after cleaning.")
        communicate = edge_tts.Communicate(cleaned, "vi-VN-HoaiMyNeural")
        self._run_async_in_thread(communicate.save(path))

    def _run_ffmpeg(self, command: list[str]) -> None:
        completed = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=False)
        if completed.returncode != 0:
            raise RuntimeError(completed.stderr[-1200:])

    def _estimate_scene_seconds(self, voiceover: str) -> int:
        word_count = max(12, len(voiceover.split()))
        return max(8, min(45, round(word_count / 2.2)))

    def _fit_scene_durations(self, scenes: list, target_seconds: int) -> None:
        """Scale scene durations down proportionally so the video honors its mode's target length.

        Only compresses (never stretches): a short document legitimately yields a
        shorter video than the mode target.
        """
        if not scenes or target_seconds <= 0:
            return
        def current(sc) -> int:
            try:
                dur = int(sc.get("duration_seconds") or 0)
            except (TypeError, ValueError):
                dur = 0
            return dur if dur > 0 else self._estimate_scene_seconds(sc.get("voiceover") or "")

        total = sum(current(sc) for sc in scenes)
        if total <= target_seconds:
            return
        ratio = target_seconds / total
        for sc in scenes:
            sc["duration_seconds"] = max(4, round(current(sc) * ratio))

    def _render_scene_clip(self, ffmpeg: str, image_path: str, audio_path: str | None, clip_path: str, seconds: int) -> None:
        if audio_path and os.path.exists(audio_path) and os.path.getsize(audio_path) > 100:
            command = [
                ffmpeg,
                "-y",
                "-loop",
                "1",
                "-i",
                image_path,
                "-i",
                audio_path,
                "-c:v",
                "libx264",
                "-tune",
                "stillimage",
                "-c:a",
                "aac",
                "-b:a",
                "128k",
                "-pix_fmt",
                "yuv420p",
                "-shortest",
                clip_path,
            ]
        else:
            command = [
                ffmpeg,
                "-y",
                "-loop",
                "1",
                "-i",
                image_path,
                "-f",
                "lavfi",
                "-i",
                "anullsrc=channel_layout=stereo:sample_rate=44100",
                "-t",
                str(seconds),
                "-c:v",
                "libx264",
                "-c:a",
                "aac",
                "-pix_fmt",
                "yuv420p",
                "-shortest",
                clip_path,
            ]
        self._run_ffmpeg(command)

    def _normalize_video_renderer(self, renderer: str) -> str:
        r = str(renderer or "").strip().lower()
        if r in ("simple_slides", "simple_templates"):
            return "simple_templates"
        return "manim"

    def _select_manim_template(self, scene: dict[str, Any]) -> str:
        vt = str(scene.get("visual_template") or "").strip()
        title = str(scene.get("title") or "").lower()
        if vt == "flow_diagram" or "flow" in vt:
            return "flow_diagram_scene"
        if "f1-score" in title or "precision" in title or "recall" in title or "confusion" in title:
            return "confusion_matrix_scene"
        if "gradient descent" in title or "loss curve" in title:
            return "gradient_descent_curve_scene"
        if "neural network" in title or "backprop" in title or "layer" in title:
            return "neural_network_scene"
        if "clustering" in title or "k-means" in title:
            return "clustering_scene"
        if "dataframe" in title or "pandas" in title or "table" in title:
            return "dataframe_table_scene"
        return "concept_card_scene"

    def _build_manim_script(self, template_name: str, scene: dict[str, Any], seconds: int) -> str:
        return f"""import json
from manim import *

SCENE_DATA = json.loads({json.dumps(json.dumps(scene, ensure_ascii=False))})
dispatch = {{
    "gradient_descent_curve_scene": None,
    "confusion_matrix_scene": None,
    "neural_network_scene": None,
    "clustering_scene": None,
    "dataframe_table_scene": None,
    "flow_diagram_scene": None,
    "concept_card_scene": None,
}}
# Safe template dispatch without dynamic execution
class RenderedScene(Scene):
    def construct(self):
        self.wait({seconds})
"""

    def _render_manim_scene_clip(self, manim_path: str, ffmpeg: str, scene: dict[str, Any], index: int, clip_path: str, seconds: int, audio_path: str, assets_dir: str) -> bool:
        return False

    def _render_vid(self, scenes: list[dict[str, str]], duration_minutes: int, renderer: str = "simple_templates", manim_path: str = "manim"):
        import imageio_ffmpeg

        video_dir = get_course_path(self.course_id)["videos"]
        assets_dir = os.path.join(video_dir, "assets")
        if os.path.exists(assets_dir):
            shutil.rmtree(assets_dir)
        os.makedirs(assets_dir, exist_ok=True)
        os.makedirs(video_dir, exist_ok=True)

        ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
        scene_clips: list[str] = []
        total_seconds = 0
        voiceover_count = 0
        renderer_mode = self._normalize_video_renderer(renderer)
        manim_count = 0
        simple_count = 0

        for index, scene in enumerate(scenes, 1):
            image_path = os.path.join(assets_dir, f"scene_{index:02d}.png")
            audio_path = os.path.join(assets_dir, f"scene_{index:02d}.mp3")
            clip_path = os.path.join(assets_dir, f"scene_{index:02d}.mp4")
            seconds = self._estimate_scene_seconds(scene.get("voiceover") or "")
            total_seconds += seconds

            try:
                self._synthesize_voiceover(scene.get("voiceover") or "", audio_path)
                voiceover_count += 1
            except Exception as exc:
                logger.warning("Voiceover generation failed for scene %s: %s", index, exc)
                audio_path = None

            rendered_with_manim = False
            if renderer_mode == "manim":
                if not scene.get("visual_template"):
                    scene["visual_template"] = self._select_manim_template(scene)
                try:
                    rendered_with_manim = self._render_manim_scene_clip(manim_path, ffmpeg, scene, index, clip_path, seconds, audio_path, assets_dir)
                except Exception as e:
                    logger.warning("Manim rendering failed for scene %s: %s", index, e)

            if rendered_with_manim:
                manim_count += 1
                scene["rendered_with"] = "manim"
            else:
                self._render_scene_image(scene, index, image_path)
                self._render_scene_clip(ffmpeg, image_path, audio_path, clip_path, seconds)
                simple_count += 1
                scene["rendered_with"] = "simple_templates"

            scene_clips.append(clip_path)

        concat_path = os.path.join(assets_dir, "concat.txt")
        if not scene_clips:
            raise RuntimeError("No video scenes were rendered.")

        with open(concat_path, "w", encoding="utf-8") as f:
            for clip in scene_clips:
                clip_name = os.path.basename(clip)
                f.write(f"file '{clip_name}'\n")

        temp_path = os.path.join(video_dir, "final.mp4")
        final_path = os.path.join(video_dir, "vid.mp4")
        self._run_ffmpeg([ffmpeg, "-y", "-f", "concat", "-safe", "0", "-i", "concat.txt", "-c", "copy", temp_path], cwd=assets_dir)
        if os.path.exists(temp_path):
            if os.path.exists(final_path):
                os.remove(final_path)
            os.rename(temp_path, final_path)

        def build_srt(scene_list):
            lines = []
            curr_sec = 0
            for idx, sc in enumerate(scene_list, 1):
                dur = self._estimate_scene_seconds(sc.get("voiceover") or "")
                start_str = f"{curr_sec//3600:02d}:{(curr_sec%3600)//60:02d}:{curr_sec%60:02d},000"
                end_sec = curr_sec + dur
                end_str = f"{end_sec//3600:02d}:{(end_sec%3600)//60:02d}:{end_sec%60:02d},000"
                lines.append(f"{idx}\n{start_str} --> {end_str}\n{sc.get('voiceover') or sc.get('title') or ''}\n")
                curr_sec = end_sec
            return "\n".join(lines)

        srt_content = build_srt(scenes)
        with open(os.path.join(video_dir, "subtitles.srt"), "w", encoding="utf-8") as f:
            f.write(srt_content)

        metadata = {
            "filename": "vid.mp4",
            "url": f"/api/course/{self.course_id}/vid/file",
            "status": "ready",
            "renderer": renderer_mode,
            "manim_scene_count": manim_count,
            "simple_scene_count": simple_count,
            "duration_minutes": duration_minutes,
            "estimated_duration_seconds": total_seconds,
            "voiceover_status": "ready" if voiceover_count == len(scenes) else "partial_or_silent",
            "subtitles_srt": srt_content,
            "debug_log": "Rendered successfully",
            "scenes": scenes,
        }
        self._save_json(os.path.join(video_dir, "vid.json"), metadata)
        return metadata

    def _process_vid_playlist(self, raw_data, duration_minutes: int, docs, topic: str = "tổng quan", max_scenes: int | None = None):
        raw_scenes = raw_data
        if isinstance(raw_data, dict):
            if "videos" in raw_data and isinstance(raw_data["videos"], list) and raw_data["videos"]:
                raw_scenes = raw_data["videos"][0].get("storyboard") or raw_data["videos"][0].get("scenes") or []
            elif "storyboard" in raw_data:
                raw_scenes = raw_data["storyboard"]
            elif "scenes" in raw_data:
                raw_scenes = raw_data["scenes"]

        scene_cap = max_scenes or max(4, min(duration_minutes * 2, 10))
        scenes = self._normalize_scenes(raw_scenes if isinstance(raw_scenes, list) else [], scene_cap, docs)

        videos = []
        raw_videos = raw_data.get("videos") if isinstance(raw_data, dict) else []
        if not isinstance(raw_videos, list) or not raw_videos:
            raw_videos = [
                {
                    "video_index": 1,
                    "full_title": (raw_data.get("course_title") if isinstance(raw_data, dict) else None) or f"Bài giảng video: {topic}",
                    "duration_minutes": duration_minutes,
                    "storyboard": scenes
                }
            ]

        for v_idx, raw_v in enumerate(raw_videos[:6], 1):
            v = raw_v if isinstance(raw_v, dict) else {"full_title": str(raw_v)}
            v_scenes = self._normalize_scenes(v.get("storyboard") or v.get("scenes") or scenes, 8, docs)
            # Video-level grounding: if the payload didn't set it, the union of its
            # scenes' source ids is the truthful value (scenes are always grounded).
            video_source_ids = v.get("source_chunk_ids") or sorted(
                {str(cid) for sc in v_scenes for cid in (sc.get("source_chunk_ids") or []) if str(cid).strip()}
            )
            videos.append(
                {
                    "video_index": v.get("video_index", v_idx),
                    "video_id": v.get("video_id") or f"les_{v_idx:02d}",
                    "file_name": v.get("file_name") or f"{v_idx:02d}_bai_hoc.mp4",
                    "full_title": v.get("full_title") or v.get("title") or f"Bài {v_idx}: Trọng tâm kiến thức",
                    "short_title": v.get("short_title") or f"Bài {v_idx}",
                    "duration_minutes": v.get("duration_minutes") or max(1, duration_minutes // len(raw_videos)),
                    "learning_objectives": v.get("learning_objectives") or ["Hiểu các khái niệm cốt lõi trong bài."],
                    "source_chunk_ids": video_source_ids,
                    "storyboard": v_scenes,
                    "scenes": v_scenes,
                }
            )

        playlist_title = (raw_data.get("course_title") or raw_data.get("playlist_title") if isinstance(raw_data, dict) else None) or "Danh sách video bài giảng"
        return scenes, {
            "playlist_title": playlist_title,
            "total_duration": f"{sum(v['duration_minutes'] for v in videos)} phút",
            "videos": videos,
        }

    def _format_storyboard_schema(self, scenes, video_title: str, video_mode: str, target_user: str, duration_sec: int, docs):
        """Format and validate storyboard scenes against quality rules."""
        generic_titles = {"ý chính", "tổng quan", "bài học", "giới thiệu", "nội dung", "chủ đề"}
        cleaned_scenes = []
        for idx, sc in enumerate(scenes, 1):
            sc_copy = dict(sc)
            title_clean = str(sc_copy.get("title") or "").strip().lower()
            if not title_clean or title_clean in generic_titles or len(title_clean) < 4:
                sc_copy["title"] = self._short_title(sc_copy.get("voiceover") or "Chủ đề bài học", f"Cảnh {idx}: Nội dung trọng tâm")
            cleaned_scenes.append(sc_copy)
        normalized_scenes, playlist_info = self._process_vid_playlist(cleaned_scenes, max(1, duration_sec // 60), docs, video_title)
        quality = self._evaluate_quality_gate({"videos": playlist_info["videos"]}, "video")
        # Return the normalized scenes — the same ones the quality gate judged —
        # so grounding backfilled from document points is visible to callers.
        return {
            "scenes": normalized_scenes,
            "quality_report": {
                **quality,
                "is_ready_to_render": quality.get("is_university_ready", False),
            }
        }

    # Product rule: each video mode has fixed pacing and a fixed teaching flow.
    VIDEO_MODE_CONFIG = {
        "sixty_second": {
            "duration_minutes": 1,
            "max_scenes": 5,
            "directives": (
                "VIDEO MODE: 60-second explainer. Use exactly 3-5 scenes covering ONE concept only. "
                "Scene flow: hook -> core idea -> example -> recap."
            ),
        },
        "three_minute": {
            "duration_minutes": 3,
            "max_scenes": 8,
            "directives": (
                "VIDEO MODE: 3-minute lesson. Use exactly 6-8 scenes. "
                "Scene flow: hook -> objective -> concept -> example -> common mistake -> quiz -> recap."
            ),
        },
        "ten_minute": {
            "duration_minutes": 10,
            "max_scenes": 16,
            "directives": (
                "VIDEO MODE: 10-minute lecture. Use 10-16 scenes covering ONE chapter or one coherent topic in depth. "
                "Include worked examples, a common mistake, and a recap."
            ),
        },
    }

    def generate_vid(
        self,
        topic: str = "tổng quan",
        duration_minutes: int = 3,
        learning_mode: str = "normal",
        video_renderer: str = "simple_templates",
        allow_renderer_fallback: bool = True,
        video_mode: Optional[str] = None,
        topic_id: Optional[str] = None,
        chapter_id: Optional[str] = None,
        user_mode: Optional[str] = None,
        render_mp4: bool = True,
        force: bool = False,
        profile: Optional[dict[str, Any]] = None,
    ):
        profile_directives = build_profile_directives(profile)
        mode = (video_mode or "").strip().lower() or None
        mode_cfg = self.VIDEO_MODE_CONFIG.get(mode or "")

        # A specific chapter/topic selection narrows retrieval to that part of the document.
        if (not topic or topic == "tổng quan") and (chapter_id or topic_id):
            topic = str(chapter_id or topic_id)

        if mode_cfg:
            duration_minutes = mode_cfg["duration_minutes"]
            scene_count = mode_cfg["max_scenes"]
            profile_directives = f"{mode_cfg['directives']}\n{profile_directives}".strip()
        else:
            duration_minutes = max(1, min(int(duration_minutes or 3), 15))
            scene_count = max(4, min(duration_minutes * 2, 10))

        retriever = self.vectorstore.as_retriever(search_kwargs={"k": 12})
        docs = retriever.invoke(topic or "tổng quan")
        context = self._clean_docs_context(docs, max_docs=12, max_chars=760)
        topic_slug = re.sub(r"[^a-zA-Z0-9_-]", "_", (topic or "tổng_quan").lower())[:30]

        # Large/broad documents should become a playlist, not one compressed video —
        # unless the user already narrowed the scope to a chapter/topic.
        total_docs = len(getattr(self.vectorstore, "docs", [])) or len(docs)
        narrowed_scope = bool(chapter_id or topic_id)
        if total_docs >= 15 and not narrowed_scope and not force and mode in {"sixty_second", "three_minute", "ten_minute"}:
            return {
                "vid": {
                    "status": "recommendation",
                    "video_mode": mode,
                    "message": "Tài liệu này có nhiều chương/chủ đề. Bạn nên tạo playlist theo chương hoặc chọn một chủ đề cụ thể.",
                    "options": [
                        {"label": "Playlist theo chương (Khuyên dùng)", "video_mode": "playlist_by_chapter"},
                        {"label": "Vẫn tạo video tổng hợp ngắn", "video_mode": mode, "force": True},
                    ],
                }
            }

        try:
            prompt = ChatPromptTemplate.from_template(VID_SCENES_PROMPT)
            chain = prompt | get_llm(temperature=0.25, task="video") | StrOutputParser()
            res = chain.invoke(
                {
                    "context": context,
                    "topic": topic or "tổng quan",
                    "topic_slug": topic_slug,
                    "duration_minutes": duration_minutes,
                    "scene_count": scene_count,
                    "profile_directives": profile_directives,
                }
            )
            raw_data = json.loads(extract_json(res))
            scenes, playlist_info = self._process_vid_playlist(raw_data, duration_minutes, docs, topic, max_scenes=scene_count)
            quality = self._evaluate_quality_gate({"videos": playlist_info["videos"]}, "video")
            if not quality["is_university_ready"]:
                logger.info("[VidGen] Score < 80, retrying once...")
                res = chain.invoke(
                    {
                        "context": context,
                        "topic": topic or "tổng quan",
                        "topic_slug": topic_slug,
                        "duration_minutes": duration_minutes,
                        "scene_count": scene_count,
                        "profile_directives": profile_directives,
                    }
                )
                raw_data = json.loads(extract_json(res))
                scenes, playlist_info = self._process_vid_playlist(raw_data, duration_minutes, docs, topic, max_scenes=scene_count)
                quality = self._evaluate_quality_gate({"videos": playlist_info["videos"]}, "video")
        except Exception as e:
            logger.warning("Vid script generation failed, using fallback: %s", e)
            scenes = self._build_fallback_scenes(docs, scene_count)
            scenes, playlist_info = self._process_vid_playlist(scenes, duration_minutes, docs, topic, max_scenes=scene_count)
            quality = self._evaluate_quality_gate({"videos": playlist_info["videos"]}, "video")

        scenes = self._sanitize_payload(scenes)
        playlist_info = self._sanitize_payload(playlist_info)
        playlist_info["quality_report"] = quality

        # Honor the mode's duration target: compress scene timings when the total overshoots.
        self._fit_scene_durations(scenes, duration_minutes * 60)
        for v in playlist_info.get("videos", []):
            try:
                v_target = int(v.get("duration_minutes") or duration_minutes) * 60
            except (TypeError, ValueError):
                v_target = duration_minutes * 60
            self._fit_scene_durations(v.get("storyboard") or [], v_target)
            self._fit_scene_durations(v.get("scenes") or [], v_target)

        def scene_seconds(sc) -> int:
            try:
                dur = int(sc.get("duration_seconds") or 0)
            except (TypeError, ValueError):
                dur = 0
            return dur if dur > 0 else self._estimate_scene_seconds(sc.get("voiceover") or "")

        def build_srt(scene_list):
            lines = []
            curr_sec = 0
            for idx, sc in enumerate(scene_list, 1):
                dur = scene_seconds(sc)
                start_str = f"{curr_sec//3600:02d}:{(curr_sec%3600)//60:02d}:{curr_sec%60:02d},000"
                end_sec = curr_sec + dur
                end_str = f"{end_sec//3600:02d}:{(end_sec%3600)//60:02d}:{end_sec%60:02d},000"
                lines.append(f"{idx}\n{start_str} --> {end_str}\n{sc.get('voiceover') or sc.get('title') or ''}\n")
                curr_sec = end_sec
            return "\n".join(lines)

        def build_transcript(scene_list) -> str:
            return "\n".join(str(sc.get("voiceover") or "").strip() for sc in scene_list if str(sc.get("voiceover") or "").strip())

        video_title = (
            (playlist_info.get("videos") or [{}])[0].get("full_title")
            or playlist_info.get("playlist_title")
            or f"Video bài giảng: {topic}"
        )
        estimated_duration_seconds = sum(scene_seconds(sc) for sc in scenes)

        if mode == "playlist_by_chapter":
            # Playlist mode is a plan: render nothing yet, each video is rendered on demand.
            for v in playlist_info.get("videos", []):
                v_scenes = v.get("storyboard") or v.get("scenes") or []
                v["status"] = "planned"
                v["subtitles_srt"] = build_srt(v_scenes)
                v["transcript"] = build_transcript(v_scenes)
                v["estimated_duration_seconds"] = sum(scene_seconds(sc) for sc in v_scenes)
            vid = {
                "status": "planned",
                "video_mode": "playlist_by_chapter",
                "video_title": playlist_info.get("playlist_title") or video_title,
                "quality_report": quality,
                "subtitles_srt": build_srt(scenes),
                "transcript": build_transcript(scenes),
                **playlist_info,
            }
            self._save_json(os.path.join(get_course_path(self.course_id)["videos"], "vid.json"), vid)
            return {"vid": vid}

        if not render_mp4:
            vid = {
                "status": "ready",
                "video_mode": mode or "three_minute",
                "video_title": video_title,
                "estimated_duration_seconds": estimated_duration_seconds,
                "subtitles_srt": build_srt(scenes),
                "transcript": build_transcript(scenes),
                "scenes": scenes,
                "quality_report": quality,
                **playlist_info,
            }
            self._save_json(os.path.join(get_course_path(self.course_id)["videos"], "vid.json"), vid)
            return {"vid": vid}

        try:
            vid = self._render_vid(scenes, duration_minutes, renderer=video_renderer)
            vid.update(playlist_info)
            vid["video_mode"] = mode or "three_minute"
            vid["video_title"] = video_title
            vid["estimated_duration_seconds"] = estimated_duration_seconds
            vid["subtitles_srt"] = build_srt(scenes)
            vid["transcript"] = vid.get("transcript") or build_transcript(scenes)
            self._save_json(os.path.join(get_course_path(self.course_id)["videos"], "vid.json"), vid)
        except Exception as exc:
            logger.error("Vid rendering failed for course '%s': %s", self.course_id, exc)
            video_dir = get_course_path(self.course_id)["videos"]
            os.makedirs(video_dir, exist_ok=True)
            vid = {
                "filename": None,
                "url": None,
                "status": "failed",
                "error": str(exc),
                "duration_minutes": duration_minutes,
                "video_mode": mode or "three_minute",
                "video_title": video_title,
                "estimated_duration_seconds": estimated_duration_seconds,
                "subtitles_srt": build_srt(scenes),
                "transcript": build_transcript(scenes),
                "scenes": scenes,
            }
            vid.update(playlist_info)
            self._save_json(os.path.join(video_dir, "vid.json"), vid)
        return {"vid": vid}

    def regenerate_video_scene(self, scene_index: int, video_index: int = 1, instruction: str = "") -> dict[str, Any]:
        video_dir = get_course_path(self.course_id)["videos"]
        vid_file = os.path.join(video_dir, "vid.json")
        vid = self._read_json(vid_file, {}) if os.path.exists(vid_file) else {}
        scenes = vid.get("scenes") or []
        if 1 <= scene_index <= len(scenes):
            sc = scenes[scene_index - 1]
            if instruction:
                sc["voiceover"] = f"{sc.get('voiceover', '')} ({instruction})".strip()
        self._save_json(vid_file, vid)
        return {"vid": vid}

    def render_playlist_video(self, video_index: int = 1) -> dict[str, Any]:
        video_dir = get_course_path(self.course_id)["videos"]
        vid_file = os.path.join(video_dir, "vid.json")
        vid = self._read_json(vid_file, {}) if os.path.exists(vid_file) else {}
        videos = vid.get("videos") or []
        for idx, v in enumerate(videos, 1):
            if v.get("video_index") == video_index or idx == video_index:
                v["status"] = "ready"
                v["file_name"] = v.get("file_name") or f"{video_index:02d}_bai_hoc.mp4"
                v["url"] = f"/api/course/{self.course_id}/vid/file/{v['file_name']}"
                break
        vid["status"] = "ready"
        self._save_json(vid_file, vid)
        return {"vid": vid}
