"""Generation Service responsible for RAG retrieval, LLM generation, scoring, and artifact storage."""

import hashlib
import json
import logging
import os
import re
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple
from app.core.config import settings
from app.models.course import Course
from app.schemas.generation import (
    GroundingData,
    QualityScoresData,
    ReadinessData,
    RegenLimitsData,
    StudyPackData,
    StudyPackResponse,
    StudyPackStats,
)
from app.schemas.generator_output import (
    BookChapter,
    BookOutput,
    QuizOutput,
    SlidesOutput,
    VidOutput,
    validate_and_score_output,
)
from app.services.llm import LLMService
from app.services.pdf_book import build_book_pdf
from app.services.text_format import clean_text
from app.services.vector_store import Document, VectorStore
from app.services.versioning import (
    AtomicArtifactDirectory,
    GenerationInFlightError,
    VERSION_CAPS,
    VersionCapReachedError,
    artifact_directory_path,
    migrate_legacy_artifact_metadata,
    remove_artifact_version,
    version_label,
    version_slug,
)

logger = logging.getLogger(__name__)

# Extra regenerations allowed per artifact type, per course, after the first generation.
# The first generation of an artifact is always free — this only caps deliberate "I'm
# not happy with the quality, try again" re-triggers, to protect the scarce Gemini quota.
MAX_REGENERATIONS = 3


_LEADING_ID_TOKEN_RE = re.compile(r"^([A-Za-z0-9-]+)[_\s]+")


def _strip_leading_id_token(text: str) -> str:
    """Strip a leading filename-style identifier code (e.g. "NLC416-14jh005357-58048_Title"
    -> "Title") from a topic-fallback string, so a generated title/topic never leaks an
    internal document ID. Only strips when the leading token has >=4 digits, so a real title
    that happens to start with a short number (e.g. "3D Printing Basics") survives untouched."""
    m = _LEADING_ID_TOKEN_RE.match(text)
    if not m:
        return text
    token = m.group(1)
    if sum(c.isdigit() for c in token) >= 4:
        return text[m.end():].strip()
    return text


def _clean_slides_output(deck: SlidesOutput) -> SlidesOutput:
    """Normalize LLM-authored prose (strip LaTeX/Markdown) in-place across a slide deck.
    Slides are always rasterized to images for the reader (no live-text consumer), so this
    stays the only artifact type cleaned before persistence; Book/Quiz keep their raw
    LLM Markdown/LaTeX in storage and clean it only at PDF-build time (see pdf_utils.prepare_pdf_text)
    so the frontend's real Markdown+KaTeX renderer gets full fidelity."""
    deck.title = clean_text(deck.title)
    for sl in deck.slides:
        sl.title = clean_text(sl.title)
        sl.bullet_points = [clean_text(b) for b in sl.bullet_points]
    return deck


def _clean_vid_output(vid: VidOutput) -> VidOutput:
    """Normalize LLM-authored prose (strip LaTeX/Markdown) in-place across a video script.
    Frames are rasterized (like Slides), so this stays the only cleanup pass before render."""
    vid.title = clean_text(vid.title)
    for sc in vid.scenes:
        sc.title = clean_text(sc.title)
        sc.on_screen_text = clean_text(sc.on_screen_text or "")
        sc.key_points = [clean_text(kp) for kp in sc.key_points]
        if sc.diagram:
            sc.diagram.title = clean_text(sc.diagram.title or "") or None
            for item in sc.diagram.items:
                item.label = clean_text(item.label)
                item.detail = clean_text(item.detail or "") or None
        sc.narration = clean_text(sc.narration)
    return vid


class Generator:
    """Orchestrates RAG retrieval, AI generation, validation, and file storage."""

    def __init__(self, vector_store: VectorStore, llm: LLMService, feature_llms: Optional[Dict[str, LLMService]] = None):
        self.vector_store = vector_store
        self.llm = llm
        # Per-feature LLM instances (book/slides/quiz/vid), each optionally configured with
        # its own Gemini API key. Falls back to `self.llm` for any feature not provided.
        self.feature_llms = feature_llms or {}
        self._generation_versions: Dict[tuple[str, str], str] = {}

    def _llm_for(self, feature: str) -> LLMService:
        return self.feature_llms.get(feature, self.llm)

    def _get_db(self, db_session_factory=None):
        if db_session_factory is None:
            from app.services.database import SessionLocal
            return SessionLocal()
        return db_session_factory()

    def _get_embedding_provider(self, course_id: str, db_session_factory=None) -> str:
        """Which Chroma collection this course's chunks actually live in — "gemini" (default)
        or "openrouter" if document_processor fell back to it on Gemini quota exhaustion.
        Retrieval must route to the same collection a course was embedded into, since vectors
        from different embedding models aren't comparable (see VectorStore._collection_for)."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            return (course.embedding_provider if course and course.embedding_provider else "gemini")
        finally:
            db.close()

    def _retrieve_context(
        self, course_id: str, query: str = "", k: int = 20, db_session_factory=None
    ) -> Tuple[str, List[str]]:
        """Retrieve relevant chunks from Chroma vector store, drop near-duplicate chunks
        (retrieval overlap / re-retrieval across Book chapters), and order the survivors
        by (source_file, page) so the LLM sees document-coherent context instead of chunks
        shuffled by raw similarity rank."""
        search_query = query or "tổng quan kiến thức khóa học các chương quan trọng"
        provider = self._get_embedding_provider(course_id, db_session_factory)
        chunks = self.vector_store.search(query=search_query, course_id=course_id, k=k, provider=provider)
        if not chunks:
            logger.warning(f"No vector chunks found for course {course_id}. RAG context will be empty.")
            return "", []

        seen_hashes = set()
        deduped = []
        for doc in chunks:
            normalized = " ".join(doc.content.strip().lower().split())[:150]
            digest = hashlib.sha1(normalized.encode("utf-8")).hexdigest()
            if digest in seen_hashes:
                continue
            seen_hashes.add(digest)
            deduped.append(doc)

        def _sort_key(doc: Document):
            page = doc.metadata.get("page", 0)
            try:
                page = int(page)
            except (TypeError, ValueError):
                page = 0
            return (str(doc.metadata.get("source_file", "")), page)

        deduped.sort(key=_sort_key)

        context_lines = []
        valid_chunk_ids = []
        for i, doc in enumerate(deduped):
            cid = doc.metadata.get("chunk_id") or f"chunk_{i+1}"
            valid_chunk_ids.append(cid)
            file_name = doc.metadata.get("source_file", "unknown")
            page_num = doc.metadata.get("page", 1)
            context_lines.append(
                f"[Chunk ID: {cid}] (Tài liệu: {file_name}, Trang: {page_num}):\n{doc.content}"
            )

        return "\n\n".join(context_lines), valid_chunk_ids

    _NO_CONTEXT_MSG = (
        "Không tìm thấy nội dung nào từ tài liệu để tạo học liệu. Tài liệu có thể là bản "
        "scan/ảnh chưa trích xuất được chữ, hoặc chưa lập chỉ mục thành công. Hãy thử xoá "
        "và tải lại tài liệu (ưu tiên PDF có lớp văn bản thật, không phải ảnh chụp)."
    )
    _PROCESSING_MSG = (
        "Tài liệu vẫn đang được xử lý (trích xuất và lập chỉ mục nội dung). Việc này có thể "
        "mất khoảng nửa phút với tài liệu dài. Vui lòng đợi giây lát rồi thử lại."
    )

    def _require_course_not_processing(self, course_id: str, db_session_factory=None) -> None:
        """Guard: refuse to run a generator while the course's document ingestion is still
        running. Without this, a generate call fired right after upload — before chunking/
        embedding finishes writing chunk_count — hits the exact same "no chunks found" path
        as a genuinely broken document, and `_require_context`'s scan/OCR-focused message is
        actively misleading here since the document is fine, just not indexed yet. Ingestion
        legitimately takes 10-30+s for large real documents (retries, OpenRouter fallback on
        Gemini quota exhaustion — see document_processor.process_course), long enough for a
        user clicking into a tab right after upload to reliably hit this race."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if course and course.status == "processing":
                raise ValueError(self._PROCESSING_MSG)
        finally:
            db.close()

    def _require_context(self, context: str) -> None:
        """Guard: refuse to run a generator on empty RAG context. Without this, the LLM
        dutifully writes a 'no context was provided' apology that gets saved and marked
        ready — the Study Guide shows an abstention essay, Quiz yields 0 questions (the
        frontend then sticks at 100% because its data never passes isReady), and Video
        narrates generic filler. Failing loud turns all of those into one clear error."""
        if not context or not context.strip():
            raise ValueError(self._NO_CONTEXT_MSG)

    def _get_artifact_dir(self, course_id: str) -> str:
        """Get or create local filesystem artifact storage directory."""
        dir_path = os.path.join(settings.UPLOAD_DIR, course_id, "artifacts")
        os.makedirs(dir_path, exist_ok=True)
        return dir_path

    def _start_version_write(self, course_id: str, artifact: str, version_id: Optional[str]):
        if not version_id:
            return None, None
        self._generation_versions[(course_id, artifact)] = version_id
        transaction = AtomicArtifactDirectory(
            artifact_directory_path(settings.UPLOAD_DIR, course_id, artifact, version_id)
        )
        return transaction, transaction.prepare()

    def _finish_version_write(self, transaction, success: bool) -> None:
        if transaction:
            transaction.commit() if success else transaction.abort()

    def _save_artifact_json(self, course_id: str, filename: str, data: Any, artifact_dir: Optional[str] = None) -> str:
        """Save generated Pydantic model or dict as JSON file."""
        dir_path = artifact_dir or self._get_artifact_dir(course_id)
        file_path = os.path.join(dir_path, filename)
        try:
            if hasattr(data, "model_dump"):
                content_dict = data.model_dump()
            elif hasattr(data, "dict"):
                content_dict = data.dict()
            else:
                content_dict = data
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(content_dict, f, ensure_ascii=False, indent=2)
            logger.info(f"Saved artifact JSON to {file_path}")
            return file_path
        except Exception as e:
            logger.error(f"Error saving artifact JSON {filename}: {e}")
            raise

    def _load_artifact_json(self, course_id: str, filename: str, artifact_dir: Optional[str] = None) -> Optional[Dict[str, Any]]:
        """Load artifact JSON from disk if exists."""
        file_path = os.path.join(artifact_dir or self._get_artifact_dir(course_id), filename)
        if not os.path.exists(file_path):
            return None
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Error loading artifact JSON {filename}: {e}")
            return None

    def _generate_pdf_book(self, course_id: str, book_data: BookOutput, artifact_dir: Optional[str] = None) -> str:
        """Generate the Study Guide PDF (cover, preface, page-numbered TOC, chapters).

        Raises on failure — callers must treat that as a hard generation error, not write a
        placeholder file in its place.
        """
        file_path = os.path.join(artifact_dir or self._get_artifact_dir(course_id), "book.pdf")
        build_book_pdf(file_path, book_data)
        return file_path

    def _generate_pdf_slides(self, course_id: str, slides_data: SlidesOutput, artifact_dir: Optional[str] = None) -> str:
        """Generate a 16:9 Widescreen PDF presentation using ReportLab."""
        file_path = os.path.join(artifact_dir or self._get_artifact_dir(course_id), "slide.pdf")
        try:
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib import colors
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, PageBreak, Table, TableStyle
            from app.services.pdf_utils import prepare_pdf_text, register_vietnamese_fonts

            font_name, font_bold = register_vietnamese_fonts()

            # Widescreen 16:9 is 960 width by 540 height
            doc = SimpleDocTemplate(
                file_path,
                pagesize=(960, 540),
                leftMargin=50,
                rightMargin=50,
                topMargin=50,
                bottomMargin=50
            )

            styles = getSampleStyleSheet()

            # Styles for Slide contents
            slide_title_style = ParagraphStyle(
                "SlideTitle",
                parent=styles["Title"],
                fontName=font_bold,
                fontSize=28,
                leading=34,
                textColor=colors.HexColor("#06b6d4"), # Cyan-500
                alignment=0,
                spaceAfter=20
            )

            slide_body_style = ParagraphStyle(
                "SlideBody",
                parent=styles["BodyText"],
                fontName=font_name,
                fontSize=16,
                leading=22,
                textColor=colors.HexColor("#f8fafc"), # slate-50
                spaceAfter=12
            )

            slide_quote_style = ParagraphStyle(
                "SlideQuote",
                parent=styles["Normal"],
                fontName=font_name,
                fontSize=22,
                leading=30,
                textColor=colors.HexColor("#38bdf8"), # light-blue-400
                alignment=1,
                spaceAfter=20
            )

            story = []

            # Page Templates: background & footer drawing
            def draw_title_slide_bg(canvas, doc):
                canvas.saveState()
                canvas.setFillColor(colors.HexColor("#020617"))
                canvas.rect(0, 0, 960, 540, fill=True, stroke=False)
                
                canvas.setFillColor(colors.HexColor("#0f172a"))
                canvas.rect(0, 0, 960, 80, fill=True, stroke=False)
                
                canvas.setFillColor(colors.HexColor("#06b6d4"))
                canvas.rect(50, 150, 8, 260, fill=True, stroke=False)
                canvas.restoreState()

            def draw_content_slide_bg(canvas, doc):
                canvas.saveState()
                canvas.setFillColor(colors.HexColor("#0f172a"))
                canvas.rect(0, 0, 960, 540, fill=True, stroke=False)

                canvas.setFillColor(colors.HexColor("#1e293b"))
                canvas.rect(0, 480, 960, 60, fill=True, stroke=False)

                canvas.setStrokeColor(colors.HexColor("#06b6d4"))
                canvas.setLineWidth(2)
                canvas.line(0, 480, 960, 480)

                canvas.setFillColor(colors.HexColor("#64748b"))
                canvas.setFont(font_name, 10)
                canvas.drawRightString(910, 20, f"Slide {canvas._pageNumber}")
                canvas.restoreState()

            # Slide 1: Cover Page
            story.append(Spacer(1, 100))
            title_p = Paragraph(f"<font color='#06b6d4'>{prepare_pdf_text(slides_data.title)}</font>", ParagraphStyle("CoverTitle", parent=slide_title_style, fontSize=36, leading=44, leftIndent=30))
            story.append(title_p)
            story.append(PageBreak())

            # Slide content pages
            for idx, item in enumerate(slides_data.slides):
                story.append(Spacer(1, 10))
                story.append(Paragraph(prepare_pdf_text(item.title), slide_title_style))
                story.append(Spacer(1, 10))

                layout_type = getattr(item, "layout_type", "default") or "default"
                if layout_type == "two_column":
                    mid = (len(item.bullet_points) + 1) // 2
                    left_bps = [f"• {prepare_pdf_text(bp)}" for bp in item.bullet_points[:mid]]
                    right_bps = [f"• {prepare_pdf_text(bp)}" for bp in item.bullet_points[mid:]]

                    left_text = "<br/><br/>".join(left_bps)
                    right_text = "<br/><br/>".join(right_bps)

                    left_p = Paragraph(left_text, slide_body_style)
                    right_p = Paragraph(right_text, slide_body_style)

                    t = Table([[left_p, right_p]], colWidths=[420, 420])
                    t.setStyle(TableStyle([
                        ('VALIGN', (0,0), (-1,-1), 'TOP'),
                        ('LEFTPADDING', (0,0), (-1,-1), 0),
                        ('RIGHTPADDING', (0,0), (-1,-1), 10),
                    ]))
                    story.append(t)
                elif layout_type == "quote":
                    quote_text = "<br/><br/>".join(prepare_pdf_text(bp) for bp in item.bullet_points)
                    story.append(Spacer(1, 40))
                    story.append(Paragraph(f"<i>“{quote_text}”</i>", slide_quote_style))
                else:
                    bullets_html = []
                    for bp in item.bullet_points:
                        bullets_html.append(f"• {prepare_pdf_text(bp)}")
                    content_text = "<br/><br/>".join(bullets_html)
                    story.append(Paragraph(content_text, slide_body_style))

                if idx < len(slides_data.slides) - 1:
                    story.append(PageBreak())

            doc.build(
                story,
                onFirstPage=draw_title_slide_bg,
                onLaterPages=draw_content_slide_bg
            )
            logger.info(f"Generated 16:9 PDF Slides at {file_path}")
            return file_path
        except Exception as e:
            logger.error(f"Error generating PDF slides: {e}", exc_info=True)
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(f"PDF Slides Placeholder for {slides_data.title}")
            return file_path

    def _convert_pdf_to_images(self, pdf_path: str, course_id: str, artifact_dir: Optional[str] = None) -> list[str]:
        """Convert a PDF file into PNG slide images using PyMuPDF (fitz)."""
        import fitz
        artifact_dir = artifact_dir or self._get_artifact_dir(course_id)
        image_paths = []
        try:
            doc = fitz.open(pdf_path)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img_name = f"slide_{i+1}.png"
                img_path = os.path.join(artifact_dir, img_name)
                pix.save(img_path)
                image_paths.append(img_path)
            logger.info(f"Converted {len(image_paths)} pages to PNG slide images for course {course_id}")
            return image_paths
        except Exception as e:
            logger.error(f"Failed to convert PDF to images: {e}", exc_info=True)
            return []

    def _generate_pptx_slides(self, course_id: str, slides_data: SlidesOutput, artifact_dir: Optional[str] = None) -> str:
        """Generate PowerPoint presentation by inserting ReportLab slide PNGs full screen."""
        artifact_dir = artifact_dir or self._get_artifact_dir(course_id)
        file_path = os.path.join(artifact_dir, "slide.pptx")
        try:
            from pptx import Presentation
            from pptx.util import Inches

            pdf_path = self._generate_pdf_slides(course_id, slides_data, artifact_dir)
            image_paths = self._convert_pdf_to_images(pdf_path, course_id, artifact_dir)

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            
            for img_path in image_paths:
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            prs.save(file_path)
            logger.info(f"Generated Slide PPTX at {file_path}")
            return file_path
        except Exception as e:
            logger.error(f"Error generating Slide PPTX: {e}", exc_info=True)
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(f"PPTX Placeholder for {slides_data.title}")
            return file_path

    def _generate_pdf_quiz_key(self, course_id: str, quiz_data: QuizOutput, artifact_dir: Optional[str] = None) -> str:
        """Generate Quiz PDF in two sections: Student Quiz Sheet and Answer Key & Explanations."""
        file_path = os.path.join(artifact_dir or self._get_artifact_dir(course_id), "quiz-key.pdf")
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib import colors
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, PageBreak
            from app.services.pdf_utils import prepare_pdf_text, register_vietnamese_fonts

            font_name, font_bold = register_vietnamese_fonts()

            doc = SimpleDocTemplate(file_path, pagesize=A4, rightMargin=50, leftMargin=50, topMargin=50, bottomMargin=50)
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle("QTitle", parent=styles["Title"], fontName=font_bold, fontSize=20, leading=24, textColor=colors.HexColor("#1e3a8a"), spaceAfter=15)
            part_style = ParagraphStyle("QPart", parent=styles["Heading1"], fontName=font_bold, fontSize=15, leading=18, textColor=colors.HexColor("#4f46e5"), spaceBefore=10, spaceAfter=15)
            q_style = ParagraphStyle("QQ", parent=styles["Heading2"], fontName=font_bold, fontSize=12, leading=16, textColor=colors.HexColor("#1e293b"), spaceBefore=10, spaceAfter=6)
            body_style = ParagraphStyle("QBody", parent=styles["BodyText"], fontName=font_name, fontSize=11, leading=15, textColor=colors.HexColor("#0f172a"), spaceAfter=5)

            story = [Paragraph(f"BỘ ĐỀ KIỂM TRA & ĐÁNH GIÁ: {prepare_pdf_text(quiz_data.title)}", title_style), Spacer(1, 10)]

            # --- Part 1: Student Quiz Sheet ---
            story.append(Paragraph("PHẦN 1: ĐỀ THI TRẮC NGHIỆM (STUDENT QUIZ SHEET)", part_style))
            for q in quiz_data.questions:
                story.append(Paragraph(f"<b>Câu {q.question_number}:</b> {prepare_pdf_text(q.question_text)}", q_style))
                for opt in q.options:
                    story.append(Paragraph(f"<b>{opt.key}.</b> {prepare_pdf_text(opt.text)}", body_style))
                story.append(Spacer(1, 10))
            story.append(PageBreak())

            # --- Part 2: Answer Key & Explanations ---
            # Same Dễ/Vừa/Khó label the web quiz badge shows — no internal "Bloom" jargon.
            difficulty_vn = {"easy": "Dễ", "medium": "Vừa", "hard": "Khó"}
            story.append(Paragraph("PHẦN 2: ĐÁP ÁN & GIẢI THÍCH CHI TIẾT (ANSWER KEY & EXPLANATIONS)", part_style))
            for q in quiz_data.questions:
                diff_str = getattr(q, "difficulty", "Medium") or "Medium"
                diff_label = difficulty_vn.get(diff_str.strip().lower(), diff_str)
                story.append(Paragraph(f"<b>Câu {q.question_number} ({diff_label}):</b> {prepare_pdf_text(q.question_text)}", q_style))
                story.append(Paragraph(f"<b>Đáp án đúng:</b> <font color='#16a34a'><b>{q.correct_answer}</b></font>", body_style))
                if q.explanation:
                    story.append(Paragraph(f"<b>Giải thích chi tiết:</b> {prepare_pdf_text(q.explanation)}", body_style))
                story.append(Spacer(1, 10))

            doc.build(story)
            logger.info(f"Generated Quiz Key PDF at {file_path}")
            return file_path
        except Exception as e:
            logger.error(f"Error generating Quiz Key PDF: {e}")
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(f"Quiz Key Placeholder for {quiz_data.title}")
            return file_path

    def _generate_video_mp4(
        self,
        course_id: str,
        vid_data: VidOutput,
        fmt: str,
        voice: str,
        progress_cb=None,
        artifact_dir: Optional[str] = None,
        scene_visual_map: Optional[Dict[int, Dict[str, Any]]] = None,
    ) -> str:
        """Render the narrated MP4 (TTS + still frames + ffmpeg concat) plus transcript.txt /
        vid.srt. Raises on failure — callers must treat that as a hard generation error, not
        write a placeholder file in its place (matches the strict invariant used by Book's PDF)."""
        from app.services.video_render import assemble_video

        artifact_dir = artifact_dir or self._get_artifact_dir(course_id)
        return assemble_video(
            vid_data,
            fmt,
            voice,
            artifact_dir,
            progress_cb=progress_cb,
            scene_visual_map=scene_visual_map,
        )

    def _resolve_course_pdf_path(self, course_id: str, source_file: str) -> Optional[str]:
        """Resolve a chunk's original filename to its timestamp-prefixed uploaded PDF."""
        safe_name = os.path.basename(source_file or "")
        if not safe_name.lower().endswith(".pdf"):
            return None
        course_dir = os.path.join(settings.UPLOAD_DIR, course_id)
        direct_path = os.path.join(course_dir, safe_name)
        if os.path.isfile(direct_path):
            return direct_path
        try:
            for name in os.listdir(course_dir):
                prefix, separator, original_name = name.partition("_")
                if separator and prefix.isdigit() and original_name == safe_name:
                    path = os.path.join(course_dir, name)
                    if os.path.isfile(path):
                        return path
        except OSError:
            return None
        return None

    def _build_scene_visual_map(self, course_id: str, vid_data: VidOutput) -> Dict[int, Dict[str, Any]]:
        """Select grounded PDF-page visuals for every other eligible middle video scene."""
        chunk_ids = [chunk_id for scene in vid_data.scenes for chunk_id in scene.source_chunk_ids]
        if not chunk_ids:
            return {}
        chunks = self.vector_store.get_course_chunks(course_id, list(dict.fromkeys(chunk_ids)))
        chunks_by_id = {str(chunk.metadata.get("chunk_id", "")): chunk for chunk in chunks}
        candidates: List[Dict[str, Any]] = []
        total_scenes = len(vid_data.scenes)
        for index, scene in enumerate(vid_data.scenes):
            if index == 0 or index == total_scenes - 1 or scene.diagram:
                continue
            for chunk_id in scene.source_chunk_ids:
                chunk = chunks_by_id.get(chunk_id)
                if not chunk:
                    continue
                pdf_path = self._resolve_course_pdf_path(course_id, str(chunk.metadata.get("source_file", "")))
                if not pdf_path:
                    continue
                try:
                    page = max(1, int(chunk.metadata.get("page", 1)))
                except (TypeError, ValueError):
                    page = 1
                candidates.append({"scene_number": scene.scene_number, "pdf_path": pdf_path, "page": page})
                break

        visual_map: Dict[int, Dict[str, Any]] = {}
        for visual_index, candidate in enumerate(candidates[::2]):
            visual_map[candidate["scene_number"]] = {
                "pdf_path": candidate["pdf_path"],
                "page": candidate["page"],
                "side": "left" if visual_index % 2 == 0 else "right",
            }
        return visual_map

    def _update_course_metadata(self, course_id: str, artifact_type: str, score: int, db_session_factory=None):
        """Update course metadata_json and readiness flags in database."""
        version_id = self._generation_versions.get((course_id, artifact_type))
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if not course:
                return

            meta = course.metadata_json or "{}"
            if isinstance(meta, str):
                try:
                    meta = json.loads(meta)
                except Exception:
                    meta = {}
            study_pack = meta.get("study_pack", {})
            if version_id:
                artifacts = dict(study_pack.get("artifacts", {}))
                entry = dict(artifacts.get(artifact_type, {}))
                versions = dict(entry.get("versions", {}))
                if version_id in versions:
                    version = dict(versions[version_id])
                    version["quality_score"] = score
                    versions[version_id] = version
                    entry["versions"] = versions
                    artifacts[artifact_type] = entry
                    study_pack["artifacts"] = artifacts
            readiness = study_pack.get("readiness", {})
            quality_scores = study_pack.get("quality_scores", {})
            grounding = study_pack.get("grounding", {"num_chunks": course.chunk_count, "quality_score": 0, "warnings": []})

            if artifact_type == "book":
                readiness["study_guide_pdf"] = True
                quality_scores["study_guide_pdf"] = score
            elif artifact_type == "slides":
                readiness["slides"] = True
                quality_scores["slides"] = score
            elif artifact_type == "quiz":
                readiness["quiz"] = True
                quality_scores["quiz"] = score
            elif artifact_type == "vid":
                readiness["vid"] = True
                quality_scores["vid"] = score


            # Update overall average score
            scores_list = [v for v in quality_scores.values() if v > 0]
            if scores_list:
                grounding["quality_score"] = sum(scores_list) // len(scores_list)
                course.quality_score = grounding["quality_score"]

            study_pack["readiness"] = readiness
            study_pack["quality_scores"] = quality_scores
            study_pack["grounding"] = grounding
            meta["study_pack"] = study_pack
            course.metadata_json = json.dumps(meta, ensure_ascii=False)
            db.commit()
            logger.info(f"Updated course {course_id} metadata for {artifact_type} with score {score}")
        except Exception as e:
            db.rollback()
            logger.error(f"Error updating course metadata: {e}")
        finally:
            db.close()

    @staticmethod
    def _metadata_dict(course: Course) -> Dict[str, Any]:
        raw = course.metadata_json or "{}"
        if isinstance(raw, str):
            try:
                raw = json.loads(raw)
            except Exception:
                raw = {}
        return raw if isinstance(raw, dict) else {}

    def prepare_artifact_version(
        self, course_id: str, artifact: str, options: Dict[str, Any], topic: Optional[str] = None,
        user_prompt: str = "", replace_version_id: Optional[str] = None, reserve: bool = True, db_session_factory=None,
    ) -> str:
        """Reserve a version slot and enforce per-artifact concurrency/caps."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if not course:
                raise ValueError("Course not found")
            meta, _ = migrate_legacy_artifact_metadata(self._metadata_dict(course))
            study_pack = dict(meta.get("study_pack", {}))
            artifacts = dict(study_pack.get("artifacts", {}))
            entry = dict(artifacts.get(artifact, {}))
            versions = dict(entry.get("versions", {}))
            now = datetime.utcnow().isoformat()
            stale_minutes = {"book": 10, "vid": 20, "slides": 8, "quiz": 8}[artifact]
            for value in versions.values():
                if not isinstance(value, dict) or value.get("status") != "processing":
                    continue
                try:
                    age = datetime.utcnow() - datetime.fromisoformat(value.get("updated_at", ""))
                except (TypeError, ValueError):
                    age = None
                if age is None or age.total_seconds() <= stale_minutes * 60:
                    raise GenerationInFlightError()
                value.update({"status": "error", "error": "Tác vụ tạo đã hết thời gian chờ.", "finished_at": now, "updated_at": now})

            version_id = version_slug(artifact, options)
            is_new = version_id not in versions
            if is_new and len(versions) >= VERSION_CAPS[artifact] and not replace_version_id:
                raise VersionCapReachedError(self._version_summaries(versions))
            if replace_version_id and replace_version_id not in versions:
                raise ValueError("Version selected for replacement does not exist")
            current = dict(versions.get(version_id, {}))
            current.update({
                "options": dict(options), "label": version_label(artifact, options), "topic": topic,
                "user_prompt": user_prompt, "path": version_id, "status": "processing", "error": None,
                "progress": 0, "created_at": current.get("created_at", now), "started_at": now, "updated_at": now,
            })
            if is_new and replace_version_id:
                current["replace_version_id"] = replace_version_id
            versions[version_id] = current
            artifacts[artifact] = {"active": entry.get("active"), "versions": versions}
            study_pack["artifacts"] = artifacts
            meta["study_pack"] = study_pack
            if reserve:
                course.metadata_json = json.dumps(meta, ensure_ascii=False)
                db.commit()
            return version_id
        except Exception:
            db.rollback()
            raise
        finally:
            db.close()

    @staticmethod
    def _version_summaries(versions: Dict[str, Any]) -> list[Dict[str, Any]]:
        return [
            {"version_id": key, "label": value.get("label", key), "options": value.get("options", {}),
             "status": value.get("status", "empty"), "created_at": value.get("created_at")}
            for key, value in versions.items() if isinstance(value, dict)
        ]

    def artifact_versions(self, course_id: str, artifact: str, db_session_factory=None) -> tuple[Optional[str], list[Dict[str, Any]]]:
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if not course:
                return None, []
            meta, _ = migrate_legacy_artifact_metadata(self._metadata_dict(course))
            entry = meta.get("study_pack", {}).get("artifacts", {}).get(artifact, {})
            return entry.get("active"), self._version_summaries(entry.get("versions", {}))
        finally:
            db.close()

    def _set_artifact_status(
        self,
        course_id: str,
        artifact: str,
        status: str,
        error: Optional[str] = None,
        progress: Optional[int] = None,
        version_id: Optional[str] = None,
        db_session_factory=None,
    ):
        """Persist per-artifact generation status (processing/ready/error) into Course.metadata_json."""
        version_id = version_id or self._generation_versions.get((course_id, artifact))
        replacement_to_remove: Optional[str] = None
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if not course:
                return

            meta, _ = migrate_legacy_artifact_metadata(self._metadata_dict(course))
            study_pack = dict(meta.get("study_pack", {}))
            artifacts = dict(study_pack.get("artifacts", {}))
            entry = dict(artifacts.get(artifact, {}))
            if not version_id and entry.get("active") and isinstance(entry.get("versions"), dict):
                version_id = entry["active"]
            versions = dict(entry.get("versions", {})) if version_id else {}
            current = dict(versions.get(version_id, {})) if version_id else entry

            now = datetime.utcnow().isoformat()
            current["status"] = status
            current["error"] = error
            if progress is not None:
                current["progress"] = progress
            if status == "processing" and "started_at" not in current:
                current["started_at"] = now
            if status in ("ready", "error"):
                current["finished_at"] = now
            current["updated_at"] = now
            if version_id:
                versions[version_id] = current
                if status == "ready":
                    victim = current.pop("replace_version_id", None)
                    if victim and victim != version_id:
                        versions.pop(victim, None)
                        replacement_to_remove = victim
                    entry["active"] = version_id
                entry["versions"] = versions
            else:
                entry = current

            artifacts[artifact] = entry
            study_pack["artifacts"] = artifacts
            meta["study_pack"] = study_pack
            course.metadata_json = json.dumps(meta, ensure_ascii=False)
            db.commit()
            if replacement_to_remove:
                try:
                    remove_artifact_version(settings.UPLOAD_DIR, course_id, artifact, replacement_to_remove)
                except OSError as exc:
                    logger.warning(
                        "Could not remove replaced %s artifact %s for course %s: %s",
                        artifact,
                        replacement_to_remove,
                        course_id,
                        exc,
                    )
            if version_id and status in ("ready", "error"):
                self._generation_versions.pop((course_id, artifact), None)
        except Exception as e:
            db.rollback()
            logger.error(f"Error setting artifact status for {artifact}: {e}")
        finally:
            db.close()

    def get_artifact_status(self, course_id: str, artifact: str, version_id: Optional[str] = None, db_session_factory=None) -> Dict[str, Any]:
        """Read per-artifact generation status from Course.metadata_json."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if not course:
                return {}
            meta, _ = migrate_legacy_artifact_metadata(self._metadata_dict(course))
            entry = meta.get("study_pack", {}).get("artifacts", {}).get(artifact, {})
            if not isinstance(entry, dict) or "versions" not in entry:
                return entry if isinstance(entry, dict) else {}
            return dict(entry.get("versions", {}).get(version_id or entry.get("active"), {}))
        finally:
            db.close()

    def get_regen_usage(self, course_id: str, db_session_factory=None) -> Dict[str, int]:
        """Read the current regeneration usage per artifact from Course.metadata_json."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if not course:
                return {}
            meta = course.metadata_json or "{}"
            if isinstance(meta, str):
                try:
                    meta = json.loads(meta)
                except Exception:
                    meta = {}
            return meta.get("study_pack", {}).get("regen_counts", {})
        finally:
            db.close()

    def check_and_record_regen_attempt(
        self, course_id: str, artifact: str, db_session_factory=None
    ) -> Tuple[bool, int, int]:
        """Gate a generate-* call as a "regeneration" only if the artifact already has a
        ready/error status (i.e. it's been generated/attempted before) — the very first
        generation is always free and never counted. Returns (allowed, used, max). When
        the call is allowed and is in fact a regeneration, the counter is incremented and
        persisted as a side effect.

        Fails open on any bookkeeping error — a DB hiccup in the counter must never block
        the user's actual ability to generate content."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if not course:
                return True, 0, MAX_REGENERATIONS

            meta = course.metadata_json or "{}"
            if isinstance(meta, str):
                try:
                    meta = json.loads(meta)
                except Exception:
                    meta = {}
            study_pack = meta.get("study_pack", {})
            artifacts = study_pack.get("artifacts", {})
            artifact_entry = artifacts.get(artifact, {})
            if isinstance(artifact_entry, dict) and "versions" in artifact_entry:
                active = artifact_entry.get("active")
                prior_status = artifact_entry.get("versions", {}).get(active, {}).get("status")
            else:
                prior_status = artifact_entry.get("status")

            if prior_status not in ("ready", "error"):
                # First-ever attempt for this artifact: not a regen, never blocked/counted.
                regen_counts = study_pack.get("regen_counts", {})
                return True, regen_counts.get(artifact, 0), MAX_REGENERATIONS

            regen_counts = study_pack.get("regen_counts", {})
            used = regen_counts.get(artifact, 0)
            if used >= MAX_REGENERATIONS:
                return False, used, MAX_REGENERATIONS

            regen_counts[artifact] = used + 1
            study_pack["regen_counts"] = regen_counts
            meta["study_pack"] = study_pack
            course.metadata_json = json.dumps(meta, ensure_ascii=False)
            db.commit()
            return True, used + 1, MAX_REGENERATIONS
        except Exception as e:
            db.rollback()
            logger.error(f"Error checking regen limit for {artifact}: {e}")
            return True, 0, MAX_REGENERATIONS
        finally:
            db.close()

    def _resolve_topic(self, course_id: str, topic: Optional[str] = None, db_session_factory=None) -> str:
        """Resolve actual course/document topic if topic is missing or hardcoded generic AI string."""
        ignore_list = ["AI Quiz", "AI Overview", "AI Video", "AI Course", "General Students", ""]
        if topic and str(topic).strip() and str(topic).strip() not in ignore_list:
            return str(topic).strip()

        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if course and course.name and course.name.strip():
                return course.name.strip()
            if course and course.filenames and len(course.filenames) > 0:
                fn = str(course.filenames[0])
                for ext in [".pdf", ".docx", ".txt", ".PPTX", ".PDF", ".DOCX", ".TXT"]:
                    if fn.lower().endswith(ext.lower()):
                        fn = fn[:-len(ext)]
                fn = _strip_leading_id_token(fn.strip())
                if fn.strip():
                    return fn.strip()
            if course and course.metadata_json:
                try:
                    meta = json.loads(course.metadata_json) if isinstance(course.metadata_json, str) else course.metadata_json
                    if isinstance(meta, dict) and meta.get("title"):
                        return str(meta["title"]).strip()
                except Exception:
                    pass
        finally:
            db.close()
        return "Nội dung tài liệu chính"

    def _get_doc_names(self, course_id: str, db_session_factory=None) -> str:
        """Fetch the course's source document filenames as a comma-separated string."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            if course and course.filenames:
                return ", ".join(str(fn) for fn in course.filenames)
            return ""
        finally:
            db.close()

    def generate_book(
        self,
        course_id: str,
        detail_level: str = "Tiêu chuẩn",
        user_prompt: str = "",
        db_session_factory=None,
        **kwargs,
    ) -> Optional[BookOutput]:
        """Execute the multi-pass generation pipeline for the Study Guide Book:
        outline pass -> one LLM call per chapter (with per-chapter retrieval) -> assemble -> validate -> PDF.

        On any failure, records an "error" artifact status and returns None instead of writing a
        partial/placeholder artifact.
        """
        logger.info(f"Starting Book generation for course {course_id}")
        transaction, artifact_dir = self._start_version_write(course_id, "book", kwargs.get("version_id"))
        try:
            self._set_artifact_status(course_id, "book", "processing", progress=5, db_session_factory=db_session_factory)
            self._require_course_not_processing(course_id, db_session_factory)

            book_llm = self._llm_for("book")
            context, base_ids = self._retrieve_context(course_id, k=20, db_session_factory=db_session_factory)
            self._require_context(context)
            doc_names = self._get_doc_names(course_id, db_session_factory)
            outline = book_llm.generate_book_outline(context, detail_level, user_prompt, doc_names)
            plans = outline.chapters[:8]
            if len(plans) < 4:
                raise ValueError(f"Dàn ý chỉ có {len(plans)} chương, cần tối thiểu 4 chương.")

            self._set_artifact_status(course_id, "book", "processing", progress=15, db_session_factory=db_session_factory)

            all_ids = set(base_ids)
            chapters: List[BookChapter] = []
            total = len(plans)
            for i, plan in enumerate(plans):
                ch_context, ch_ids = self._retrieve_context(
                    course_id, query=plan.retrieval_query, k=10, db_session_factory=db_session_factory
                )
                if not ch_context:
                    ch_context, ch_ids = context, base_ids
                all_ids.update(ch_ids)

                content = book_llm.generate_book_chapter(
                    outline.title, plan, total, ch_context, detail_level, ch_ids
                )

                chapters.append(
                    BookChapter(
                        chapter_title=content.chapter_title or plan.chapter_title,
                        introduction=content.introduction,
                        objectives=content.objectives,
                        sections=content.sections,
                        key_points=content.key_points,
                        review_questions=content.review_questions,
                        source_chunk_ids=content.source_chunk_ids,
                    )
                )
                progress = 15 + int(75 * (i + 1) / total)
                self._set_artifact_status(course_id, "book", "processing", progress=progress, db_session_factory=db_session_factory)

            book = BookOutput(title=outline.title, summary=outline.summary, preface=outline.preface, chapters=chapters)
            validated_output, score, warnings = validate_and_score_output(book, "book", list(all_ids))
            if warnings:
                logger.warning(f"Book generation warnings for {course_id}: {warnings}")

            self._save_artifact_json(course_id, "book.json", validated_output, artifact_dir)
            self._generate_pdf_book(course_id, validated_output, artifact_dir)
            self._finish_version_write(transaction, True)
            self._update_course_metadata(course_id, "book", score, db_session_factory)
            self._set_artifact_status(course_id, "book", "ready", progress=100, db_session_factory=db_session_factory)
            return validated_output
        except Exception as e:
            self._finish_version_write(transaction, False)
            logger.error(f"Book generation failed for course {course_id}: {e}", exc_info=True)
            self._set_artifact_status(
                course_id, "book", "error", error=str(e)[:500], db_session_factory=db_session_factory
            )
            return None

    def generate_slides(self, course_id: str, topic: str = "AI Overview", num_slides: int = 15, db_session_factory=None, **kwargs) -> Optional[SlidesOutput]:
        """Execute full generation pipeline for Presentation Slides.

        On any failure, records an "error" artifact status and returns None instead of
        letting a background-task exception vanish silently.
        """
        logger.info(f"Starting Slides generation for course {course_id}")
        transaction, artifact_dir = self._start_version_write(course_id, "slides", kwargs.get("version_id"))
        try:
            self._set_artifact_status(course_id, "slides", "processing", progress=10, db_session_factory=db_session_factory)
            self._require_course_not_processing(course_id, db_session_factory)
            resolved_topic = self._resolve_topic(course_id, topic, db_session_factory=db_session_factory)
            context, valid_chunk_ids = self._retrieve_context(
                course_id, query=resolved_topic, db_session_factory=db_session_factory
            )
            self._require_context(context)
            raw_output = self._llm_for("slides").generate_slides(context, resolved_topic, num_slides, valid_chunk_ids)
            validated_output, score, warnings = validate_and_score_output(raw_output, "slides", valid_chunk_ids)
            if warnings:
                logger.warning(f"Slides generation warnings for {course_id}: {warnings}")
            validated_output = _clean_slides_output(validated_output)

            self._set_artifact_status(course_id, "slides", "processing", progress=70, db_session_factory=db_session_factory)
            self._save_artifact_json(course_id, "slides.json", validated_output, artifact_dir)
            self._generate_pptx_slides(course_id, validated_output, artifact_dir)
            self._finish_version_write(transaction, True)
            self._update_course_metadata(course_id, "slides", score, db_session_factory)
            self._set_artifact_status(course_id, "slides", "ready", progress=100, db_session_factory=db_session_factory)
            return validated_output
        except Exception as e:
            self._finish_version_write(transaction, False)
            logger.error(f"Slides generation failed for course {course_id}: {e}", exc_info=True)
            self._set_artifact_status(
                course_id, "slides", "error", error=str(e)[:500], db_session_factory=db_session_factory
            )
            return None

    def generate_quiz(self, course_id: str, topic: str = "AI Quiz", quantity: int = 5, difficulty: str = "mixed", db_session_factory=None, **kwargs) -> Optional[QuizOutput]:
        """Execute full generation pipeline for Multiple Choice Quiz.

        On any failure, records an "error" artifact status and returns None instead of
        letting a background-task exception vanish silently.
        """
        logger.info(f"Starting Quiz generation for course {course_id} (quantity={quantity}, difficulty={difficulty})")
        transaction, artifact_dir = self._start_version_write(course_id, "quiz", kwargs.get("version_id"))
        try:
            self._set_artifact_status(course_id, "quiz", "processing", progress=10, db_session_factory=db_session_factory)
            self._require_course_not_processing(course_id, db_session_factory)
            resolved_topic = self._resolve_topic(course_id, topic, db_session_factory=db_session_factory)
            context, valid_chunk_ids = self._retrieve_context(
                course_id, query=resolved_topic, db_session_factory=db_session_factory
            )
            self._require_context(context)
            raw_output = self._llm_for("quiz").generate_quiz(context, resolved_topic, quantity, valid_chunk_ids, difficulty=difficulty)
            validated_output, score, warnings = validate_and_score_output(raw_output, "quiz", valid_chunk_ids)
            if warnings:
                logger.warning(f"Quiz generation warnings for {course_id}: {warnings}")

            self._set_artifact_status(course_id, "quiz", "processing", progress=70, db_session_factory=db_session_factory)
            self._save_artifact_json(course_id, "quiz.json", validated_output, artifact_dir)
            self._generate_pdf_quiz_key(course_id, validated_output, artifact_dir)
            self._finish_version_write(transaction, True)
            self._update_course_metadata(course_id, "quiz", score, db_session_factory)
            self._set_artifact_status(course_id, "quiz", "ready", progress=100, db_session_factory=db_session_factory)
            return validated_output
        except Exception as e:
            self._finish_version_write(transaction, False)
            logger.error(f"Quiz generation failed for course {course_id}: {e}", exc_info=True)
            self._set_artifact_status(
                course_id, "quiz", "error", error=str(e)[:500], db_session_factory=db_session_factory
            )
            return None

    def generate_vid(
        self,
        course_id: str,
        topic: str = "AI Video",
        fmt: str = "standard",
        voice: str = "female",
        user_prompt: str = "",
        db_session_factory=None,
        **kwargs,
    ) -> Optional[VidOutput]:
        """Execute full generation pipeline for the narrated Video: LLM script (1 call) ->
        per-scene TTS narration + still frame -> ffmpeg mux/concat into vid.mp4.

        On any failure, records an "error" artifact status and returns None instead of
        letting a background-task exception vanish silently.
        """
        logger.info(f"Starting Video generation for course {course_id} (format={fmt}, voice={voice})")
        transaction, artifact_dir = self._start_version_write(course_id, "vid", kwargs.get("version_id"))
        try:
            self._set_artifact_status(course_id, "vid", "processing", progress=10, db_session_factory=db_session_factory)
            self._require_course_not_processing(course_id, db_session_factory)
            resolved_topic = self._resolve_topic(course_id, topic, db_session_factory=db_session_factory)
            context, valid_chunk_ids = self._retrieve_context(
                course_id, query=resolved_topic, db_session_factory=db_session_factory
            )
            self._require_context(context)
            raw_output = self._llm_for("vid").generate_vid(
                context, resolved_topic, fmt, user_prompt, valid_chunk_ids
            )
            validated_output, score, warnings = validate_and_score_output(raw_output, "vid", valid_chunk_ids)
            if warnings:
                logger.warning(f"Vid generation warnings for {course_id}: {warnings}")
            validated_output = _clean_vid_output(validated_output)
            scene_visual_map = self._build_scene_visual_map(course_id, validated_output)

            self._set_artifact_status(course_id, "vid", "processing", progress=25, db_session_factory=db_session_factory)

            def _progress_cb(fraction: float) -> None:
                self._set_artifact_status(
                    course_id, "vid", "processing", progress=25 + int(60 * fraction),
                    db_session_factory=db_session_factory,
                )

            self._generate_video_mp4(
                course_id,
                validated_output,
                fmt,
                voice,
                progress_cb=_progress_cb,
                artifact_dir=artifact_dir,
                scene_visual_map=scene_visual_map,
            )

            self._set_artifact_status(course_id, "vid", "processing", progress=90, db_session_factory=db_session_factory)
            self._save_artifact_json(course_id, "vid.json", validated_output, artifact_dir)
            self._finish_version_write(transaction, True)
            self._update_course_metadata(course_id, "vid", score, db_session_factory)
            self._set_artifact_status(course_id, "vid", "ready", progress=100, db_session_factory=db_session_factory)
            return validated_output
        except Exception as e:
            self._finish_version_write(transaction, False)
            logger.error(f"Vid generation failed for course {course_id}: {e}", exc_info=True)
            self._set_artifact_status(
                course_id, "vid", "error", error=str(e)[:500], db_session_factory=db_session_factory
            )
            return None

    def get_study_pack(self, course_id: str, db_session_factory=None) -> StudyPackResponse:
        """Build and return full StudyPackResponse from filesystem artifacts and DB state."""
        db = self._get_db(db_session_factory)
        try:
            course = db.query(Course).filter(Course.id == course_id).first()
            status_val = course.status if course else "ready"
            chunk_cnt = course.chunk_count if course else 0
            q_score = course.quality_score if course else 0
            meta = course.metadata_json if course and course.metadata_json else "{}"
            if isinstance(meta, str):
                try:
                    meta = json.loads(meta)
                except Exception:
                    meta = {}
            sp_meta = meta.get("study_pack", {})
        finally:
            db.close()

        def active_artifact_dir(artifact: str) -> str:
            entry = sp_meta.get("artifacts", {}).get(artifact, {})
            if isinstance(entry, dict) and entry.get("active"):
                return artifact_directory_path(settings.UPLOAD_DIR, course_id, artifact, entry["active"])
            return self._get_artifact_dir(course_id)

        book_dir = active_artifact_dir("book")
        slides_dir = active_artifact_dir("slides")
        quiz_dir = active_artifact_dir("quiz")
        vid_dir = active_artifact_dir("vid")
        book_json = self._load_artifact_json(course_id, "book.json", book_dir)
        slides_json = self._load_artifact_json(course_id, "slides.json", slides_dir)
        quiz_json = self._load_artifact_json(course_id, "quiz.json", quiz_dir)
        vid_json = self._load_artifact_json(course_id, "vid.json", vid_dir)
        has_book = book_json is not None
        has_book_pdf = os.path.exists(os.path.join(book_dir, "book.pdf"))
        has_slide = slides_json is not None
        has_slide_pptx = os.path.exists(os.path.join(slides_dir, "slide.pptx"))
        has_quiz = quiz_json is not None
        has_quiz_key = os.path.exists(os.path.join(quiz_dir, "quiz-key.pdf"))
        has_vid = vid_json is not None

        readiness_meta = sp_meta.get("readiness", {})
        readiness = ReadinessData(
            study_guide_pdf=has_book or readiness_meta.get("study_guide_pdf", False),
            slides=has_slide or readiness_meta.get("slides", False),
            quiz=has_quiz or readiness_meta.get("quiz", False),
            vid=has_vid or readiness_meta.get("vid", False),
        )

        quality_scores_meta = sp_meta.get("quality_scores", {})
        quality_scores = QualityScoresData(
            study_guide_pdf=quality_scores_meta.get("study_guide_pdf", 0),
            slides=quality_scores_meta.get("slides", 0),
            quiz=quality_scores_meta.get("quiz", 0),
            vid=quality_scores_meta.get("vid", 0),
        )

        grounding_meta = sp_meta.get("grounding", {})
        grounding = GroundingData(
            num_chunks=grounding_meta.get("num_chunks", chunk_cnt if (has_book or has_slide or has_quiz or has_vid) else 0),
            quality_score=grounding_meta.get("quality_score", q_score if (has_book or has_slide or has_quiz or has_vid) else 0),
            warnings=grounding_meta.get("warnings", []),
        )

        regen_limits = RegenLimitsData(max=MAX_REGENERATIONS, used=sp_meta.get("regen_counts", {}))

        return StudyPackResponse(
            course_id=course_id,
            stats=StudyPackStats(
                course_id=course_id,
                status=status_val,
                has_book=has_book,
                has_book_pdf=has_book_pdf,
                has_slide=has_slide,
                has_slide_pptx=has_slide_pptx,
                has_quiz=has_quiz,
                has_quiz_answer_key=has_quiz_key,
                has_vid=has_vid,
                quality_score=q_score,
                num_chunks=chunk_cnt,
            ),
            study_pack=StudyPackData(
                title=book_json.get("title", "Course Title") if book_json else "Course Title",
                book=book_json,
                slides=slides_json,
                quiz=quiz_json.get("questions", []) if quiz_json else [],
                vid=vid_json,
                readiness=readiness,
                quality_scores=quality_scores,
                grounding=grounding,
                regen_limits=regen_limits,
            ),
        )
